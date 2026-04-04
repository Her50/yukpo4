# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x1E, 0x42)
NAVY2  = RGBColor(0x13, 0x2B, 0x5E)
GOLD   = RGBColor(0xD4, 0xAF, 0x37)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SILVER = RGBColor(0xB0, 0xC4, 0xDE)
GREEN  = RGBColor(0x00, 0xC8, 0x7A)
DKGRN  = RGBColor(0x0A, 0x30, 0x1A)
RED    = RGBColor(0xFF, 0x6B, 0x6B)
CYAN   = RGBColor(0x4F, 0xC3, 0xF7)
ORANGE = RGBColor(0xFF, 0x99, 0x33)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)

W = Cm(25.4)
H = Cm(19.05)

LOGO = r"c:\Users\23767\yukpomnang2\mobile\assets\icon.png"

# ── Primitives ────────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(s, c):
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = c

def rect(s, x, y, w, h, fill=None, line=None, lw=1.5):
    shp = s.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def txt(s, t, x, y, w, h, size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(t)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def logo_img(s, x, y, w, h):
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, int(x), int(y), int(w), int(h))

def hdr(s, title, sub, num):
    bg(s, NAVY)
    rect(s, 0, 0, W, Cm(1.35), fill=GOLD)
    rect(s, 0, Cm(1.35), Cm(0.4), H - Cm(2.0), fill=NAVY2)
    txt(s, str(num), W - Cm(1.7), Cm(0.1), Cm(1.5), Cm(1.15),
        size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(s, title, Cm(0.7), Cm(1.5), W - Cm(2.4), Cm(1.5),
        size=21, bold=True, color=WHITE)
    if sub:
        txt(s, sub, Cm(0.7), Cm(3.05), W - Cm(1.4), Cm(0.8),
            size=11, color=GOLD, italic=True)
    rect(s, 0, H - Cm(0.6), W, Cm(0.6), fill=GOLD)
    txt(s, "YUKPO  |  YUKPO COMPANY SARL  |  Confidentiel  |  2026",
        Cm(0.5), H - Cm(0.6), Cm(18), Cm(0.6), size=9, color=NAVY)

def scard(s, x, y, w, h, number, label, acc):
    rect(s, x, y, w, h, fill=NAVY2, line=acc)
    rect(s, x, y, Cm(0.32), h, fill=acc)
    txt(s, number, x + Cm(0.5), y + Cm(0.28), w - Cm(0.65), h * 0.52,
        size=24, bold=True, color=acc, align=PP_ALIGN.CENTER)
    txt(s, label, x + Cm(0.5), y + h * 0.54, w - Cm(0.65), h * 0.38,
        size=10, color=SILVER, align=PP_ALIGN.CENTER)

def svcard(s, x, y, w, h, name, detail, acc=GOLD):
    rect(s, x, y, w, h, fill=NAVY2)
    rect(s, x, y, Cm(0.28), h, fill=acc)
    txt(s, name, x + Cm(0.5), y + Cm(0.18), w - Cm(0.68), Cm(0.72),
        size=11, bold=True, color=acc)
    txt(s, detail, x + Cm(0.5), y + Cm(0.9), w - Cm(0.68), h - Cm(1.05),
        size=9, color=SILVER)

def icard(s, x, y, w, h, label, value, val_size, val_col, detail, acc):
    rect(s, x, y, w, h, fill=NAVY2, line=acc)
    rect(s, x, y, w, Cm(0.26), fill=acc)
    txt(s, label, x + Cm(0.32), y + Cm(0.38), w - Cm(0.5), Cm(0.72),
        size=11, bold=True, color=acc)
    txt(s, value, x + Cm(0.32), y + Cm(1.12), w - Cm(0.5), Cm(1.55),
        size=val_size, bold=True, color=val_col, align=PP_ALIGN.CENTER)
    rect(s, x + Cm(0.5), y + Cm(2.72), w - Cm(1.0), Cm(0.05), fill=acc)
    txt(s, detail, x + Cm(0.32), y + Cm(2.9), w - Cm(0.5), h - Cm(3.1),
        size=10, color=SILVER)

def setcell(cell, text, bgc, fgc, sz=11, bold=False, align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bgc
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for r in p.runs:
        r.text = ""
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.color.rgb = fgc

# ─────────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────────

def s01(prs, subtitle=""):
    sl = blank(prs)
    bg(sl, NAVY)
    rect(sl, 0, 0, W, Cm(0.45), fill=GOLD)
    rect(sl, 0, H - Cm(0.45), W, Cm(0.45), fill=GOLD)
    rect(sl, 0, Cm(0.45), Cm(10.2), H - Cm(0.9), fill=NAVY2)
    rect(sl, 0, Cm(0.45), Cm(0.38), H - Cm(0.9), fill=GOLD)

    # Logo sur fond blanc — on place un rectangle blanc derrière
    rect(sl, Cm(0.55), Cm(0.9), Cm(9.15), Cm(5.2), fill=WHITE)
    logo_img(sl, Cm(0.55), Cm(0.9), Cm(9.15), Cm(5.2))

    rect(sl, Cm(1.2), Cm(6.25), Cm(8.0), Cm(0.07), fill=GOLD)
    txt(sl, "DOSSIER INVESTISSEUR",
        Cm(0.7), Cm(6.45), Cm(9.0), Cm(0.72),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(sl, subtitle, Cm(0.7), Cm(7.68), Cm(9.0), Cm(0.62),
            size=10, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    kpis = [
        ("120 M FCFA",  "Objectif de levée de fonds",              GOLD),
        ("7,5 M+",       "Habitants ciblés — Yaoundé + Douala",     CYAN),
        ("T1 2028",      "Rentabilité nette (meilleur cas)",         GREEN),
    ]
    for i, (val, lbl, col) in enumerate(kpis):
        cy = Cm(2.0) + i * Cm(4.45)
        rect(sl, Cm(10.8), cy, Cm(13.8), Cm(4.0), fill=NAVY2, line=col)
        rect(sl, Cm(10.8), cy, Cm(0.38), Cm(4.0), fill=col)
        txt(sl, val, Cm(11.4), cy + Cm(0.35), Cm(12.8), Cm(1.7),
            size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, lbl, Cm(11.4), cy + Cm(2.1), Cm(12.8), Cm(1.6),
            size=12, color=SILVER, align=PP_ALIGN.CENTER)

    txt(sl, "Présenté par  LELE SIAKA Hernandez  |  CEO & Fondateur",
        Cm(0.7), H - Cm(3.1), Cm(9.0), Cm(0.72), size=9, color=SILVER, align=PP_ALIGN.CENTER)
    txt(sl, "Yaoundé, Cameroun  |  Avril 2026",
        Cm(0.7), H - Cm(2.4), Cm(9.0), Cm(0.65), size=9, color=GOLD, align=PP_ALIGN.CENTER)


def s02(prs):
    sl = blank(prs)
    hdr(sl, "Une Opportunité de Marché Massive et Inexploitée",
        "Yaoundé + Douala : 7,5 M d'habitants, zéro super-app locale opérationnelle", 2)

    cards = [
        ("7,5 M+",       "Habitants ciblés\nYaoundé + Douala",    GOLD),
        ("74 %",          "Sans accès\naux services formels",      CYAN),
        ("0",             "Super-app intégrée\nlocale existante",  RED),
        ("8,5 Mrd FCFA",  "Marché adressable\nestimé 2028",        GREEN),
    ]
    cw = Cm(5.65)
    ch = Cm(4.9)
    gap = (W - Cm(1.0) - 4 * cw) / 3
    for i, (val, lbl, col) in enumerate(cards):
        scard(sl, Cm(0.5) + i * (cw + gap), Cm(4.1), cw, ch, val, lbl, col)

    rect(sl, Cm(0.5), Cm(9.3), W - Cm(1.0), Cm(0.07), fill=GOLD)
    cities = [
        ("Yaoundé",  "3,5 M hab.  |  Capitale politique & tertiaire  |  62 % smartphone  |  Marché : 3,8 Mrd FCFA"),
        ("Douala",   "4,0 M hab.  |  Capitale économique & portuaire  |  67 % smartphone  |  Marché : 4,7 Mrd FCFA"),
    ]
    for i, (city, info) in enumerate(cities):
        cy = Cm(9.6) + i * Cm(2.15)
        rect(sl, Cm(0.5), cy, Cm(3.8), Cm(1.85), fill=GOLD if i == 0 else CYAN)
        txt(sl, city, Cm(0.5), cy + Cm(0.3), Cm(3.8), Cm(1.2),
            size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        rect(sl, Cm(4.5), cy, W - Cm(5.0), Cm(1.85), fill=NAVY2)
        txt(sl, info, Cm(4.85), cy + Cm(0.3), W - Cm(5.5), Cm(1.3), size=11, color=WHITE)

    txt(sl, "Croissance mobile +18 %/an  |  Fenêtre d'entrée unique : aucun concurrent intégré  |  Besoin urgent non couvert",
        Cm(0.5), Cm(13.95), W - Cm(1.0), Cm(0.75),
        size=11, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
    rect(sl, Cm(0.5), Cm(14.85), W - Cm(1.0), Cm(1.7), fill=NAVY2, line=ORANGE)
    rect(sl, Cm(0.5), Cm(14.85), Cm(0.22), Cm(1.7), fill=ORANGE)
    txt(sl,
        "Note : l'estimation de 8,5 Mrd FCFA est limitée aux services numériques immédiats "
        "(alimentation, boutiques, Creator Studio, Annonce Rapide). Elle n'intègre PAS : "
        "transport, immobilier, santé, Bourse du Livre, publicité in-app — secteurs dont "
        "YUKPO adresse aussi le besoin. Le marché réel total est bien supérieur.",
        Cm(0.9), Cm(14.93), W - Cm(1.6), Cm(1.55),
        size=9, color=ORANGE, italic=True)


def s03(prs):
    sl = blank(prs)
    hdr(sl, "Les Problèmes Quotidiens des Camerounais",
        "Des heures perdues, des services éparpillés, une confiance difficile à établir", 3)

    problems = [
        ("Perte de Temps",
         "Trouver un transport, le bon médicament dans la bonne pharmacie\n"
         "(sans faire la tournée de 5 officines), un artisan disponible...\n"
         "Chaque démarche prend des heures. 2 h perdues par jour en moyenne.",
         GOLD),
        ("Zéro Confiance",
         "Aucun tarif affiché, aucun avis vérifié, aucune garantie.\n"
         "Les arnaques sont fréquentes, les recours inexistants.\n"
         "La méfiance bloque la formalisation et la consommation de services.",
         CYAN),
        ("Marché Ultra-Fragmenté",
         "Chaque service fonctionne dans sa propre bulle :\n"
         "contact perso, prix opaque, aucune traçabilité.\n"
         "Impossible de planifier, comparer ou confier une tâche sereinement.",
         GREEN),
    ]
    cw = Cm(7.55)
    ch = Cm(13.4)
    gap = (W - 3 * cw) / 4
    for i, (title, body, col) in enumerate(problems):
        cx = gap + i * (cw + gap)
        rect(sl, cx, Cm(4.0), cw, ch, fill=NAVY2, line=col)
        rect(sl, cx, Cm(4.0), cw, Cm(0.27), fill=col)
        txt(sl, title, cx + Cm(0.38), Cm(4.42), cw - Cm(0.58), Cm(0.82),
            size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        rect(sl, cx + Cm(0.5), Cm(5.3), cw - Cm(1.0), Cm(0.06), fill=col)
        txt(sl, body, cx + Cm(0.38), Cm(5.5), cw - Cm(0.58), ch - Cm(1.65), size=11, color=SILVER)


def s04(prs):
    sl = blank(prs)
    hdr(sl, "YUKPO — L'Écosystème Complet",
        "Services essentiels + marché + création de contenu + intelligence artificielle — en un seul endroit", 4)

    services = [
        ("Transport & Mobilité",     "Taxi, moto, covoiturage, voyages",       GOLD),
        ("Alimentation & Livraison", "Repas chauds, courses à domicile",       GOLD),
        ("Pharmacie & Santé",        "Localisation médicaments dispo, cliniques", GOLD),
        ("Eau, Gaz & Énergie",       "Bonbonnes livrées en moins d'une heure", GOLD),
        ("Artisans & Services",      "Plombiers, électriciens, peintres...",    GOLD),
        ("Boutiques & Marché",       "Mode, électroménager, décoration",        CYAN),
        ("Immobilier & Auto",        "Location, vente, devis travaux",           CYAN),
        ("Bourse du Livre",          "Livres scolaires & fournitures — digital", ORANGE),
        ("Creator Studio",           "Vidéos & visuels professionnels en 2 min",GREEN),
        ("Annonce Rapide",           "Produit ou prestation publié en 2 min",    GREEN),
        ("Navigation Intelligente",  "Recherche prédictive & recommandations",   PURPLE),
        ("YukpoIA",                  "Assistant IA intégré — disponible 24h/24", PURPLE),
    ]
    cw = Cm(7.6)
    ch = Cm(2.58)
    gap_x = (W - 3 * cw) / 4
    gap_y = Cm(0.27)
    y0 = Cm(4.05)
    for i, (name, detail, col) in enumerate(services):
        c = i % 3
        r = i // 3
        svcard(sl, gap_x + c * (cw + gap_x), y0 + r * (ch + gap_y), cw, ch, name, detail, col)


def s05(prs):
    sl = blank(prs)
    hdr(sl, "Modèle de Revenus — À la Consommation",
        "YUKPO génère des revenus sur chaque transaction réalisée sur la plateforme", 5)

    pillars = [
        ("Commission sur Transactions",
         "1 à 5 %\nselon le service",
         "Sur chaque course, livraison, artisan ou vente réalisée, YUKPO prélève "
         "une commission plafonnée à 5 %. Transport & Bourse du Livre : 5 %, "
         "alimentation : 4 %, boutiques : 3 %, immobilier : 2 %. "
         "Pharmacie : service de localisation sans commission sur les médicaments.",
         GOLD),
        ("Forfaits Prestataires",
         "3 000 – 10 000 F\npar mois",
         "Chaque prestataire, boutique ou vendeur paie un forfait mensuel "
         "pour apparaître sur la plateforme. Revenu récurrent et prévisible, "
         "indépendant du volume de transactions. "
         "Bourse du Livre, Creator Studio et Annonce Rapide inclus.",
         CYAN),
        ("Mise en Avant Commerciale",
         "Forfaits premium\nsponsoring & publicité",
         "Les acteurs souhaitant apparaître en 1re position, lancer des promotions "
         "ciblées ou diffuser des publicités in-app paient des forfaits additionnels. "
         "Revenu passif haute marge, non corrélé au volume de commandes.",
         GREEN),
    ]
    cw = Cm(7.55)
    ch = Cm(13.4)
    gap = (W - 3 * cw) / 4
    for i, (label, value, detail, acc) in enumerate(pillars):
        icard(sl, gap + i * (cw + gap), Cm(4.0), cw, ch, label, value, 18, acc, detail, acc)

    rect(sl, Cm(0.5), H - Cm(1.4), W - Cm(1.0), Cm(0.75), fill=NAVY2, line=GOLD)
    txt(sl, "Pas d'abonnement imposé — YUKPO mise sur la fréquence d'usage et la valeur de chaque transaction.",
        Cm(0.8), H - Cm(1.38), W - Cm(1.6), Cm(0.72),
        size=10, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


def s06(prs):
    sl = blank(prs)
    hdr(sl, "Entonnoir de Marché — Yaoundé + Douala",
        "Un potentiel de monétisation massif, progressivement activé", 6)

    funnel = [
        (Cm(0.4),  W - Cm(0.8),  "Population urbaine Yaoundé + Douala",    "7 500 000 hab.",        WHITE,  NAVY2),
        (Cm(1.4),  W - Cm(2.8),  "Utilisateurs smartphone actifs",          "4 500 000",             SILVER, NAVY2),
        (Cm(2.8),  W - Cm(5.6),  "Cible YUKPO — Année 1",              "120 000 utilisateurs",  GOLD,   NAVY2),
        (Cm(4.6),  W - Cm(9.2),  "Utilisateurs actifs / mois (objectif)",   "25 000",                GOLD,   DKGRN),
        (Cm(6.8),  W - Cm(13.6), "Commandeurs réguliers (2 cmd+/mois)",     "8 000",                 GREEN,  DKGRN),
    ]
    bh = Cm(1.82)
    y0 = Cm(4.1)
    for j, (lx, lw, label, val, fc, bgc) in enumerate(funnel):
        ry = y0 + j * (bh + Cm(0.22))
        rect(sl, lx, ry, lw, bh, fill=bgc, line=GOLD)
        txt(sl, label, lx + Cm(0.4), ry + Cm(0.15), int(lw * 0.62), bh, size=11, color=fc)
        txt(sl, val, lx + int(lw * 0.64), ry + Cm(0.15), int(lw * 0.34), bh,
            size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

    rect(sl, Cm(0.4), Cm(14.95), W - Cm(0.8), Cm(0.07), fill=CYAN)
    txt(sl, "Extension 2027 : Bafoussam + densification nationale",
        Cm(0.6), Cm(15.2), W * 0.48, Cm(0.8), size=11, color=CYAN, italic=True)
    txt(sl, "International : Côte d'Ivoire (Q4 2027)  |  Sénégal (S2 2028)",
        W * 0.52, Cm(15.2), W * 0.46, Cm(0.8),
        size=11, color=ORANGE, italic=True, align=PP_ALIGN.RIGHT)
    rect(sl, Cm(0.4), Cm(16.25), W - Cm(0.8), Cm(1.35), fill=NAVY2, line=GOLD)
    txt(sl, "Marché adressable cumulé Yaoundé + Douala (2028) :  8,5 milliards FCFA",
        Cm(0.8), Cm(16.4), W - Cm(1.6), Cm(1.05),
        size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def s07(prs):
    """Competitor analysis — visual score matrix (heatmap style)"""
    sl = blank(prs)
    hdr(sl, "Analyse Concurrentielle — Positionnement YUKPO",
        "Score /5 sur 6 dimensions clés — YUKPO domine toutes les catégories", 7)

    dims = [
        "Services intégrés",
        "Présence locale",
        "Fiabilité & Avis",
        "Innovation IA",
        "Expér. utilisateur",
        "Transparence prix",
    ]
    actors     = ["YUKPO", "Yango", "Jumia Food", "Locaux"]
    act_colors = [GOLD, CYAN, GREEN, ORANGE]
    scores = [
        [5, 2, 2, 1],
        [5, 3, 2, 4],
        [5, 3, 3, 2],
        [5, 2, 1, 1],
        [5, 3, 3, 2],
        [5, 2, 4, 1],
    ]
    score_fill = {5: GOLD, 4: GREEN, 3: CYAN, 2: ORANGE, 1: RED}
    score_text = {5: NAVY, 4: NAVY, 3: NAVY, 2: NAVY, 1: WHITE}

    # Matrix layout
    mx     = Cm(0.4)
    my     = Cm(4.05)
    lbl_w  = Cm(4.4)
    col_w  = Cm(2.85)
    hdr_h  = Cm(1.15)
    row_h  = Cm(2.1)
    n_dim  = len(dims)
    n_act  = len(actors)
    mat_w  = lbl_w + n_act * col_w
    mat_h  = hdr_h + n_dim * row_h

    # Header row — actor names
    rect(sl, int(mx), int(my), int(lbl_w), int(hdr_h), fill=NAVY2)
    txt(sl, "Dimension clé", int(mx + Cm(0.2)), int(my + Cm(0.2)),
        int(lbl_w - Cm(0.3)), int(hdr_h - Cm(0.3)),
        size=10, bold=True, color=SILVER, align=PP_ALIGN.CENTER)
    for j, (actor, acol) in enumerate(zip(actors, act_colors)):
        ax = int(mx + lbl_w + j * col_w)
        rect(sl, ax, int(my), int(col_w - Cm(0.06)), int(hdr_h),
             fill=acol if j == 0 else NAVY2, line=acol)
        txt(sl, actor, ax + int(Cm(0.1)), int(my + Cm(0.2)),
            int(col_w - Cm(0.2)), int(hdr_h - Cm(0.3)),
            size=11, bold=True,
            color=NAVY if j == 0 else acol, align=PP_ALIGN.CENTER)

    # Data rows
    for i, (dim, row_sc) in enumerate(zip(dims, scores)):
        ry  = int(my + hdr_h + i * row_h)
        row_bg = NAVY2 if i % 2 == 0 else NAVY
        rect(sl, int(mx), ry, int(lbl_w), int(row_h - Cm(0.05)), fill=row_bg)
        txt(sl, dim, int(mx + Cm(0.25)), ry + int(Cm(0.15)),
            int(lbl_w - Cm(0.4)), int(row_h - Cm(0.25)),
            size=10, color=WHITE)
        for j, sc in enumerate(row_sc):
            ax    = int(mx + lbl_w + j * col_w)
            fc    = score_fill[sc]
            tc    = score_text[sc]
            rect(sl, ax, ry, int(col_w - Cm(0.06)), int(row_h - Cm(0.05)), fill=row_bg)
            # Score bubble
            bsz = Cm(1.15)
            bx  = int(ax + (col_w - bsz) / 2)
            by  = int(ry  + (row_h - bsz) / 2)
            rect(sl, bx, by, int(bsz), int(bsz), fill=fc)
            txt(sl, str(sc), bx, by, int(bsz), int(bsz),
                size=14, bold=True, color=tc, align=PP_ALIGN.CENTER)

    # Legend row (score key)
    leg_y = int(my + mat_h + Cm(0.3))
    leg_x = int(mx)
    txt(sl, "Légende :", leg_x, leg_y, int(Cm(2.0)), int(Cm(0.72)),
        size=9, color=SILVER, italic=True)
    for k, (sc, label) in enumerate([(5, "Leader"), (4, "Fort"), (3, "Moyen"), (2, "Faible"), (1, "Absent")]):
        bx2 = int(mx + Cm(2.2) + k * Cm(2.4))
        rect(sl, bx2, leg_y, int(Cm(0.75)), int(Cm(0.72)), fill=score_fill[sc])
        txt(sl, str(sc), bx2, leg_y, int(Cm(0.75)), int(Cm(0.72)),
            size=10, bold=True, color=score_text[sc], align=PP_ALIGN.CENTER)
        txt(sl, label, bx2 + int(Cm(0.82)), leg_y, int(Cm(1.5)), int(Cm(0.72)),
            size=9, color=SILVER)

    # Advantage cards — right panel
    advantages = [
        ("12 services\nunifiés",  "Aucun concurrent\nn'en offre plus de 2",        GOLD),
        ("YukpoIA",               "Assistant IA intégré\n24h/24 — unique Afrique",  PURPLE),
        ("Creator Studio",        "Vidéo & visuels\ndirectement dans l'app",        GREEN),
        ("Bourse du Livre",       "Marché scolaire num.\nexclusif au Cameroun",      ORANGE),
    ]
    ax0, aw, ah = Cm(16.6), Cm(8.4), Cm(2.9)
    for i, (val, detail, col) in enumerate(advantages):
        ay = Cm(4.05) + i * (ah + Cm(0.45))
        rect(sl, int(ax0), int(ay), int(aw), int(ah), fill=NAVY2, line=col)
        rect(sl, int(ax0), int(ay), int(Cm(0.27)), int(ah), fill=col)
        txt(sl, val, int(ax0 + Cm(0.45)), int(ay + Cm(0.18)), int(aw - Cm(0.65)), int(Cm(1.2)),
            size=13, bold=True, color=col)
        txt(sl, detail, int(ax0 + Cm(0.45)), int(ay + Cm(1.4)), int(aw - Cm(0.65)), int(ah - Cm(1.5)),
            size=10, color=SILVER)

    txt(sl, "YUKPO = le seul acteur 360° du marché camerounais.",
        Cm(0.4), Cm(17.4), W - Cm(0.8), Cm(0.72),
        size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER, italic=True)


def s08(prs):
    sl = blank(prs)
    hdr(sl, "Expansion — Leader en Afrique d'ici 2030",
        "Du Cameroun à l'Afrique de l'Ouest, puis continentale : une stratégie étape par étape", 8)

    phases = [
        ("Phase 1",  "T2 – T4 2026",  "Yaoundé\n+ Douala",
         "Lancement simultané des 2 villes.\n9 services + Bourse du Livre.\n25 000 actifs/mois visés.\nBase opérationnelle solide.",
         GOLD),
        ("Phase 2",  "2027",           "Bafoussam\n+ Cameroun",
         "Extension Bafoussam, Buea,\nGaroua, Bamenda.\nDensification nationale.\nLeader camerounais confirmé.",
         CYAN),
        ("Phase 3",  "Q4 2027",        "Côte d'Ivoire\n(Abidjan)",
         "1er marché international.\nAbidjan : 5,5 M hab.\nFort taux smartphone.\nModèle YUKPO adapté.",
         GREEN),
        ("Phase 4",  "S2 2028 +",      "Sénégal &\nAfrique",
         "Dakar (3,8 M) puis expansion\npan-africaine progressive.\nObjectif : leader incontournable\nen Afrique d'ici 2030.",
         ORANGE),
    ]
    cw = Cm(5.75)
    ch = Cm(13.2)
    gap = (W - 4 * cw) / 5
    cy0 = Cm(4.05)
    for i, (ph, dates, city, desc, col) in enumerate(phases):
        cx = gap + i * (cw + gap)
        rect(sl, cx, cy0, cw, ch, fill=NAVY2, line=col)
        rect(sl, cx, cy0, cw, Cm(0.27), fill=col)
        rect(sl, cx + Cm(0.3), cy0 + Cm(0.38), Cm(0.85), Cm(0.85), fill=col)
        txt(sl, str(i + 1), cx + Cm(0.3), cy0 + Cm(0.38), Cm(0.85), Cm(0.85),
            size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(sl, ph,    cx + Cm(1.28), cy0 + Cm(0.4),  cw - Cm(1.45), Cm(0.72), size=10, bold=True, color=col)
        txt(sl, dates, cx + Cm(0.3),  cy0 + Cm(1.38), cw - Cm(0.5),  Cm(0.65), size=9, color=GOLD, italic=True)
        rect(sl, cx + Cm(0.3), cy0 + Cm(2.1), cw - Cm(0.6), Cm(0.05), fill=col)
        txt(sl, city, cx + Cm(0.3), cy0 + Cm(2.25), cw - Cm(0.5), Cm(1.4),  size=14, bold=True, color=WHITE)
        txt(sl, desc, cx + Cm(0.3), cy0 + Cm(3.75), cw - Cm(0.5), ch - Cm(4.0), size=10, color=SILVER)

    rect(sl, Cm(0.4), cy0 + ch + Cm(0.3), W - Cm(0.8), Cm(1.0), fill=NAVY2, line=ORANGE)
    rect(sl, Cm(0.4), cy0 + ch + Cm(0.3), Cm(0.27), Cm(1.0), fill=ORANGE)
    txt(sl, "Vision 2030 :  YUKPO — la super-app de référence en Afrique subsaharienne",
        Cm(0.9), cy0 + ch + Cm(0.38), W - Cm(1.5), Cm(0.85),
        size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def s09(prs):
    sl = blank(prs)
    hdr(sl, "Stratégie d'Acquisition — Marketing = 1er Poste Budgétaire",
        "La Bourse du Livre est le levier de lancement stratégique du 1er semestre 2026", 9)

    # Featured block: Bourse du Livre
    rect(sl, Cm(0.4), Cm(4.0), W - Cm(0.8), Cm(3.8), fill=NAVY2, line=ORANGE)
    rect(sl, Cm(0.4), Cm(4.0), W - Cm(0.8), Cm(0.27), fill=ORANGE)
    rect(sl, Cm(0.4), Cm(4.0), Cm(0.27), Cm(3.8), fill=ORANGE)
    txt(sl, "LEVIER CLÉ — BOURSE DU LIVRE  (T2 – T3 2026)",
        Cm(0.9), Cm(4.38), Cm(14.0), Cm(0.72), size=13, bold=True, color=ORANGE)
    txt(sl,
        "La rentrée scolaire mobilise des millions de familles camerounaises chaque année.\n"
        "YUKPO propose la première bourse numérique du livre scolaire : achat, revente, troc de livres "
        "et fournitures entre particuliers et librairies.\n"
        "Ce service unique crée une adoption massive, rapide et organique dès l'ouverture.",
        Cm(0.9), Cm(5.1), Cm(14.0), Cm(2.5), size=11, color=SILVER)
    rect(sl, Cm(15.1), Cm(4.15), Cm(9.8), Cm(3.5), fill=NAVY, line=ORANGE)
    txt(sl, "Impact prévu T2–T3 2026",
        Cm(15.3), Cm(4.35), Cm(9.4), Cm(0.72), size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    for j, (v, l) in enumerate([("50 000+", "Familles ciblées"), ("×3", "Boost inscriptions"), ("Viral", "Bouche à oreille")]):
        txt(sl, v + "   " + l, Cm(15.3), Cm(5.25) + j * Cm(0.88), Cm(9.4), Cm(0.82),
            size=11, color=WHITE, align=PP_ALIGN.CENTER)

    channels = [
        ("Présence Digitale",
         "Facebook, WhatsApp, TikTok\ncampagnes hyper-ciblées.\nSEO Google Maps local.\nInfluenceurs & micro-influenceurs.",
         "60 % inscriptions estimées", GOLD),
        ("Ambassadeurs Terrain",
         "100+ ambassadeurs formés\ndans quartiers, marchés,\nuniversités — Yaoundé + Douala.\nBouche-à-oreille structuré.",
         "25 % inscriptions estimées", CYAN),
        ("Partenariats B2B",
         "Hôtels, cliniques, écoles,\nentreprises — YUKPO intégré\ncomme service institutionnel.\nContrats récurrents.",
         "15 % + fidélisation forte", GREEN),
    ]
    cw = Cm(7.55)
    ch = Cm(7.8)
    gap = (W - 3 * cw) / 4
    cy0 = Cm(8.2)
    for i, (title, body, result, col) in enumerate(channels):
        cx = gap + i * (cw + gap)
        rect(sl, cx, cy0, cw, ch, fill=NAVY2, line=col)
        rect(sl, cx, cy0, cw, Cm(0.27), fill=col)
        txt(sl, title, cx + Cm(0.35), cy0 + Cm(0.42), cw - Cm(0.55), Cm(0.78), size=12, bold=True, color=col)
        rect(sl, cx + Cm(0.35), cy0 + Cm(1.25), cw - Cm(0.7), Cm(0.05), fill=col)
        txt(sl, body, cx + Cm(0.35), cy0 + Cm(1.42), cw - Cm(0.55), Cm(4.2), size=10, color=SILVER)
        rect(sl, cx + Cm(0.35), cy0 + ch - Cm(1.65), cw - Cm(0.55), Cm(1.4),
             fill=DKGRN if col == GREEN else NAVY, line=col)
        txt(sl, result, cx + Cm(0.35), cy0 + ch - Cm(1.6), cw - Cm(0.55), Cm(1.35),
            size=12, bold=True, color=col, align=PP_ALIGN.CENTER)


def s10(prs):
    sl = blank(prs)
    hdr(sl, "Utilisation des Fonds — 120 M FCFA",
        "L'acquisition clients est le 1er investissement — le produit est déjà très avancé", 10)

    items = [
        ("Marketing & Acquisition",    45, GOLD,
         "54 M FCFA  —  Campagnes digitales, ambassadeurs, Bourse du Livre, événements de lancement Yaoundé + Douala"),
        ("Opérations & Équipe",        25, CYAN,
         "30 M FCFA  —  Recrutement des profils clés, bureaux dans les 2 villes, logistique opérationnelle"),
        ("Finalisation & Déploiement", 15, ORANGE,
         "18 M FCFA  —  Derniers réglages produit, tests charge, déploiement (app déjà très avancée : 20 000+ h dev)"),
        ("Réserve de Trésorerie",      15, GREEN,
         "18 M FCFA  —  Sécurité financière 6 mois + gestion des imprévus opérationnels"),
    ]
    bar_x     = Cm(10.5)
    bar_max_w = Cm(13.7)
    row_h     = Cm(3.5)
    y0        = Cm(4.1)
    for i, (label, pct, col, detail) in enumerate(items):
        cy = y0 + i * row_h
        rect(sl, Cm(0.4), cy, W - Cm(0.8), row_h - Cm(0.15), fill=NAVY2)
        rect(sl, Cm(0.4), cy, Cm(0.27), row_h - Cm(0.15), fill=col)
        txt(sl, label, Cm(0.9), cy + Cm(0.22), Cm(7.3), Cm(0.82), size=12, bold=True, color=WHITE)
        rect(sl, Cm(8.4), cy + Cm(0.2), Cm(1.85), Cm(0.95), fill=NAVY, line=col)
        txt(sl, str(pct) + " %", Cm(8.4), cy + Cm(0.2), Cm(1.85), Cm(0.95),
            size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        rect(sl, bar_x, cy + Cm(0.2), int(bar_max_w * pct / 100), Cm(0.95), fill=col)
        txt(sl, detail, Cm(0.9), cy + Cm(1.32), W - Cm(1.6), Cm(0.92), size=10, color=SILVER)

    rect(sl, Cm(0.4), y0 + 4 * row_h + Cm(0.12), W - Cm(0.8), Cm(1.2), fill=NAVY2, line=GOLD)
    txt(sl, "TOTAL : 120 000 000 FCFA  |  Marketing = 1er poste (45 %)  |  Runway : 18 mois  |  Rentabilité cible : T1 2028",
        Cm(0.7), y0 + 4 * row_h + Cm(0.22), W - Cm(1.4), Cm(1.0),
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ── Central financial data ────────────────────────────────────────────────────
def _fin_data():
    """
    Modèle financier prudent (hypothèse basse) — vérifications Python :

    Revenus : commissions 1-5 % (hors pharmacie) + forfaits prestataires.
    L'hypothèse basse exclut : transport longue-distance, immobilier,
    Bourse du Livre, publicité → marché réel bien supérieur.

    Charges : coûts opérationnels COMPLETS (équipe, IT/cloud, marketing,
    acquisitions partenaires services, logistique, admin, légal).
    Elles intègrent les coûts d'acquisition des services (transport, santé,
    immobilier, BdL, etc.) et croissent avec l'activité + les expansions.

    L'investissement (120M) couvre le déficit des phases de lancement et
    constitue une réserve de croissance. Les revenus générés RÉDUISENT
    progressivement le recours au capital investi.

    Garanties numériques :
      cumul_q[-2] < 0   → T4'27 encore déficitaire (cumulatif)
      cumul_q[-1] > 0   → T1'28 premier trimestre cumulatif positif
      min(cash_q) > 0   → solde trésorerie jamais négatif sur 2 ans
      sum(disb_q) == 120 → 120M entièrement programmés
    """
    periods = ["T2'26", "T3'26", "T4'26", "T1'27",
               "T2'27", "T3'27", "T4'27", "T1'28"]

    # ── Revenus trimestriels (M FCFA) ─────────────────────────────────────────
    # Commissions 1-5 % + forfaits prestataires — hypothèse basse minimisée
    # T3'26 : pic Bourse du Livre (rentrée scolaire)
    # T1'27 : hausse post-lancement Bafoussam
    # T2-T4'27 : CI (Abidjan) + densification Cameroun
    rev_q = [3, 8, 14, 22, 33, 48, 65, 88]
    # Vérification growth cohérente
    assert all(rev_q[i] < rev_q[i+1] for i in range(len(rev_q)-1)), "Revenus non croissants"

    # ── Charges trimestrielles (M FCFA) ──────────────────────────────────────
    # Incluent : salaires équipe, marketing campagnes, IT/cloud/infra,
    # partenariats services (transport, santé, BdL, immobilier),
    # logistique, juridique/admin, coûts d'acquisition marchés.
    # Croissent avec l'activité — spikes aux extensions géographiques.
    # T2-T3'26 : lancement intense Yaoundé+Douala (marketing blitz + setup)
    # T1'27 : ouverture Bafoussam (+marketing +bureau +IT)
    # T2-T4'27 : CI Abidjan = plus grande dépense (nouveau pays)
    # T1'28 : normalisation — activité auto-financée par les revenus
    chg_q = [28, 25, 22, 26, 30, 34, 38, 42]
    assert all(c > 0 for c in chg_q), "Charges doivent être positives"
    # Note : sum(chg_q) = 245M — total opex sur 2 ans
    # Financé par : 120M investment + 125M revenus générés progressivement

    # ── Résultats et cumulatifs ───────────────────────────────────────────────
    net_q = [r - c for r, c in zip(rev_q, chg_q)]
    # = [-25, -17, -8, -4, +3, +14, +27, +46]

    cr = 0; cumul_q = []
    for n in net_q:
        cr += n; cumul_q.append(cr)
    # = [-25, -42, -50, -54, -51, -37, -10, +36]
    assert cumul_q[-2] < 0,  f"T4'27 devrait être négatif : {cumul_q[-2]}"
    assert cumul_q[-1] > 0,  f"T1'28 devrait être positif : {cumul_q[-1]}"

    # ── Position de trésorerie (départ = 120M investis) ───────────────────────
    # Montre que la trésorerie n'est JAMAIS négative sur 2 ans
    cash = 120; cash_q = []
    for n in net_q:
        cash += n; cash_q.append(cash)
    # = [95, 78, 70, 66, 69, 83, 110, 156]
    assert min(cash_q) > 0, f"Trésorerie négative ! min={min(cash_q)}"

    # ── Budget d'investissement — 120M programmés par ligne ───────────────────
    # Chaque ligne : [valeurs par trimestre], total ligne, couleur
    # Somme de chaque colonne = décaissement trimestriel d'investissement
    # Somme totale = 120M exactement
    disb = [
        ("Marketing & Acquisition (45%=54M)",
         [12, 10,  8,  8,  7,  5,  3,  1], 54, GOLD),
        ("Opérations & Équipe (25%=30M)",
         [ 7,  6,  5,  5,  4,  2,  1,  0], 30, CYAN),
        ("Acq. Services & Infra (15%=18M)",
         [ 7,  5,  3,  2,  1,  0,  0,  0], 18, ORANGE),
        ("Réserve de Trésorerie (15%=18M)",
         [ 2,  4,  6,  1,  3,  1,  0,  1], 18, GREEN),
    ]
    disb_q = [sum(d[1][j] for d in disb) for j in range(len(periods))]
    assert sum(disb_q) == 120, f"Budget total ≠ 120M : {sum(disb_q)}"
    # = [28, 25, 22, 16, 15, 8, 4, 2] — décroissant après lancement

    cd_acc = 0; cumul_disb = []
    for d in disb_q:
        cd_acc += d; cumul_disb.append(cd_acc)

    # ── Actifs utilisateurs (fin de trimestre) ────────────────────────────────
    actifs_q = [2_000, 5_500, 11_000, 18_000, 28_000, 40_000, 55_000, 70_000]

    return (periods, rev_q, chg_q, net_q, cumul_q,
            cash_q, disb, disb_q, cumul_disb, actifs_q)


def _draw_table(sl, tx, ty, lw, cw, hh, rh, periods, rows_data):
    """Helper — renders the financial grid shared by s11 and s11b."""
    n_p = len(periods)
    # Header row
    rect(sl, int(tx), int(ty), int(lw), int(hh), fill=NAVY)
    txt(sl, "Indicateur", int(tx + Cm(0.2)), int(ty + Cm(0.1)),
        int(lw - Cm(0.3)), int(hh - Cm(0.18)),
        size=9, bold=True, color=SILVER, italic=True)
    for j, per in enumerate(periods):
        px   = int(tx + lw + j * cw)
        is_t = (j == n_p - 1)
        hcol = GOLD if is_t else NAVY2
        tc   = NAVY if is_t else GOLD
        rect(sl, px, int(ty), int(cw - Cm(0.06)), int(hh), fill=hcol)
        txt(sl, per, px + int(Cm(0.04)), int(ty + Cm(0.08)),
            int(cw - Cm(0.1)), int(hh - Cm(0.15)),
            size=9, bold=True, color=tc, align=PP_ALIGN.CENTER)
    # Data rows
    for i, (label, values, row_col) in enumerate(rows_data):
        ry     = int(ty + hh + i * rh)
        row_bg = NAVY2 if i % 2 == 0 else NAVY
        rect(sl, int(tx), ry, int(lw), int(rh - Cm(0.06)), fill=row_bg, line=row_col)
        rect(sl, int(tx), ry, int(Cm(0.22)), int(rh - Cm(0.06)), fill=row_col)
        txt(sl, label, int(tx + Cm(0.32)), ry + int(Cm(0.2)),
            int(lw - Cm(0.42)), int(rh - Cm(0.3)),
            size=9, bold=True, color=row_col)
        for j, val in enumerate(values):
            px     = int(tx + lw + j * cw)
            is_t   = (j == n_p - 1)
            s      = str(val)
            neg    = s.startswith("−") or s.startswith("-")
            pos    = s.startswith("+")
            if row_col in (GREEN, PURPLE):
                fc      = RED if neg else (GREEN if pos else SILVER)
                cell_bg = DKGRN if (is_t and pos) else row_bg
            else:
                fc      = GOLD if is_t else WHITE
                cell_bg = NAVY if is_t else row_bg
            rect(sl, px, ry, int(cw - Cm(0.06)), int(rh - Cm(0.06)), fill=cell_bg)
            txt(sl, s, px + int(Cm(0.04)), ry + int(Cm(0.22)),
                int(cw - Cm(0.1)), int(rh - Cm(0.36)),
                size=10, bold=(is_t and row_col in (CYAN, GREEN, PURPLE)),
                color=fc, align=PP_ALIGN.CENTER)


def s11(prs):
    """Tableau P&L trimestriel compact — 6 lignes, données Python vérifiées"""
    sl = blank(prs)
    hdr(sl, "Projection Financière — Hypothèse Basse · 2026-2028",
        "Revenus minimisés · Charges élargies (acquisitions services incluses) · Seuil cumulatif T1 2028", 11)

    (periods, rev_q, chg_q, net_q, cumul_q,
     cash_q, disb, disb_q, cumul_disb, actifs_q) = _fin_data()

    def sg(v):
        s = f"+{v}" if v >= 0 else f"−{abs(v)}"
        return s

    actifs_s  = [f"{a:,}".replace(",", " ") for a in actifs_q]
    rev_s     = [str(v) for v in rev_q]
    chg_s     = [str(v) for v in chg_q]
    net_s     = [sg(v)  for v in net_q]
    cumul_s   = [sg(v)  for v in cumul_q]
    cash_s    = [str(v) for v in cash_q]

    rows_data = [
        ("Util. actifs / fin trimestre", actifs_s, GOLD),
        ("Revenus trim. (M FCFA)",       rev_s,    CYAN),
        ("Charges trim. (M FCFA) *",     chg_s,    ORANGE),
        ("Résultat net trim. (M)",        net_s,    GREEN),
        ("Cumulatif P&L (M FCFA)",        cumul_s,  PURPLE),
        ("Trésorerie disponible (M) **",  cash_s,   SILVER),
    ]

    tx = Cm(0.4); ty = Cm(4.1); lw = Cm(5.45)
    cw = (W - tx - lw - Cm(0.4)) / len(periods)
    hh = Cm(0.88); rh = Cm(1.98)

    _draw_table(sl, tx, ty, lw, cw, hh, rh, periods, rows_data)

    n_r = len(rows_data)
    fy  = int(ty + hh + n_r * rh + Cm(0.18))
    rect(sl, int(tx), fy, int(W - Cm(0.8)), int(Cm(1.1)), fill=NAVY2, line=GOLD)
    rect(sl, int(tx), fy, int(Cm(0.22)), int(Cm(1.1)), fill=GOLD)
    note = (
        f"*  Charges incluent : acquisitions partenaires (transport, santé, immobilier, BdL...) "
        f"+ marketing, IT, équipe, logistique.  Total 2 ans = {sum(chg_q)} M FCFA "
        f"(financé par 120M invest. + {sum(rev_q)} M revenus générés).\n"
        f"**  Tréso = 120M invest. + cumulatif P&L.  Minimum = {min(cash_q)} M (T1'27) — "
        f"jamais négatif.  Hypothèse basse sans Bourse du Livre ni immobilier/transport."
    )
    txt(sl, note, int(tx + Cm(0.35)), fy + int(Cm(0.1)),
        int(W - Cm(1.4)), int(Cm(1.0)),
        size=8, color=GOLD, italic=True)


def s11b(prs):
    """Suivi des décaissements — cadence de déploiement des 120M FCFA"""
    sl = blank(prs)
    hdr(sl, "Décaissements — Cadence de Déploiement des 120 M FCFA",
        "Plan de déploiement par poste budgétaire · Marketing ↑ aux extensions · IT ↑ aux nouveaux marchés", 12)

    (periods, rev_q, chg_q, net_q, cumul_q,
     cash_q, disb, disb_q, cumul_disb, _) = _fin_data()
    restant = [120 - cc for cc in cumul_disb]

    n_p = len(periods)
    n_r = len(disb) + 3   # 4 lignes + total + cumul + restant
    tx  = Cm(0.4)
    ty  = Cm(4.1)
    lw  = Cm(5.8)
    cw  = (W - tx - lw - Cm(0.4)) / n_p
    hh  = Cm(0.9)
    rh  = Cm(1.72)

    def _hdr_row():
        rect(sl, int(tx), int(ty), int(lw), int(hh), fill=NAVY)
        txt(sl, "Ligne budgétaire", int(tx + Cm(0.2)), int(ty + Cm(0.1)),
            int(lw - Cm(0.3)), int(hh - Cm(0.18)),
            size=9, bold=True, color=SILVER, italic=True)
        for j, per in enumerate(periods):
            px   = int(tx + lw + j * cw)
            is_t = (j == 7)
            hcol = GOLD if is_t else NAVY2
            tc   = NAVY if is_t else GOLD
            rect(sl, px, int(ty), int(cw - Cm(0.06)), int(hh), fill=hcol)
            txt(sl, per, px + int(Cm(0.04)), int(ty + Cm(0.08)),
                int(cw - Cm(0.1)), int(hh - Cm(0.15)),
                size=9, bold=True, color=tc, align=PP_ALIGN.CENTER)
    _hdr_row()

    # Budget line rows
    for i, (label, vals, total, col) in enumerate(disb):
        ry     = int(ty + hh + i * rh)
        row_bg = NAVY2 if i % 2 == 0 else NAVY
        rect(sl, int(tx), ry, int(lw), int(rh - Cm(0.06)), fill=row_bg, line=col)
        rect(sl, int(tx), ry, int(Cm(0.22)), int(rh - Cm(0.06)), fill=col)
        txt(sl, label, int(tx + Cm(0.32)), ry + int(Cm(0.18)),
            int(lw - Cm(0.42)), int(rh - Cm(0.28)),
            size=9, bold=True, color=col)
        for j, v in enumerate(vals):
            px  = int(tx + lw + j * cw)
            vc  = col if v > 0 else SILVER
            rect(sl, px, ry, int(cw - Cm(0.06)), int(rh - Cm(0.06)), fill=row_bg)
            txt(sl, str(v) if v > 0 else "—",
                px + int(Cm(0.04)), ry + int(Cm(0.25)),
                int(cw - Cm(0.1)), int(rh - Cm(0.38)),
                size=11, bold=(v > 0), color=vc, align=PP_ALIGN.CENTER)

    # TOTAL row
    tr_y = int(ty + hh + len(disb) * rh)
    rect(sl, int(tx), tr_y, int(lw + n_p * cw), int(rh - Cm(0.06)), fill=NAVY2, line=GOLD)
    rect(sl, int(tx), tr_y, int(Cm(0.22)), int(rh - Cm(0.06)), fill=GOLD)
    txt(sl, "TOTAL décaissé / trimestre (M)",
        int(tx + Cm(0.32)), tr_y + int(Cm(0.18)),
        int(lw - Cm(0.42)), int(rh - Cm(0.28)),
        size=9, bold=True, color=GOLD)
    for j, v in enumerate(disb_q):
        px = int(tx + lw + j * cw)
        rect(sl, px, tr_y, int(cw - Cm(0.06)), int(rh - Cm(0.06)), fill=NAVY2)
        txt(sl, str(v), px + int(Cm(0.04)), tr_y + int(Cm(0.25)),
            int(cw - Cm(0.1)), int(rh - Cm(0.38)),
            size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # CUMUL DÉCAISSÉ row
    cr_y = int(ty + hh + (len(disb) + 1) * rh)
    rect(sl, int(tx), cr_y, int(lw + n_p * cw), int(rh - Cm(0.06)), fill=NAVY, line=CYAN)
    rect(sl, int(tx), cr_y, int(Cm(0.22)), int(rh - Cm(0.06)), fill=CYAN)
    txt(sl, "Cumulé déployé (M FCFA)",
        int(tx + Cm(0.32)), cr_y + int(Cm(0.18)),
        int(lw - Cm(0.42)), int(rh - Cm(0.28)),
        size=9, bold=True, color=CYAN)
    for j, v in enumerate(cumul_disb):
        px  = int(tx + lw + j * cw)
        pct = int(v / 1.2)
        rect(sl, px, cr_y, int(cw - Cm(0.06)), int(rh - Cm(0.06)), fill=NAVY)
        txt(sl, f"{v}M\n({pct}%)", px + int(Cm(0.04)), cr_y + int(Cm(0.15)),
            int(cw - Cm(0.1)), int(rh - Cm(0.25)),
            size=9, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # BUDGET RESTANT row
    br_y = int(ty + hh + (len(disb) + 2) * rh)
    rect(sl, int(tx), br_y, int(lw + n_p * cw), int(rh - Cm(0.06)), fill=NAVY2, line=GREEN)
    rect(sl, int(tx), br_y, int(Cm(0.22)), int(rh - Cm(0.06)), fill=GREEN)
    txt(sl, "Budget restant (M FCFA)",
        int(tx + Cm(0.32)), br_y + int(Cm(0.18)),
        int(lw - Cm(0.42)), int(rh - Cm(0.28)),
        size=9, bold=True, color=GREEN)
    for j, v in enumerate(restant):
        px   = int(tx + lw + j * cw)
        col2 = GREEN if v > 40 else (ORANGE if v > 0 else SILVER)
        rect(sl, px, br_y, int(cw - Cm(0.06)), int(rh - Cm(0.06)), fill=NAVY2)
        txt(sl, f"{v}M" if v > 0 else "0",
            px + int(Cm(0.04)), br_y + int(Cm(0.25)),
            int(cw - Cm(0.1)), int(rh - Cm(0.38)),
            size=11, bold=True, color=col2, align=PP_ALIGN.CENTER)

    # TRÉSORERIE (= 120M + cumul P&L) — clé pour l'investisseur
    csh_y = int(ty + hh + (len(disb) + 3) * rh)
    n_r   = len(disb) + 4
    rect(sl, int(tx), csh_y, int(lw + n_p * cw), int(rh - Cm(0.06)), fill=DKGRN, line=GREEN)
    rect(sl, int(tx), csh_y, int(Cm(0.22)), int(rh - Cm(0.06)), fill=GREEN)
    txt(sl, "Position trésorerie (M FCFA)",
        int(tx + Cm(0.32)), csh_y + int(Cm(0.18)),
        int(lw - Cm(0.42)), int(rh - Cm(0.28)),
        size=9, bold=True, color=GREEN)
    for j, v in enumerate(cash_q):
        px   = int(tx + lw + j * cw)
        col2 = GREEN if v > 100 else (CYAN if v > 70 else ORANGE)
        rect(sl, px, csh_y, int(cw - Cm(0.06)), int(rh - Cm(0.06)), fill=DKGRN)
        txt(sl, str(v), px + int(Cm(0.04)), csh_y + int(Cm(0.25)),
            int(cw - Cm(0.1)), int(rh - Cm(0.38)),
            size=11, bold=True, color=col2, align=PP_ALIGN.CENTER)

    # Footer
    fy = int(ty + hh + n_r * rh + Cm(0.18))
    rect(sl, int(tx), fy, int(W - Cm(0.8)), int(Cm(0.85)), fill=NAVY2, line=GOLD)
    txt(sl,
        f"120M FCFA programmés sur 8 trimestres  |  Pic de décaissement T2'26 : {disb_q[0]}M (lancement)  |  "
        f"T1'28 : {disb_q[-1]}M seulement — activité auto-financée  |  "
        f"Trésorerie min. = {min(cash_q)}M (T1'27) · max = {max(cash_q)}M (T1'28)",
        int(tx + Cm(0.3)), fy + int(Cm(0.1)), int(W - Cm(1.4)), int(Cm(0.72)),
        size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def s12(prs):
    """Graphique Revenus/Charges + courbe trésorerie — aligné sur _fin_data()"""
    sl = blank(prs)
    hdr(sl, "Rentabilité & Trésorerie — Visualisation Trimestrielle",
        "Revenus vs Charges · Trésorerie jamais négative · Seuil cumulatif T1 2028 · 120M invest. préservé", 13)

    (periods, rev_q, chg_q, net_q, cumul_q,
     cash_q, disb, disb_q, cumul_disb, _) = _fin_data()

    # Graphique barres groupées : Revenus vs Charges
    cd = ChartData()
    cd.categories = periods
    cd.add_series("Revenus trimestriels (M FCFA)",   tuple(rev_q))
    cd.add_series("Charges trimestrielles (M FCFA)", tuple(chg_q))
    cd.add_series("Trésorerie disponible (M FCFA)",  tuple(cash_q))

    cs = sl.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        int(Cm(0.4)), int(Cm(4.0)),
        int(Cm(19.5)), int(Cm(13.3)), cd)
    chart = cs.chart
    chart.has_legend = False   # légende native désactivée — on dessine la nôtre en blanc

    plot = chart.plots[0]
    for idx, col in enumerate([GOLD, RED, CYAN]):
        plot.series[idx].format.fill.solid()
        plot.series[idx].format.fill.fore_color.rgb = col

    # Légende manuelle sous le graphique — texte blanc sur fond navy
    leg_y = Cm(17.0)
    leg_items = [
        ("Revenus (M FCFA)",    GOLD),
        ("Charges (M FCFA)",    RED),
        ("Trésorerie (M FCFA)", CYAN),
    ]
    leg_total_w = Cm(19.0)
    leg_item_w  = leg_total_w / len(leg_items)
    for i, (lbl, col) in enumerate(leg_items):
        lx = Cm(0.5) + i * leg_item_w
        rect(sl, int(lx), int(leg_y), int(Cm(0.55)), int(Cm(0.45)), fill=col)
        txt(sl, lbl, int(lx + Cm(0.65)), int(leg_y), int(leg_item_w - Cm(0.75)), int(Cm(0.48)),
            size=10, bold=True, color=WHITE)

    # Callout panel (right)
    co_x, co_w, co_h = Cm(20.4), Cm(4.6), Cm(13.3)
    rect(sl, co_x, Cm(4.0), co_w, co_h, fill=DKGRN, line=GREEN)
    rect(sl, co_x, Cm(4.0), co_w, Cm(0.27), fill=GREEN)
    txt(sl, "T1 2028",
        co_x + Cm(0.2), Cm(4.45), co_w - Cm(0.35), Cm(0.72),
        size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(sl, "Seuil\ncumulatif",
        co_x + Cm(0.2), Cm(5.22), co_w - Cm(0.35), Cm(1.1),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, co_x + Cm(0.4), Cm(6.45), co_w - Cm(0.8), Cm(0.06), fill=GREEN)

    kpis = [
        ("Revenus T1'28",    f"{rev_q[-1]} M",            GOLD),
        ("Charges T1'28",    f"{chg_q[-1]} M",            RED),
        ("Net trim. T1'28",  f"+{net_q[-1]} M FCFA",      GREEN),
        ("Tréso T1'28",      f"{cash_q[-1]} M FCFA",      CYAN),
        ("Cumul P&L T1'28",  f"+{cumul_q[-1]} M FCFA",    PURPLE),
    ]
    for i, (lbl, val, col) in enumerate(kpis):
        ky = Cm(6.7) + i * Cm(1.35)
        txt(sl, lbl, co_x + Cm(0.2), ky,
            co_w - Cm(0.35), Cm(0.55), size=8, color=SILVER, align=PP_ALIGN.CENTER)
        txt(sl, val, co_x + Cm(0.2), ky + Cm(0.55),
            co_w - Cm(0.35), Cm(0.75), size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        rect(sl, co_x + Cm(0.4), ky + Cm(1.28), co_w - Cm(0.8), Cm(0.04), fill=NAVY2)

    rect(sl, co_x + Cm(0.3), Cm(14.05), co_w - Cm(0.6), Cm(0.78), fill=NAVY2, line=ORANGE)
    txt(sl, "Tréso min. = 66M\n(T1'27 — jamais < 0)",
        co_x + Cm(0.3), Cm(14.08), co_w - Cm(0.6), Cm(0.75),
        size=8, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)


def s13(prs):
    sl = blank(prs)
    hdr(sl, "Valorisation — Estimation Partielle : 600+ M FCFA",
        "200 M FCFA de valeur directe (dev) · Innovation, IP et supports technologiques non valorisés", 14)

    val_blocks = [
        ("200 M\nFCFA",        "Valeur directe\nde développement", GOLD,
         "20 000 h de dev.\nà 10 000 FCFA/h\n= 200 M FCFA\nde valeur directe mesurée"),
        ("> 600 M\nFCFA",      "Valorisation totale\n(estimation partielle)", GREEN,
         "200 M code + algorithmes\n+ YukpoIA + Creator Studio\n+ IP, marque & potentiel\n→ 400+ M non encore valorisés"),
        ("Max 20 %",           "Part cédée\naux investisseurs",    CYAN,
         "120 M investis\n/ 600+ M post-money\n≤ 20 % de dilution\nÉquité fondateur préservée"),
    ]
    bw = Cm(7.55)
    bh = Cm(7.8)
    bgap = (W - 3 * bw) / 4
    by0 = Cm(4.1)
    for i, (val, title, col, detail) in enumerate(val_blocks):
        cx = bgap + i * (bw + bgap)
        rect(sl, cx, by0, bw, bh, fill=NAVY2, line=col)
        rect(sl, cx, by0, bw, Cm(0.27), fill=col)
        txt(sl, val, cx + Cm(0.3), by0 + Cm(0.4), bw - Cm(0.5), Cm(2.4),
            size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, title, cx + Cm(0.3), by0 + Cm(2.85), bw - Cm(0.5), Cm(1.1),
            size=11, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, cx + Cm(0.5), by0 + Cm(4.0), bw - Cm(1.0), Cm(0.06), fill=col)
        txt(sl, detail, cx + Cm(0.3), by0 + Cm(4.2), bw - Cm(0.5), bh - Cm(4.35),
            size=11, color=SILVER, align=PP_ALIGN.CENTER)

    # Note partielle
    note_y = by0 + bh + Cm(0.25)
    rect(sl, Cm(0.4), note_y, W - Cm(0.8), Cm(0.72), fill=NAVY2, line=ORANGE)
    rect(sl, Cm(0.4), note_y, Cm(0.22), Cm(0.72), fill=ORANGE)
    txt(sl, "Note : estimation partielle — l'innovation, les supports technologiques et le potentiel marché ne sont pas intégrés dans cette base de calcul.",
        Cm(0.75), note_y + Cm(0.08), W - Cm(1.3), Cm(0.6),
        size=9, italic=True, color=ORANGE)

    rois = [
        ("ROI cible",  "3× à 5×",           GOLD),
        ("Horizon",    "3 – 5 ans",           CYAN),
        ("Structure",  "YUKPO COMPANY SARL",  WHITE),
        ("Sortie",     "Rachat / Dividendes", GREEN),
    ]
    rw = (W - Cm(0.8)) / 4
    ry = note_y + Cm(0.9)
    rh = Cm(2.5)
    for i, (lbl, val, col) in enumerate(rois):
        rx = Cm(0.4) + i * rw
        rect(sl, rx, ry, rw - Cm(0.12), rh, fill=NAVY2, line=col)
        txt(sl, lbl, rx + Cm(0.2), ry + Cm(0.18), rw - Cm(0.5), Cm(0.6), size=10, color=SILVER)
        txt(sl, val, rx + Cm(0.2), ry + Cm(0.72), rw - Cm(0.5), rh - Cm(0.85),
            size=13, bold=True, color=col, align=PP_ALIGN.CENTER)


def s14(prs):
    """Fondateur basé sur le CV réel + profils techniques & commerciaux terrain"""
    sl = blank(prs)
    hdr(sl, "L'Équipe — Fondateur & Profils à Recruter",
        "Un fondateur aux compétences rares + une équipe terrain opérationnelle dès le lancement", 15)

    # ── Founder card (left) ───────────────────────────────────────────────────
    fc_x, fc_y = Cm(0.4), Cm(4.0)
    fc_w, fc_h = Cm(9.5), Cm(13.35)
    rect(sl, fc_x, fc_y, fc_w, fc_h, fill=NAVY2, line=GOLD)
    rect(sl, fc_x, fc_y, fc_w, Cm(0.27), fill=GOLD)

    # Avatar
    rect(sl, fc_x + Cm(2.7), fc_y + Cm(0.45), Cm(4.0), Cm(4.0), fill=NAVY, line=GOLD)
    txt(sl, "LSH", fc_x + Cm(2.7), fc_y + Cm(0.85), Cm(4.0), Cm(2.5),
        size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    txt(sl, "LELE SIAKA Hernandez",
        fc_x + Cm(0.3), fc_y + Cm(4.65), fc_w - Cm(0.5), Cm(0.85),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "CEO & Fondateur — YUKPO",
        fc_x + Cm(0.3), fc_y + Cm(5.5), fc_w - Cm(0.5), Cm(0.6),
        size=10, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
    txt(sl, "+237 697 49 06 61  /  +237 674 54 68 95",
        fc_x + Cm(0.3), fc_y + Cm(6.15), fc_w - Cm(0.5), Cm(0.6),
        size=10, color=CYAN, align=PP_ALIGN.CENTER)
    txt(sl, "lelehernandez2007@gmail.com",
        fc_x + Cm(0.3), fc_y + Cm(6.78), fc_w - Cm(0.5), Cm(0.58),
        size=9, color=SILVER, align=PP_ALIGN.CENTER)

    rect(sl, fc_x + Cm(0.6), fc_y + Cm(7.45), fc_w - Cm(1.2), Cm(0.06), fill=GOLD)

    bio = (
        "Gestion de programmes & pilotage stratégique\n"
        "CEPS-BID / BIsD-UNICEF — 5 ans :\n"
        "Coordination multi-acteurs, 6 régions, 61 districts\n"
        "Budgets bailleurs, reporting, planification\n\n"
        "Gouvernance financière & risque\n"
        "Assurances / CIMA — 6 ans :\n"
        "Conformité réglementaire, analyse de risque\n"
        "Gouvernance des données, pilotage financier\n\n"
        "Développeur Full Stack — 20 000+ h :\n"
        "Rust/Axum · React/TypeScript · PostgreSQL"
    )
    txt(sl, bio, fc_x + Cm(0.3), fc_y + Cm(7.65), fc_w - Cm(0.5), fc_h - Cm(7.85),
        size=10, color=SILVER)

    # ── 4 technical + commercial profiles (right, 2×2) ───────────────────────
    roles = [
        ("Développeur Mobile Senior",
         "Développement des apps iOS & Android.\n"
         "React Native / Expo.\n"
         "Intégration API, performance, UX mobile.",
         CYAN),
        ("Ingénieur Backend & DevOps",
         "Rust/Axum, PostgreSQL, Redis.\n"
         "Déploiement Cloud, CI/CD, sécurité.\n"
         "Scalabilité infrastructure.",
         GOLD),
        ("Commercial Terrain — Yaoundé",
         "Acquisition prestataires & boutiques.\n"
         "Démarchage quartiers, marchés.\n"
         "Formation et onboarding partenaires.",
         GREEN),
        ("Commercial Terrain — Douala",
         "Même mission sur Douala.\n"
         "Réseau B2B entreprises & hôtels.\n"
         "Ambassadeurs & animation terrain.",
         ORANGE),
    ]
    rc_w, rc_h, rc_gap = Cm(7.3), Cm(6.3), Cm(0.38)
    rc_x0, rc_y0 = Cm(10.3), Cm(4.0)

    for i, (role, desc, col) in enumerate(roles):
        cx = rc_x0 + (i % 2) * (rc_w + rc_gap)
        cy = rc_y0 + (i // 2) * (rc_h + rc_gap)
        rect(sl, cx, cy, rc_w, rc_h, fill=NAVY2, line=col)
        rect(sl, cx, cy, rc_w, Cm(0.27), fill=col)
        rect(sl, cx + rc_w - Cm(2.8), cy + Cm(0.3), Cm(2.5), Cm(0.52),
             fill=NAVY, line=col)
        txt(sl, "À recruter", cx + rc_w - Cm(2.8), cy + Cm(0.3), Cm(2.5), Cm(0.52),
            size=9, color=col, italic=True, align=PP_ALIGN.CENTER)
        txt(sl, role, cx + Cm(0.3), cy + Cm(0.92), rc_w - Cm(0.5), Cm(0.85),
            size=12, bold=True, color=WHITE)
        rect(sl, cx + Cm(0.3), cy + Cm(1.87), rc_w - Cm(0.6), Cm(0.05), fill=col)
        txt(sl, desc, cx + Cm(0.3), cy + Cm(2.05), rc_w - Cm(0.5), rc_h - Cm(2.2),
            size=11, color=SILVER)


def s15(prs):
    sl = blank(prs)
    hdr(sl, "Modalités de Financement — 3 Options Disponibles",
        "Prêt bancaire, caution bancaire ou entrée au capital — chaque option est indépendante", 16)

    modes = [
        (GOLD,  "Prêt Bancaire",
         "Montant :\n80 – 120 M FCFA",
         "Durée : 3 à 5 ans\nTaux négocié avec la banque\nRemboursement mensuel\nsur les revenus YUKPO",
         "Modalités clés",
         "Garantie par le chiffre d'affaires\nprévu et les actifs de la société.\nOption idéale pour\npréserver la structure du capital."),
        (CYAN,  "Caution Bancaire",
         "Montant :\n50 – 120 M FCFA",
         "L'investisseur apporte sa\ncaution personnelle à une\nbanque, qui octroie en\nretour un prêt à YUKPO.",
         "Avantages",
         "YUKPO accède au crédit\nsans diluer son capital.\nL'investisseur valorise\nsa relation bancaire.\nRisque limité si l'activité\ncouvre les échéances."),
        (GREEN, "Entrée dans le Capital",
         "Part max : 20 %\n(16,7 % pour 120 M)",
         "Valorisation : > 600 M FCFA\n20 000+ heures dev.\nPost-money : 720 M FCFA\nROI cible : 3× à 5×",
         "Conditions",
         "Structure : YUKPO COMPANY SARL\nHorizon sortie : 3 à 5 ans\nReporting trimestriel\nSiège observateur optionnel"),
    ]
    cw = Cm(7.55)
    ch = Cm(12.5)
    cgap = (W - 3 * cw) / 4
    cy0 = Cm(4.05)
    for i, (col, title, value, d1, d1t, d2) in enumerate(modes):
        cx = cgap + i * (cw + cgap)
        rect(sl, cx, cy0, cw, ch, fill=NAVY2, line=col)
        rect(sl, cx, cy0, cw, Cm(0.27), fill=col)
        txt(sl, title, cx + Cm(0.33), cy0 + Cm(0.4), cw - Cm(0.5), Cm(0.82),
            size=13, bold=True, color=col)
        txt(sl, value, cx + Cm(0.33), cy0 + Cm(1.25), cw - Cm(0.5), Cm(1.5),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, cx + Cm(0.4), cy0 + Cm(2.82), cw - Cm(0.8), Cm(0.05), fill=col)
        txt(sl, d1, cx + Cm(0.33), cy0 + Cm(3.0), cw - Cm(0.5), Cm(3.0), size=10, color=SILVER)
        rect(sl, cx + Cm(0.4), cy0 + Cm(6.1), cw - Cm(0.8), Cm(0.05), fill=col)
        txt(sl, d1t, cx + Cm(0.33), cy0 + Cm(6.25), cw - Cm(0.5), Cm(0.62),
            size=10, bold=True, color=col)
        txt(sl, d2, cx + Cm(0.33), cy0 + Cm(6.9), cw - Cm(0.5), ch - Cm(7.1), size=10, color=SILVER)

    by = cy0 + ch + Cm(0.35)
    rect(sl, Cm(0.4), by, W - Cm(0.8), Cm(1.75), fill=NAVY2, line=GOLD)
    rect(sl, Cm(0.4), by, Cm(0.27), Cm(1.75), fill=GOLD)
    txt(sl, "Combinaison possible",
        Cm(0.9), by + Cm(0.18), Cm(5.5), Cm(0.72), size=12, bold=True, color=GOLD)
    txt(sl, "Ces trois options ne sont pas exclusives. Selon votre profil et vos contraintes, "
        "vous pouvez opter pour l'une, l'autre, ou une combinaison personnalisée. "
        "Toute proposition sera étudiée avec ouverture et flexibilité.",
        Cm(6.5), by + Cm(0.15), W - Cm(7.1), Cm(1.55), size=11, color=SILVER)


def s16(prs):
    sl = blank(prs)
    hdr(sl, "Feuille de Route — 36 Mois",
        "Des jalons concrets et géographiquement orientés vers le leadership africain", 17)

    milestones = [
        ("T2 2026",  "Lancement\nYaoundé+Douala", "9 services + Bourse\ndu Livre actifs",         GOLD),
        ("T3 2026",  "Bourse\ndu Livre",           "Levier scolaire :\n50 000+ familles ciblées",  ORANGE),
        ("T4 2026",  "30 000\nutilisateurs",        "Croissance organique\ncampagnes actives",       CYAN),
        ("T1 2027",  "Bafoussam\n+ densif.",        "3e ville + réseau\nnational camerounais",       GOLD),
        ("Q4 2027",  "Côte\nd'Ivoire",             "Abidjan — 1er marché\ninternational",            GREEN),
        ("T1 2028",  "Rentabilité\nNette",          "+27 M FCFA/mois\nmeilleur des cas",             GREEN),
        ("S2 2028+", "Sénégal &\nAfrique",          "Dakar + leadership\npan-africain 2030",          ORANGE),
    ]
    n = len(milestones)
    x0, x_end = Cm(1.0), W - Cm(1.0)
    step   = (x_end - x0) / (n - 1)
    axis_y = Cm(10.3)
    cw, ch = Cm(3.1), Cm(5.2)
    rect(sl, x0, axis_y - Cm(0.08), x_end - x0, Cm(0.16), fill=GOLD)

    for i, (date, title, detail, col) in enumerate(milestones):
        cx_mid  = x0 + i * step
        cx_card = cx_mid - cw / 2
        above   = (i % 2 == 0)
        rect(sl, int(cx_mid - Cm(0.52)), int(axis_y - Cm(0.52)),
             int(Cm(1.04)), int(Cm(1.04)), fill=col)
        txt(sl, str(i + 1), int(cx_mid - Cm(0.52)), int(axis_y - Cm(0.52)),
            int(Cm(1.04)), int(Cm(1.04)), size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        cy = (axis_y - Cm(0.52) - ch - Cm(0.25)) if above else (axis_y + Cm(0.52) + Cm(0.25))
        rect(sl, int(cx_card), int(cy), int(cw), int(ch), fill=NAVY2, line=col)
        rect(sl, int(cx_card), int(cy), int(cw), int(Cm(0.22)), fill=col)
        txt(sl, date, int(cx_card + Cm(0.18)), int(cy + Cm(0.3)),
            int(cw - Cm(0.32)), int(Cm(0.55)), size=8, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, title, int(cx_card + Cm(0.18)), int(cy + Cm(0.88)),
            int(cw - Cm(0.32)), int(Cm(1.35)), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, detail, int(cx_card + Cm(0.18)), int(cy + Cm(2.28)),
            int(cw - Cm(0.32)), int(ch - Cm(2.45)), size=9, color=SILVER, align=PP_ALIGN.CENTER)


def s17(prs):
    sl = blank(prs)
    bg(sl, NAVY)
    rect(sl, 0, 0, W, Cm(0.45), fill=GOLD)
    rect(sl, 0, H - Cm(0.45), W, Cm(0.45), fill=GOLD)
    rect(sl, 0, Cm(0.45), Cm(0.22), H - Cm(0.9), fill=GOLD)
    rect(sl, W - Cm(0.22), Cm(0.45), Cm(0.22), H - Cm(0.9), fill=GOLD)

    rect(sl, Cm(9.5), Cm(0.6), Cm(6.4), Cm(3.2), fill=WHITE)
    logo_img(sl, Cm(9.5), Cm(0.6), Cm(6.4), Cm(3.2))

    txt(sl, "Investissons Ensemble dans l'Avenir Numérique de l'Afrique",
        Cm(1.0), Cm(3.8), W - Cm(2.0), Cm(2.0),
        size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rect(sl, Cm(4.5), Cm(5.85), W - Cm(9.0), Cm(0.08), fill=GOLD)
    txt(sl, "CONNECT. CREATE. SOLVE.",
        Cm(1.0), Cm(6.1), W - Cm(2.0), Cm(0.75),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    contacts = [
        ("Email",        "lelehernandez2007@gmail.com"),
        ("Téléphone",    "+237 697 49 06 61  /  +237 674 54 68 95"),
        ("Web",          "yukpomnang.com"),
        ("Localisation", "Yaoundé, Cameroun"),
    ]
    cw = Cm(10.8)
    ch = Cm(2.9)
    cgx = (W - 2 * cw) / 3
    cy0 = Cm(7.2)
    for i, (label, val) in enumerate(contacts):
        cx = cgx + (i % 2) * (cw + cgx)
        cy = cy0 + (i // 2) * (ch + Cm(0.4))
        rect(sl, cx, cy, cw, ch, fill=NAVY2, line=GOLD)
        rect(sl, cx, cy, cw, Cm(0.22), fill=GOLD)
        txt(sl, label, cx + Cm(0.4), cy + Cm(0.35), cw - Cm(0.65), Cm(0.65),
            size=11, bold=True, color=GOLD)
        txt(sl, val, cx + Cm(0.4), cy + Cm(1.0), cw - Cm(0.65), Cm(1.65),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "LELE SIAKA Hernandez  |  CEO & Fondateur — YUKPO COMPANY SARL",
        Cm(1.0), Cm(14.0), W - Cm(2.0), Cm(0.82),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, '"CONNECT. CREATE. SOLVE."',
        Cm(1.0), Cm(14.9), W - Cm(2.0), Cm(0.82),
        size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER, italic=True)


# ── Build ─────────────────────────────────────────────────────────────────────
def build(filename, cover_subtitle=""):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    for fn in [s01, s02, s03, s04, s05, s06, s07, s08,
               s09, s10, s11, s11b, s12, s13, s14, s15, s16, s17]:
        fn(prs, cover_subtitle) if fn == s01 else fn(prs)
    prs.save(filename)
    print("Saved: " + filename + "  (" + str(len(prs.slides)) + " slides)")


if __name__ == "__main__":
    base = r"c:\Users\23767\yukpomnang2"
    build(os.path.join(base, "Yukpo_Investisseur_Presentation.pptx"),           "")
    build(os.path.join(base, "Yukpo_Investisseur_Presentation_Boardroom.pptx"), "| BOARDROOM EDITION")
    build(os.path.join(base, "Yukpo_Investisseur_Presentation_Finale.pptx"),    "| VERSION FINALE")
    print("Done — 3 présentations reconstruites.")
