#!/usr/bin/env python3
"""
document_generator.py — Génère un fichier PPTX, DOCX ou PDF depuis un outline JSON (stdin).
Écrit les octets bruts du fichier sur stdout.
Appelé par document_generation_service.rs (YukpoIA).
"""

import sys
import json
import io

# ─── Couleurs par thème ────────────────────────────────────────────────────────
THEMES = {
    "blue":   {"primary": (0x1A, 0x73, 0xE8), "secondary": (0x34, 0xA8, 0x53), "accent": (0xFB, 0xBC, 0x04), "bg": (0xFF, 0xFF, 0xFF), "text": (0x20, 0x21, 0x24)},
    "green":  {"primary": (0x0F, 0x9D, 0x58), "secondary": (0x18, 0x72, 0xD6), "accent": (0xFF, 0x63, 0x00), "bg": (0xFF, 0xFF, 0xFF), "text": (0x20, 0x21, 0x24)},
    "purple": {"primary": (0x6C, 0x3F, 0xC2), "secondary": (0xE9, 0x1E, 0x63), "accent": (0xFF, 0xD7, 0x00), "bg": (0xFF, 0xFF, 0xFF), "text": (0x20, 0x21, 0x24)},
    "orange": {"primary": (0xF5, 0x7C, 0x00), "secondary": (0xE5, 0x39, 0x35), "accent": (0x0D, 0x47, 0xA1), "bg": (0xFF, 0xFF, 0xFF), "text": (0x20, 0x21, 0x24)},
    "dark":   {"primary": (0x20, 0x21, 0x24), "secondary": (0x1A, 0x73, 0xE8), "accent": (0xFB, 0xBC, 0x04), "bg": (0x12, 0x13, 0x14), "text": (0xE8, 0xEA, 0xED)},
    "light":  {"primary": (0x5F, 0x6B, 0x7C), "secondary": (0x17, 0x65, 0xC5), "accent": (0xFF, 0x63, 0x00), "bg": (0xF8, 0xF9, 0xFA), "text": (0x20, 0x21, 0x24)},
}


def get_theme(name):
    return THEMES.get(name, THEMES["blue"])


# ─── PPTX ─────────────────────────────────────────────────────────────────────

def rgb(t):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    return RGBColor(*t)


def pptx_set_bg(slide, color_tuple):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color_tuple)


def pptx_add_text_box(slide, text, left, top, width, height, font_size, bold=False, color=(0x20, 0x21, 0x24), wrap=True, align_center=False):
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return txBox


def generate_pptx(outline):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    theme = get_theme(outline.get("theme", "blue"))
    pr_color = theme["primary"]
    bg_color = theme["bg"]
    txt_color = theme["text"]
    acc_color = theme["accent"]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height
    MARGIN = Inches(0.6)

    blank_layout = prs.slide_layouts[6]  # blank

    slides_data = outline.get("slides") or outline.get("pages") or []
    title_text = outline.get("title", "Présentation")
    subtitle_text = outline.get("subtitle", "")
    author_text = outline.get("author", "")

    # ── Slide de titre ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    pptx_set_bg(slide, pr_color)
    # Bande décorative
    from pptx.util import Emu
    rect = slide.shapes.add_shape(1, Inches(0), Inches(6.2), W, Inches(1.3))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(*acc_color)
    rect.line.fill.background()

    pptx_add_text_box(slide, title_text, MARGIN, Inches(1.8), W - 2*MARGIN, Inches(2.5),
                      font_size=40, bold=True, color=(255,255,255), align_center=True)
    if subtitle_text:
        pptx_add_text_box(slide, subtitle_text, MARGIN, Inches(4.4), W - 2*MARGIN, Inches(0.8),
                          font_size=22, bold=False, color=(220,220,220), align_center=True)
    if author_text:
        pptx_add_text_box(slide, author_text, MARGIN, Inches(5.4), W - 2*MARGIN, Inches(0.6),
                          font_size=16, bold=False, color=(200,200,200), align_center=True)

    # ── Slides de contenu ──────────────────────────────────────────────────────
    for sl in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        pptx_set_bg(slide, bg_color)

        # Bande titre
        hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.1))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = RGBColor(*pr_color)
        hdr.line.fill.background()

        sl_title = sl.get("title", "")
        pptx_add_text_box(slide, sl_title, MARGIN, Inches(0.12), W - 2*MARGIN, Inches(0.9),
                          font_size=26, bold=True, color=(255,255,255))

        layout = sl.get("layout", "content")
        body_top = Inches(1.25)
        body_h = H - body_top - MARGIN

        if layout == "kpi":
            kpis = sl.get("kpis", [])
            n = len(kpis) or 1
            kpi_w = (W - 2*MARGIN) / n
            for i, kpi in enumerate(kpis):
                x = MARGIN + i * kpi_w
                # Carte KPI
                card = slide.shapes.add_shape(1, x + Inches(0.1), body_top + Inches(0.2),
                                              kpi_w - Inches(0.2), Inches(2.8))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*pr_color)
                card.line.fill.background()
                val = kpi.get("value", "")
                label = kpi.get("label", "")
                trend = kpi.get("trend", "")
                pptx_add_text_box(slide, val, x + Inches(0.2), body_top + Inches(0.5),
                                  kpi_w - Inches(0.4), Inches(1.2),
                                  font_size=38, bold=True, color=(255,255,255), align_center=True)
                pptx_add_text_box(slide, label, x + Inches(0.2), body_top + Inches(1.7),
                                  kpi_w - Inches(0.4), Inches(0.6),
                                  font_size=16, bold=False, color=(220,220,220), align_center=True)
                if trend:
                    pptx_add_text_box(slide, trend, x + Inches(0.2), body_top + Inches(2.3),
                                      kpi_w - Inches(0.4), Inches(0.5),
                                      font_size=18, bold=True, color=acc_color, align_center=True)

        elif layout == "two_column":
            mid = W / 2
            left_text = sl.get("left", "")
            right_text = sl.get("right", "")
            # Séparateur vertical
            sep = slide.shapes.add_shape(1, mid - Inches(0.02), body_top, Inches(0.04), body_h)
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(*pr_color)
            sep.line.fill.background()
            pptx_add_text_box(slide, left_text, MARGIN, body_top, mid - MARGIN - Inches(0.1),
                              body_h, font_size=16, color=txt_color)
            pptx_add_text_box(slide, right_text, mid + Inches(0.1), body_top,
                              mid - MARGIN - Inches(0.1), body_h, font_size=16, color=txt_color)

        elif layout == "table":
            tbl_data = sl.get("table", {})
            headers = tbl_data.get("headers", [])
            rows = tbl_data.get("rows", [])
            if headers and rows:
                _add_table_to_slide(slide, headers, rows, MARGIN, body_top,
                                    W - 2*MARGIN, min(body_h, Inches(5)),
                                    pr_color, txt_color, bg_color)

        else:  # content / title / blank
            content = sl.get("content", "")
            bullets = sl.get("bullets", [])
            y = body_top
            if content:
                pptx_add_text_box(slide, content, MARGIN, y, W - 2*MARGIN, Inches(1.8),
                                  font_size=18, color=txt_color)
                y += Inches(1.9)
            for b in bullets:
                bx = slide.shapes.add_textbox(MARGIN + Inches(0.2), y, W - 2*MARGIN - Inches(0.4), Inches(0.55))
                tf = bx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "• " + b
                run.font.size = Pt(17)
                run.font.color.rgb = RGBColor(*txt_color)
                y += Inches(0.6)
                if y + Inches(0.6) > H - MARGIN:
                    break

        # Table optionnelle en bas d'une slide content
        if layout not in ("table", "kpi", "two_column"):
            tbl_data = sl.get("table")
            if tbl_data:
                headers = tbl_data.get("headers", [])
                rows = tbl_data.get("rows", [])
                if headers and rows:
                    _add_table_to_slide(slide, headers, rows, MARGIN, H - Inches(3.2),
                                        W - 2*MARGIN, Inches(2.8), pr_color, txt_color, bg_color)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_table_to_slide(slide, headers, rows, left, top, width, height, pr_color, txt_color, bg_color):
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor

    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    tbl.first_row = True

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*pr_color)
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.text = h
        run.font.bold = True
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(255, 255, 255)

    for ri, row in enumerate(rows):
        alt = (ri % 2 == 1)
        for ci, val in enumerate(row[:n_cols]):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            if alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 244, 255)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.text = str(val)
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(*txt_color)


# ─── DOCX ─────────────────────────────────────────────────────────────────────

def generate_docx(outline):
    from docx import Document
    from docx.shared import Pt, RGBColor as DocxRGB, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    theme = get_theme(outline.get("theme", "blue"))
    pr_color = theme["primary"]

    doc = Document()

    # Marges
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    # Titre
    title = doc.add_heading(outline.get("title", "Document"), level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.runs[0] if title.runs else title.add_run(outline.get("title", "Document"))
    run.font.color.rgb = DocxRGB(*pr_color)

    if outline.get("subtitle"):
        sub = doc.add_paragraph(outline["subtitle"])
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = sub.runs[0] if sub.runs else sub.add_run(outline["subtitle"])
        run.font.italic = True
        run.font.size = Pt(14)
        run.font.color.rgb = DocxRGB(100, 100, 100)

    if outline.get("author"):
        auth = doc.add_paragraph(outline["author"])
        auth.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = auth.runs[0] if auth.runs else auth.add_run(outline["author"])
        run.font.size = Pt(12)

    doc.add_paragraph()

    for sect in outline.get("sections", []):
        level = sect.get("level", 1)
        heading_text = sect.get("heading", "")
        h = doc.add_heading(heading_text, level=min(level, 3))
        if h.runs:
            h.runs[0].font.color.rgb = DocxRGB(*pr_color)

        for para in sect.get("paragraphs", []):
            p = doc.add_paragraph(para)
            p.paragraph_format.space_after = Pt(6)

        for bullet in sect.get("bullets", []):
            doc.add_paragraph(bullet, style="List Bullet")

        tbl_data = sect.get("table")
        if tbl_data:
            headers = tbl_data.get("headers", [])
            rows = tbl_data.get("rows", [])
            if headers:
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = "Table Grid"
                hdr_cells = table.rows[0].cells
                for i, h_text in enumerate(headers):
                    hdr_cells[i].text = h_text
                    run = hdr_cells[i].paragraphs[0].runs[0]
                    run.font.bold = True
                    run.font.color.rgb = DocxRGB(*pr_color)
                for ri, row in enumerate(rows):
                    cells = table.rows[ri + 1].cells
                    for ci, val in enumerate(row[:len(headers)]):
                        cells[ci].text = str(val)
                doc.add_paragraph()

        doc.add_paragraph()

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─── PDF ──────────────────────────────────────────────────────────────────────

def generate_pdf(outline):
    from fpdf import FPDF

    theme = get_theme(outline.get("theme", "blue"))
    pr_color = theme["primary"]
    acc_color = theme["accent"]
    txt_color = theme["text"]

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)

    pages = outline.get("pages") or outline.get("sections") or []
    title_text = outline.get("title", "Rapport")
    subtitle_text = outline.get("subtitle", "")
    author_text = outline.get("author", "")

    # ── Page de couverture ─────────────────────────────────────────────────────
    pdf.add_page()
    pdf.set_fill_color(*pr_color)
    pdf.rect(0, 0, 210, 297, "F")

    pdf.set_font("Helvetica", "B", 32)
    pdf.set_text_color(255, 255, 255)
    pdf.set_y(80)
    pdf.multi_cell(210, 14, title_text, align="C")

    if subtitle_text:
        pdf.set_font("Helvetica", "", 18)
        pdf.set_text_color(220, 220, 220)
        pdf.ln(6)
        pdf.multi_cell(210, 10, subtitle_text, align="C")

    if author_text:
        pdf.set_font("Helvetica", "", 13)
        pdf.set_text_color(200, 200, 200)
        pdf.set_y(240)
        pdf.multi_cell(210, 8, author_text, align="C")

    # Bande décorative bas
    pdf.set_fill_color(*acc_color)
    pdf.rect(0, 270, 210, 27, "F")

    # ── Pages de contenu ───────────────────────────────────────────────────────
    for pg in pages:
        pdf.add_page()
        pg_title = pg.get("title") or pg.get("heading", "")
        layout = pg.get("layout", "content")

        # En-tête de page
        pdf.set_fill_color(*pr_color)
        pdf.rect(0, 0, 210, 18, "F")
        pdf.set_font("Helvetica", "B", 14)
        pdf.set_text_color(255, 255, 255)
        pdf.set_xy(10, 4)
        pdf.cell(190, 10, pg_title, align="L")

        pdf.set_y(26)
        pdf.set_text_color(*txt_color)

        if layout == "kpi":
            kpis = pg.get("kpis", [])
            n = len(kpis)
            if n:
                box_w = 170 / n
                for i, kpi in enumerate(kpis):
                    x = 20 + i * (box_w + 2)
                    pdf.set_fill_color(*pr_color)
                    pdf.rect(x, 30, box_w, 45, "F")
                    pdf.set_font("Helvetica", "B", 24)
                    pdf.set_text_color(255, 255, 255)
                    pdf.set_xy(x, 38)
                    pdf.cell(box_w, 14, kpi.get("value", ""), align="C")
                    pdf.set_font("Helvetica", "", 11)
                    pdf.set_xy(x, 53)
                    pdf.cell(box_w, 8, kpi.get("label", ""), align="C")
                    trend = kpi.get("trend", "")
                    if trend:
                        pdf.set_font("Helvetica", "B", 13)
                        pdf.set_text_color(*acc_color)
                        pdf.set_xy(x, 62)
                        pdf.cell(box_w, 8, trend, align="C")

        elif layout == "table":
            tbl = pg.get("table", {})
            _pdf_table(pdf, tbl, pr_color, txt_color, acc_color)

        else:  # content
            content = pg.get("content") or pg.get("paragraphs", [])
            if isinstance(content, list):
                content = "\n".join(content)
            if content:
                pdf.set_font("Helvetica", "", 12)
                pdf.set_text_color(*txt_color)
                pdf.multi_cell(190, 7, content)
                pdf.ln(4)

            for bullet in pg.get("bullets", []):
                pdf.set_font("Helvetica", "", 12)
                pdf.set_text_color(*txt_color)
                pdf.set_x(20)
                pdf.multi_cell(175, 7, "• " + bullet)

            tbl = pg.get("table")
            if tbl:
                pdf.ln(4)
                _pdf_table(pdf, tbl, pr_color, txt_color, acc_color)

    buf = io.BytesIO()
    pdf.output(buf)
    return buf.getvalue()


def _pdf_table(pdf, tbl_data, pr_color, txt_color, acc_color):
    from fpdf import FPDF
    headers = tbl_data.get("headers", [])
    rows = tbl_data.get("rows", [])
    if not headers:
        return
    n_cols = len(headers)
    col_w = 170 / n_cols

    # Headers
    pdf.set_fill_color(*pr_color)
    pdf.set_text_color(255, 255, 255)
    pdf.set_font("Helvetica", "B", 11)
    pdf.set_x(20)
    for h in headers:
        pdf.cell(col_w, 9, str(h)[:20], border=0, align="C", fill=True)
    pdf.ln()

    # Rows
    pdf.set_font("Helvetica", "", 10)
    for ri, row in enumerate(rows):
        if ri % 2 == 0:
            pdf.set_fill_color(240, 244, 255)
        else:
            pdf.set_fill_color(255, 255, 255)
        pdf.set_text_color(*txt_color)
        pdf.set_x(20)
        for val in row[:n_cols]:
            pdf.cell(col_w, 8, str(val)[:25], border=0, align="C", fill=True)
        pdf.ln()


# ─── Entrée principale ────────────────────────────────────────────────────────

def main():
    raw = sys.stdin.buffer.read()
    outline = json.loads(raw)

    doc_type = outline.get("document_type", "pptx").lower()

    if doc_type in ("docx", "word"):
        data = generate_docx(outline)
    elif doc_type in ("pdf",):
        data = generate_pdf(outline)
    else:
        data = generate_pptx(outline)

    sys.stdout.buffer.write(data)


if __name__ == "__main__":
    main()
