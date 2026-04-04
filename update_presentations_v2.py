#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Version 2 - patch amélioré:
- Reconstruction des runs fragmentés avant substitution
- Couverture de tous les patterns résiduels
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import shutil, os, re, copy

NAVY   = RGBColor(0x0B, 0x1E, 0x42)
GOLD   = RGBColor(0xD4, 0xAF, 0x37)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x00, 0xC8, 0x7A)

SUBS = [
    ("Besoin de financement : 200 millions FCFA \u00b7 d\u00e9ploiement mai 2026 \u2013 fin 2027",
     "Besoin de financement : 120 millions FCFA \u00b7 rentabilit\u00e9 d\u00e8s d\u00e9cembre 2027"),
    ("Entr\u00e9e au capital possible jusqu'\u00e0 15 % (discussion encadr\u00e9e)",
     "Entr\u00e9e au capital possible jusqu'\u00e0 20 % \u2014 ticket optimis\u00e9, ROI acc\u00e9l\u00e9r\u00e9"),
    ("Entr\u00e9e au capital possible jusqu\u2019\u00e0 15 % (discussion encadr\u00e9e)",
     "Entr\u00e9e au capital possible jusqu\u2019\u00e0 20 % \u2014 ticket optimis\u00e9, ROI acc\u00e9l\u00e9r\u00e9"),
    ("200 millions FCFA sur 20 mois \u00b7 lancement effectif juillet 2026",
     "120 millions FCFA sur 20 mois \u00b7 lancement effectif juillet 2026"),
    ("D\u00e9tail du besoin \u2014 200 millions FCFA", "D\u00e9tail du besoin \u2014 120 millions FCFA"),
    ("Total\u00a0: 200 millions FCFA \u2014 hypoth\u00e8ses arrondies pour la lecture",
     "Total\u00a0: 120 millions FCFA \u2014 plan optimis\u00e9 \u00b7 rentabilit\u00e9 T4 2027"),
    ("Total : 200 millions FCFA \u2014 hypoth\u00e8ses arrondies pour la lecture",
     "Total : 120 millions FCFA \u2014 plan optimis\u00e9 \u00b7 rentabilit\u00e9 T4 2027"),
    ("Poste acquisition & notori\u00e9t\u00e9 : 118 M FCFA (enveloppe 200 M)",
     "Poste acquisition & notori\u00e9t\u00e9 : 70 M FCFA (enveloppe 120 M)"),
    ("Mai 2026 \u2013 avril 2028 (coh\u00e9rent avec l\u2019enveloppe 200 M)",
     "Mai 2026 \u2013 d\u00e9cembre 2027 \u00b7 break-even mensuel T4 2027"),
    ("mai 2026 \u2013 avril 2028", "mai 2026 \u2013 d\u00e9cembre 2027"),
    ("(coh\u00e9rent avec l\u2019enveloppe 200 M)", "(enveloppe 120 M)"),
    ("Hypoth\u00e8ses de synth\u00e8se (mai 2026 \u2013 avril 2028)",
     "Hypoth\u00e8ses de synth\u00e8se (mai 2026 \u2013 d\u00e9cembre 2027)"),
    ("Sur les 8 trimestres : revenus cumul\u00e9s \u2248 168 M ; charges d\u2019activit\u00e9 \u2248 200 M (coh\u00e9rent avec les deux bilans annuels). Hypoth\u00e8se de revenu conservatrice : commissions moyennes 8\u201312 % selon segment, panier progressif avec densit\u00e9 utilisateurs.",
     "Sur 18 mois : revenus cumul\u00e9s \u2248 45 M FCFA \u00b7 charges 120 M \u00b7 break-even mensuel T4 2027. Commissions 8-12 % selon segment. 2028 projet\u00e9 : 90-110 M FCFA revenus annuels."),
    ("d\u00e9passer le seuil de rentabilit\u00e9 en 2\u1d49 ann\u00e9e",
     "atteindre la rentabilit\u00e9 mensuelle d\u00e8s d\u00e9cembre 2027"),
    ("peut \u00eatre rentable sur cette base", "est rentable d\u00e8s T4 2027 (d\u00e9cembre 2027)"),
    ("\u00c0 la fin de la 2\u1d49 ann\u00e9e compl\u00e8te, les revenus d\u00e9passent les charges : l\u2019entreprise peut \u00eatre rentable sur cette base (hors nouveaux tours de financement de croissance).",
     "Break-even mensuel atteint en d\u00e9cembre 2027 :"),
    ("Ordres de grandeur pour discussion \u2014 \u00e0 confirmer avec la comptabilit\u00e9 de gestion.",
     "\u2022 Revenus/mois T4 2027 : ~6,5 M FCFA  |  Charges r\u00e9currentes : ~4,5 M FCFA\n\u2022 Profit mensuel d\u00e8s d\u00e9cembre 2027 : +2 M FCFA\n\u2022 2028 : premi\u00e8re ann\u00e9e pleinement rentable \u00b7 90-110 M FCFA projet\u00e9s"),
    ("42 millions FCFA sur 200 \u2014 charges salariales croissantes en 2\u1d49 ann\u00e9e",
     "26 millions FCFA sur 120 \u2014 mont\u00e9e en charge ma\u00eetris\u00e9e"),
    ("Budget global tech (18 M sur 200 M) : les d\u00e9caissements augmentent avec les utilisateurs (cloud, monitoring, mises \u00e0 jour).",
     "Budget global tech (14 M sur 120 M) : infra cloud scalable \u2014 co\u00fbts index\u00e9s sur la croissance r\u00e9elle."),
    ("Comment financer les 200 millions ?", "Comment financer les 120 millions ?"),
    ("plafonn\u00e9e \u00e0 15 % pour encadrer la dilution",
     "jusqu'\u00e0 20 % \u2014 meilleur ratio risque/rendement pour l\u2019investisseur"),
    ("Juil. 2027 \u2013 juin 2028 \u00b7 acc\u00e9l\u00e9ration des revenus",
     "Juil. 2027 \u2013 d\u00e9c. 2027 \u00b7 acc\u00e9l\u00e9ration et rentabilit\u00e9 mensuelle"),
    ("revenus cumul\u00e9s \u2248 168 M", "revenus cumul\u00e9s \u2248 45 M"),
    ("charges d\u2019activit\u00e9 \u2248 200 M", "charges \u2248 120 M"),
    # Generic patterns - LAST
    ("200 millions FCFA", "120 millions FCFA"),
    ("200 millions", "120 millions"),
    ("enveloppe 200 M", "enveloppe 120 M"),
    ("sur 200 M", "sur 120 M"),
    ("sur 200 \u2014", "sur 120 \u2014"),
    ("les 200 M", "les 120 M"),
    (" 200 M", " 120 M"),
    ("200 M\n", "120 M\n"),
    ("200M", "120M"),
    ("118 M FCFA", "70 M FCFA"),
    ("118 M", "70 M"),
    ("42 millions FCFA", "26 millions FCFA"),
    ("42 millions", "26 millions"),
    (" 42 M", " 26 M"),
    ("18 M sur 200", "14 M sur 120"),
    ("(18 M", "(14 M"),
    # Extra patterns for residual cases
    ("jusqu\u2019\u00e0 15\u00a0%", "jusqu\u2019\u00e0 20\u00a0%"),
    ("jusqu'\u00e0 15\u00a0%", "jusqu'\u00e0 20\u00a0%"),
    ("jusqu'\u00e0 15 %", "jusqu'\u00e0 20 %"),
    ("jusqu\u2019\u00e0 15 %", "jusqu\u2019\u00e0 20 %"),
    ("15 % (discussion encadr\u00e9e)", "20 % \u2014 ticket optimis\u00e9"),
    ("15\u00a0% (discussion encadr\u00e9e)", "20\u00a0% \u2014 ticket optimis\u00e9"),
]

GOLD_RE = re.compile(
    r'(120\s*M|120\s*millions|70\s*M|26\s*M|14\s*M|10\s*M\s*FCFA|'
    r'6[,\.]\s*5\s*M|45\s*M|90[-\u2013]110\s*M|\+2\s*M|20\s*%|break-even|'
    r'rentabilit[e\u00e9]\s*T4)', re.IGNORECASE
)

def apply_subs(text):
    for old, new in SUBS:
        text = text.replace(old, new)
    return text


def merge_runs_in_paragraph(para_el):
    """
    Merge all <a:r> children in a paragraph into a single run
    to avoid split-run substitution failures.
    Keep the formatting of the first run.
    """
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    runs = para_el.findall(f'{{{ns}}}r')
    if len(runs) < 2:
        return
    # Collect full text
    full_text = ''.join(
        (r.find(f'{{{ns}}}t').text or '') for r in runs
        if r.find(f'{{{ns}}}t') is not None
    )
    # Apply substitutions on the merged text
    full_text = apply_subs(full_text)
    # Keep first run, set its text, remove others
    first_t = runs[0].find(f'{{{ns}}}t')
    if first_t is not None:
        first_t.text = full_text
    for r in runs[1:]:
        para_el.remove(r)


def patch_xml(slide_el):
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    # First pass: merge runs per paragraph for reliable substitution
    for para in slide_el.iter(f'{{{ns}}}p'):
        merge_runs_in_paragraph(para)
    # Second pass: apply subs on any remaining <a:t>
    for t in slide_el.iter(f'{{{ns}}}t'):
        if t.text:
            t.text = apply_subs(t.text)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def fix_text_colors(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text:
                    continue
                try:
                    fc = run.font.color
                    make_white = False
                    if fc.type is None:
                        make_white = True
                    elif str(fc.type) == 'RGB (1)' and hasattr(fc, 'rgb') and fc.rgb:
                        r, g, b = fc.rgb.red, fc.rgb.green, fc.rgb.blue
                        if (r + g + b) / 3 < 130:
                            make_white = True
                    if make_white:
                        if GOLD_RE.search(run.text):
                            fc.rgb = GOLD
                            run.font.bold = True
                        else:
                            fc.rgb = WHITE
                except Exception:
                    pass


def verify_file(path):
    prs = Presentation(path)
    remaining = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                issues = []
                if ' 200 M' in txt or '200M' in txt or '200 millions' in txt:
                    issues.append('200M residual')
                if '15 %' in txt or '15\u00a0%' in txt:
                    issues.append('15% residual')
                if '118 M' in txt or '42 millions' in txt:
                    issues.append('old figures residual')
                if issues:
                    remaining.append(f"  Slide {si+1} [{shape.name}]: {issues} -> {repr(txt[:150])}")
    return remaining


def process(path):
    print(f'\nProcessing {path}...')
    backup = path.replace('.pptx', '_BACKUP.pptx')
    if not os.path.exists(backup):
        shutil.copy2(path, backup)
        print(f'  Backup created: {backup}')
    else:
        print(f'  Backup already exists: {backup}')

    prs = Presentation(path)
    print(f'  Slides: {len(prs.slides)}')
    for slide in prs.slides:
        patch_xml(slide._element)
        set_bg(slide)
        fix_text_colors(slide)

    prs.save(path)
    print(f'  Saved: {path}')

    # Verification
    print('  --- Verifying residuals ---')
    issues = verify_file(path)
    if issues:
        print('  REMAINING ISSUES:')
        for i in issues:
            print(i)
    else:
        print('  OK - No 200M/15% residuals found.')

    # Show slide 1
    prs2 = Presentation(path)
    print('  --- Slide 1 content ---')
    for shape in prs2.slides[0].shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                print(f'    [{shape.name}]: {repr(t[:200])}')

    # Show slide 32 or last
    target_idx = 31 if len(prs2.slides) >= 32 else len(prs2.slides) - 1
    print(f'  --- Slide {target_idx+1} content ---')
    for shape in prs2.slides[target_idx].shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                print(f'    [{shape.name}]: {repr(t[:200])}')


BASE = 'C:/Users/23767/yukpomnang2'
files = [
    f'{BASE}/Yukpo_Investisseur_Presentation.pptx',
    f'{BASE}/Yukpo_Investisseur_Presentation_Boardroom.pptx',
    f'{BASE}/Yukpo_Investisseur_Presentation_Finale.pptx',
]

for f in files:
    if os.path.exists(f):
        process(f)
    else:
        print(f'NOT FOUND: {f}')

print('\nDone!')
