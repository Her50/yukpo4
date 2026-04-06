#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génération automatique de la formation PowerPoint:
"Remise à Niveau en Gestion de Sinistres — Toutes Branches"
Automobile | Incendie | Transports (TFT/A/M) | Santé | Individuel Accident
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PALETTE COULEURS ────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
BLUE   = RGBColor(0x2E, 0x86, 0xAB)
ORANGE = RGBColor(0xF1, 0x8F, 0x01)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
TEAL   = RGBColor(0x16, 0xA0, 0x85)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)

# ── DIMENSIONS 16:9 ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]   # layout entièrement vide


# ══════════════════════════════════════════════════════════════════════════════
# UTILITAIRES
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp


def txb(slide, text, x, y, w, h, sz=16, color=WHITE, bold=False,
        italic=False, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text         = text
    r.font.size    = Pt(sz)
    r.font.color.rgb = color
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.name    = name
    return tb


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def bullets(slide, items, x, y, w, h, sz=14, color=DGRAY, prefix="▸  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = prefix + item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def rich_bullets(slide, items, x, y, w, h):
    """items = list of dicts: {text, level, bold, color, size}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, str):
            item = {"text": item}
        lvl   = item.get("level", 0)
        text  = item.get("text", "")
        bold  = item.get("bold", False)
        col   = item.get("color", DGRAY)
        sz    = item.get("size", 16 if lvl == 0 else 14)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7 if lvl == 0 else 3)
        r = p.add_run()
        pfx = "◆  " if lvl == 0 else "      ▸  "
        r.text         = pfx + text
        r.font.size    = Pt(sz)
        r.font.color.rgb = col
        r.font.bold    = bold
        r.font.name    = "Calibri"
    return tb


# ══════════════════════════════════════════════════════════════════════════════
# TYPES DE SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs, title, subtitle, date="Avril 2026", version="Version 1.0"):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)

    # Fond gauche navy
    add_rect(slide, 0, 0, Inches(5.1), H, fill=NAVY)
    # Liseré orange vertical
    add_rect(slide, Inches(5.1), 0, Inches(0.13), H, fill=ORANGE)
    # Bande basse orange
    add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ORANGE)

    # Texte gauche
    txb(slide, "FORMATION\nPROFESSIONNELLE\nEN ASSURANCE",
        Inches(0.35), Inches(0.8), Inches(4.4), Inches(1.8),
        sz=18, bold=True, align=PP_ALIGN.CENTER, name="Calibri Light")

    txb(slide, "GESTION\nDES SINISTRES",
        Inches(0.3), Inches(2.7), Inches(4.55), Inches(1.9),
        sz=42, bold=True, align=PP_ALIGN.CENTER, name="Calibri")

    add_rect(slide, Inches(0.5), Inches(4.7), Inches(4.1), Inches(0.06), fill=ORANGE)

    txb(slide, "De la déclaration au règlement\nToutes branches confondues",
        Inches(0.3), Inches(4.85), Inches(4.55), Inches(1.2),
        sz=14, italic=True, align=PP_ALIGN.CENTER, color=RGBColor(0xCC, 0xDD, 0xEE))

    txb(slide, f"{date}  |  {version}",
        Inches(0.3), H - Inches(0.65), Inches(4.55), Inches(0.5),
        sz=11, align=PP_ALIGN.CENTER, color=RGBColor(0x99, 0xBB, 0xDD))

    # Nom du formateur (zone navy, bas gauche)
    txb(slide, "Formateur :", Inches(0.35), H - Inches(1.35), Inches(4.4), Inches(0.38),
        sz=11, color=ORANGE, bold=False, align=PP_ALIGN.CENTER)
    txb(slide, "TALOM Eric", Inches(0.35), H - Inches(1.05), Inches(4.4), Inches(0.48),
        sz=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, name="Calibri")

    # Zone droite — titre principal
    txb(slide, title,
        Inches(5.5), Inches(1.5), Inches(7.5), Inches(2.6),
        sz=34, bold=True, color=NAVY, align=PP_ALIGN.LEFT, name="Calibri")

    add_rect(slide, Inches(5.5), Inches(4.2), Inches(5.5), Inches(0.07), fill=ORANGE)

    txb(slide, subtitle,
        Inches(5.5), Inches(4.4), Inches(7.4), Inches(2.4),
        sz=16, color=DGRAY, align=PP_ALIGN.LEFT)

    return slide


def slide_sommaire(prs, modules):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), W, Inches(0.07), fill=ORANGE)
    add_rect(slide, 0, H - Inches(0.15), W, Inches(0.15), fill=NAVY)
    txb(slide, "PLAN DE LA FORMATION — CONDUCTEUR",
        Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.9),
        sz=26, bold=True, align=PP_ALIGN.LEFT)

    half = (len(modules) + 1) // 2
    col1, col2 = modules[:half], modules[half:]

    for ci, col in enumerate([col1, col2]):
        cx = Inches(0.35) if ci == 0 else Inches(6.9)
        for ri, (num, title, dur) in enumerate(col):
            ry = Inches(1.45) + ri * Inches(0.83)
            bg = LGRAY if ri % 2 == 0 else WHITE
            add_rect(slide, cx, ry, Inches(6.2), Inches(0.73), fill=bg,
                     line=RGBColor(0xCC, 0xCC, 0xCC))
            add_rect(slide, cx, ry, Inches(0.65), Inches(0.73), fill=BLUE)
            txb(slide, num, cx + Inches(0.02), ry + Inches(0.1),
                Inches(0.61), Inches(0.52), sz=17, bold=True, align=PP_ALIGN.CENTER)
            txb(slide, title, cx + Inches(0.72), ry + Inches(0.1),
                Inches(4.6), Inches(0.55), sz=13, color=NAVY, bold=True)
            txb(slide, dur, cx + Inches(5.35), ry + Inches(0.2),
                Inches(0.75), Inches(0.35), sz=12, color=ORANGE, bold=True,
                align=PP_ALIGN.CENTER)
    return slide


def slide_section(prs, num, title, subtitle="", color=NAVY):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, color)
    add_rect(slide, 0, 0, Inches(0.28), H, fill=ORANGE)
    add_rect(slide, 0, H - Inches(0.15), W, Inches(0.15), fill=ORANGE)
    txb(slide, f"MODULE  {num}", Inches(0.8), Inches(1.8),
        Inches(11.7), Inches(0.75), sz=18, color=ORANGE, bold=True,
        align=PP_ALIGN.CENTER, name="Calibri Light")
    txb(slide, title, Inches(0.8), Inches(2.55), Inches(11.7), Inches(1.8),
        sz=44, bold=True, align=PP_ALIGN.CENTER, name="Calibri")
    add_rect(slide, Inches(4.2), Inches(4.55), Inches(5.0), Inches(0.08), fill=ORANGE)
    if subtitle:
        txb(slide, subtitle, Inches(1.5), Inches(4.75), Inches(10.3), Inches(1.5),
            sz=18, italic=True, align=PP_ALIGN.CENTER,
            color=RGBColor(0xCC, 0xDD, 0xEE))
    return slide


def slide_content(prs, title, items, label="", accent=BLUE):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), W, Inches(0.07), fill=ORANGE)
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=accent)
    if label:
        txb(slide, label, Inches(0.3), Inches(0.07), Inches(8), Inches(0.42),
            sz=11, color=RGBColor(0xBB, 0xCC, 0xDD), name="Calibri Light")
    txb(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.95),
        sz=24, bold=True)
    rich_bullets(slide, items, Inches(0.45), Inches(1.5), Inches(12.5), Inches(5.7))
    return slide


def slide_two_col(prs, title, t1, i1, t2, i2, c1=BLUE, c2=TEAL, label=""):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), W, Inches(0.07), fill=ORANGE)
    if label:
        txb(slide, label, Inches(0.3), Inches(0.07), Inches(8), Inches(0.42),
            sz=11, color=RGBColor(0xBB, 0xCC, 0xDD))
    txb(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.95),
        sz=24, bold=True)

    CW = Inches(6.15)
    CY = Inches(1.42)
    CH = Inches(5.85)
    for ci, (cx, t, items, col) in enumerate([
        (Inches(0.22), t1, i1, c1),
        (Inches(6.85), t2, i2, c2),
    ]):
        add_rect(slide, cx, CY, CW, CH, fill=LGRAY)
        add_rect(slide, cx, CY, CW, Inches(0.52), fill=col)
        txb(slide, t, cx + Inches(0.1), CY + Inches(0.04),
            CW - Inches(0.2), Inches(0.46), sz=15, bold=True,
            align=PP_ALIGN.CENTER)
        bullets(slide, items,
                cx + Inches(0.15), CY + Inches(0.58),
                CW - Inches(0.3), CH - Inches(0.7), sz=13, color=DGRAY)
    return slide


def slide_timeline(prs, title, steps, label=""):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), W, Inches(0.07), fill=ORANGE)
    if label:
        txb(slide, label, Inches(0.3), Inches(0.07), Inches(8), Inches(0.42),
            sz=11, color=RGBColor(0xBB, 0xCC, 0xDD))
    txb(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.95),
        sz=24, bold=True)

    colors = [NAVY, BLUE, TEAL, GREEN, ORANGE, PURPLE, RED]
    n = len(steps)
    sw = (Inches(12.88) / n) - Inches(0.09)
    for i, (st, sitems) in enumerate(steps):
        sx = Inches(0.22) + i * (sw + Inches(0.09))
        sy = Inches(1.47)
        add_rect(slide, sx, sy, sw, Inches(0.62), fill=colors[i % len(colors)])
        txb(slide, f"{i+1}. {st}", sx + Inches(0.05), sy + Inches(0.06),
            sw - Inches(0.1), Inches(0.53), sz=12, bold=True,
            align=PP_ALIGN.CENTER)
        bullets(slide, sitems,
                sx + Inches(0.06), sy + Inches(0.68),
                sw - Inches(0.12), Inches(5.1),
                sz=10, color=DGRAY, prefix="• ")
    return slide


def slide_table(prs, title, headers, rows, label=""):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), W, Inches(0.07), fill=ORANGE)
    if label:
        txb(slide, label, Inches(0.3), Inches(0.07), Inches(8), Inches(0.42),
            sz=11, color=RGBColor(0xBB, 0xCC, 0xDD))
    txb(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.95),
        sz=24, bold=True)

    nc = len(headers)
    nr = len(rows) + 1
    tbl = slide.shapes.add_table(
        nr, nc, Inches(0.38), Inches(1.5), Inches(12.6), Inches(5.7)).table

    cw = Inches(12.6) // nc
    for j in range(nc):
        tbl.columns[j].width = cw

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Calibri"

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LGRAY if i % 2 == 0 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val); r.font.size = Pt(12)
            r.font.color.rgb = DGRAY; r.font.name = "Calibri"
    return slide


def slide_evaluation(prs):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), W, Inches(0.07), fill=ORANGE)
    txb(slide, "ÉVALUATION & VALIDATION DES ACQUIS",
        Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.95),
        sz=24, bold=True)

    blocs = [
        (BLUE,   "QCM THÉORIQUE",    "30 questions — 30 min", [
            "20 questions de connaissances générales toutes branches",
            "10 questions sur la réglementation CIMA",
            "Score minimum requis : 14/20 (70%)",
            "Correction collective commentée",
        ]),
        (TEAL,   "CAS PRATIQUES",    "3 dossiers — 60 min", [
            "Dossier Automobile : sinistre RC + dommages propres",
            "Dossier Incendie : risque habitation + règle proportionnelle",
            "Dossier Santé : hospitalisation avec tiers payant",
            "Grille d'évaluation fournie aux participants",
        ]),
        (ORANGE, "MISE EN SITUATION", "Jeux de rôle — 45 min", [
            "Réception d'une déclaration de sinistre par téléphone",
            "Gestion d'un assuré mécontent (refus de garantie)",
            "Rédaction d'un courrier de règlement d'indemnité",
            "Évaluation comportementale et communication",
        ]),
    ]

    for i, (col, t, sub, items) in enumerate(blocs):
        bx = Inches(0.35) + i * Inches(4.32)
        by = Inches(1.48)
        bw = Inches(4.15)
        bh = Inches(5.82)
        add_rect(slide, bx, by, bw, bh, fill=LGRAY)
        add_rect(slide, bx, by, bw, Inches(0.85), fill=col)
        txb(slide, t, bx + Inches(0.1), by + Inches(0.05),
            bw - Inches(0.2), Inches(0.48), sz=15, bold=True,
            align=PP_ALIGN.CENTER)
        txb(slide, sub, bx + Inches(0.1), by + Inches(0.54),
            bw - Inches(0.2), Inches(0.32), sz=11, italic=True,
            align=PP_ALIGN.CENTER, color=WHITE)
        bullets(slide, items,
                bx + Inches(0.18), by + Inches(0.95),
                bw - Inches(0.35), bh - Inches(1.1),
                sz=13, color=DGRAY, prefix="✓  ")
    return slide


def slide_conclusion(prs):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, NAVY)
    add_rect(slide, 0, 0, Inches(0.28), H, fill=ORANGE)
    add_rect(slide, 0, H - Inches(0.15), W, Inches(0.15), fill=ORANGE)

    txb(slide, "MERCI DE VOTRE ATTENTION",
        Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.1),
        sz=38, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.3), Inches(2.35), Inches(6.7), Inches(0.08), fill=ORANGE)

    txb(slide, "Formation — Gestion des Sinistres Toutes Branches",
        Inches(0.8), Inches(2.55), Inches(11.7), Inches(0.75),
        sz=19, italic=True, align=PP_ALIGN.CENTER,
        color=RGBColor(0xCC, 0xDD, 0xEE))

    txb(slide, "Formateur : TALOM Eric",
        Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.5),
        sz=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    txb(slide, "◈  POINTS CLÉS À RETENIR",
        Inches(2.0), Inches(3.45), Inches(9.3), Inches(0.6),
        sz=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    retenir = [
        "Rigueur et célérité dans le traitement des déclarations de sinistre",
        "Respect strict des délais réglementaires CIMA à chaque étape",
        "Communication transparente et proactive avec l'assuré",
        "Archivage rigoureux pour la traçabilité et la conformité légale",
        "Actualisation permanente des connaissances réglementaires et techniques",
    ]
    bullets(slide, retenir,
            Inches(2.2), Inches(4.15), Inches(8.9), Inches(3.0),
            sz=14, color=RGBColor(0xEE, 0xF5, 0xFF), prefix="◈  ")
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# CONSTRUCTION DE LA PRÉSENTATION
# ══════════════════════════════════════════════════════════════════════════════

# ── PAGE DE COUVERTURE ───────────────────────────────────────────────────────
slide_cover(
    prs,
    title="GESTION DES SINISTRES\nToutes Branches",
    subtitle=(
        "Formation de remise à niveau — Professionnels de l'Assurance\n\n"
        "Automobile  •  Incendie  •  Transports Facultés\n"
        "Santé  •  Individuel Accident"
    ),
)

# ── CONDUCTEUR / SOMMAIRE ────────────────────────────────────────────────────
slide_sommaire(prs, [
    ("0",  "Introduction & Cadre Réglementaire",    "1h"),
    ("1",  "Processus Général de Gestion",          "2h"),
    ("2",  "Branche Automobile",                    "3h"),
    ("3",  "Branche Incendie",                      "2h"),
    ("4",  "Transports Facultés (TFT / A / M)",     "3h"),
    ("5",  "Branche Santé",                         "2h"),
    ("6",  "Individuel Accident",                   "2h"),
    ("7",  "Classement & Archivage",                "1h"),
    ("8",  "Évaluation & Cas Pratiques",            "2h"),
])

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 0 — INTRODUCTION & CADRE RÉGLEMENTAIRE
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "0", "INTRODUCTION\n& CADRE RÉGLEMENTAIRE",
              "Fondements juridiques — Code CIMA — Obligations des parties")

slide_content(prs, "Objectifs Pédagogiques & Public Cible", [
    {"text": "Objectifs pédagogiques", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Maîtriser les procédures de réception, instruction et règlement des sinistres",           "level": 1},
    {"text": "Appliquer le cadre réglementaire CIMA à chaque stade du traitement",                     "level": 1},
    {"text": "Gérer efficacement les relations avec les assurés, experts et prestataires",              "level": 1},
    {"text": "Produire et archiver une documentation conforme aux exigences légales",                  "level": 1},
    {"text": "Analyser les sinistres par branche avec rigueur et méthode",                             "level": 1},
    {"text": "Public cible & prérequis", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Gestionnaires sinistres en activité souhaitant se remettre à niveau",                    "level": 1},
    {"text": "Inspecteurs, rédacteurs, agents généraux et courtiers",                                  "level": 1},
    {"text": "Toute personne impliquée dans le traitement ou le suivi de dossiers sinistres",          "level": 1},
    {"text": "Prérequis : connaissances de base en assurance IARD et Personnes",                       "level": 1},
], label="MODULE 0 — Introduction")

slide_two_col(prs, "Cadre Réglementaire — Code CIMA & Principes Incontournables",
    "Textes Fondamentaux", [
        "Code CIMA — Traité du 10 juillet 1992 (14 États membres CIMA)",
        "Art. 12 à 18 : Obligations de l'assureur en cas de sinistre",
        "Art. 15 : Délai d'instruction — accusé de réception 48h",
        "Art. 16 : Délai de règlement — 30 jours après accord parties",
        "Art. 18 : Pénalités pour retard — intérêts légaux + 50%",
        "Art. 28 : Prescription biennale des actions nées du contrat",
        "Circulaires CRCA et lois nationales complémentaires",
        "Code des assurances national selon pays (Sénégal, Côte d'Ivoire…)",
    ],
    "Principes Juridiques Fondamentaux", [
        "Principe indemnitaire — pas d'enrichissement sans cause",
        "Principe de bonne foi — obligation des deux parties",
        "Déclaration exacte — sanctions si fausse déclaration (nullité)",
        "Subrogation — l'assureur se substitue à l'assuré pour recours",
        "Déchéance — perte du droit à garantie si non-respect des obligations",
        "Contribution entre assureurs (coassurance / double assurance)",
        "Règle proportionnelle de capitaux (insuffisance d'assurance)",
        "Expertise contradictoire — droit de contre-expertise assuré",
    ],
    label="MODULE 0 — Introduction")

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 1 — PROCESSUS GÉNÉRAL
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "1", "PROCESSUS GÉNÉRAL\nDE GESTION DES SINISTRES",
              "Du premier appel au règlement définitif et à l'archivage")

slide_timeline(prs, "Les 6 Étapes Clés du Traitement d'un Sinistre", [
    ("DÉCLARATION", [
        "Réception de l'avis de sinistre",
        "Vérification identité assuré",
        "Enregistrement dans le SIG",
        "Attribution du n° de dossier",
        "Accusé de réception 24–48h",
    ]),
    ("INSTRUCTION", [
        "Vérification des garanties",
        "Analyse contrat / police",
        "Recherche de responsabilité",
        "Demande pièces justificatives",
        "Désignation d'un expert",
    ]),
    ("EXPERTISE", [
        "Mission d'expert agréé",
        "Rapport contradictoire",
        "Contre-expertise éventuelle",
        "Chiffrage des dommages",
        "Validation du rapport",
    ]),
    ("DÉCISION", [
        "Acceptation / refus garantie",
        "Courrier motivé obligatoire",
        "Délai : 30 jours (CIMA)",
        "Accord amiable ou litige",
        "Notification à l'assuré",
    ]),
    ("RÈGLEMENT", [
        "Décompte d'indemnité",
        "Franchises / vétustés",
        "Ordre de paiement",
        "Quittance de règlement",
        "Paiement dans les délais",
    ]),
    ("CLÔTURE", [
        "Classement du dossier",
        "Archivage physique/GED",
        "Recours / subrogation",
        "Mise à jour provisions",
        "Reporting statistique",
    ]),
], label="MODULE 1 — Processus Général")

slide_content(prs, "Réception & Enregistrement de la Déclaration de Sinistre", [
    {"text": "Canaux de déclaration acceptés", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Téléphone (standard sinistres) — enregistrement de l'appel recommandé",    "level": 1},
    {"text": "Courrier recommandé avec accusé de réception",                              "level": 1},
    {"text": "Courriel horodaté (valeur probante selon législation nationale)",           "level": 1},
    {"text": "Déclaration en agence — formulaire papier daté et signé",                  "level": 1},
    {"text": "Application mobile / espace client en ligne (horodatage automatique)",     "level": 1},
    {"text": "Informations obligatoires à collecter dès la déclaration", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Numéro de police / contrat + identité complète de l'assuré",               "level": 1},
    {"text": "Date, heure, lieu et circonstances détaillées du sinistre",                "level": 1},
    {"text": "Nature et estimation provisoire des dommages (matériels, corporels)",      "level": 1},
    {"text": "Identité des tiers impliqués, témoins et intervenants (pompiers, police)", "level": 1},
    {"text": "Existence d'autres assurances couvrant le même risque (déclaration obligatoire)", "level": 1},
    {"text": "Mesures conservatoires prises ou à prendre",                               "level": 1},
], label="MODULE 1 — Processus Général")

slide_content(prs, "Instruction, Décision & Règlement de l'Indemnité", [
    {"text": "Phase d'instruction — vérification & investigation", "level": 0, "bold": True, "color": BLUE, "size": 18},
    {"text": "Contrôle validité contrat : primes payées, délai de carence respecté",     "level": 1},
    {"text": "Vérification que le sinistre est dans l'objet et les garanties souscrites","level": 1},
    {"text": "Recherche des exclusions (art. contractuels + exclusions légales CIMA)",   "level": 1},
    {"text": "Constitution du dossier probatoire : PV police, photos, attestations",     "level": 1},
    {"text": "Phase de décision — position de l'assureur", "level": 0, "bold": True, "color": BLUE, "size": 18},
    {"text": "Délai de 30 jours à compter de la réception de toutes pièces (CIMA art. 16)", "level": 1},
    {"text": "Lettre de garantie OU refus de garantie motivé — obligatoirement écrit",   "level": 1},
    {"text": "Proposition d'indemnisation chiffrée et détaillée",                        "level": 1},
    {"text": "Phase de règlement — paiement de l'indemnité", "level": 0, "bold": True, "color": BLUE, "size": 18},
    {"text": "Calcul indemnité : valeur assurée, franchise, vétusté, règle proportionnelle", "level": 1},
    {"text": "Mode de règlement : virement bancaire, chèque, paiement direct prestataire", "level": 1},
    {"text": "Quittance finale signée par l'assuré — clôture et archivage du dossier",  "level": 1},
], label="MODULE 1 — Processus Général")

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 2 — BRANCHE AUTOMOBILE
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "2", "BRANCHE AUTOMOBILE",
              "RC Obligatoire  •  Dommages  •  Vol  •  Bris de Glace  •  Assistance")

slide_content(prs, "Automobile — Vue d'Ensemble des Garanties", [
    {"text": "Garanties obligatoires (RC Auto)", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Responsabilité Civile Obligatoire (RCO) — dommages corporels et matériels aux tiers", "level": 1},
    {"text": "Garantie défense-recours (protection juridique de l'assuré)",                         "level": 1},
    {"text": "Garanties facultatives courantes", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Dommages Tous Accidents (DTA) — véhicule assuré quelle que soit la responsabilité",   "level": 1},
    {"text": "Dommages Collision — sinistre avec tiers identifié et responsable",                   "level": 1},
    {"text": "Vol / Incendie / Bris de Glace (BDG) — couvertures distinctes paramétrables",         "level": 1},
    {"text": "Catastrophes naturelles et événements climatiques (TCG)",                             "level": 1},
    {"text": "Assistance et remorquage 24h/24 — véhicule de remplacement",                         "level": 1},
    {"text": "Garanties Corporelles Conducteur (GCC) — capital décès/invalidité",                  "level": 1},
    {"text": "Conventions applicables entre assureurs", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Convention IDA (Indemnisation Directe de l'Assuré) — règlement sans attendre recours","level": 1},
    {"text": "Convention IRSA (recours inter-assureurs) — forfaits par catégorie de sinistre",      "level": 1},
], label="MODULE 2 — Automobile", accent=BLUE)

slide_content(prs, "Automobile — Instruction d'un Sinistre RC & Dommages", [
    {"text": "Documents à collecter impérativement", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Constat amiable d'accident signé par les deux conducteurs (double exemplaire)",        "level": 1},
    {"text": "Permis de conduire, carte grise, attestation d'assurance valide à la date du sinistre","level": 1},
    {"text": "PV de police / gendarmerie (obligatoire si blessés ou sinistre grave)",               "level": 1},
    {"text": "Photos du lieu, des véhicules et de l'ensemble des dommages visibles",                "level": 1},
    {"text": "Devis de réparation des garages agréés ou mandatés",                                  "level": 1},
    {"text": "Expertise automobile — processus et outils", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Mandatement de l'expert agréé — délai de mission : 5 jours ouvrables",               "level": 1},
    {"text": "Rapport d'expertise : valeur à neuf, coefficient vétusté, valeur vénale (VV)",        "level": 1},
    {"text": "Destruction totale (DT) : VV – valeur épave = indemnité nette",                      "level": 1},
    {"text": "Contre-expertise possible à la demande de l'assuré (frais partagés)",                "level": 1},
    {"text": "Recours & subrogation après règlement", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Droit de recours contre le tiers responsable ou son assureur",                       "level": 1},
    {"text": "Application Convention IRSA pour le règlement inter-assureurs",                      "level": 1},
], label="MODULE 2 — Automobile", accent=BLUE)

slide_table(prs, "Automobile — Tableau de Bord des Délais Réglementaires",
    ["Étape", "Délai réglementaire", "Responsable", "Conséquence si non-respecté"],
    [
        ["Accusé de réception de la déclaration",   "24 à 48 h",                  "Gestionnaire",       "Engagement de responsabilité de l'assureur"],
        ["Mandatement de l'expert",                 "5 jours ouvrables",          "Gestionnaire",       "Pénalités CIMA — signalement possible"],
        ["Dépôt du rapport d'expertise",            "15 jours ouvrables",         "Expert agréé",       "Mise en cause de l'expert par la CRCA"],
        ["Proposition d'indemnisation à l'assuré",  "30 jours après accord",      "Assureur",           "Intérêts légaux majorés de 50% (art. 18 CIMA)"],
        ["Paiement effectif de l'indemnité",        "10 jours après quittance",   "Compta / Assureur",  "Pénalités CIMA — réclamation possible"],
        ["Recours contre tiers responsable",        "Prescription biennale",      "Juriste / Contentieux", "Forclusion — perte du droit de recours"],
    ],
    label="MODULE 2 — Automobile")

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 3 — BRANCHE INCENDIE
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "3", "BRANCHE INCENDIE",
              "Risques simples  •  Risques d'entreprise  •  Règle proportionnelle")

slide_content(prs, "Incendie — Périls Couverts, Extensions et Exclusions", [
    {"text": "Périls de base (garantis automatiquement)", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Incendie, explosion, implosion — Foudre directe et ses suites immédiates",             "level": 1},
    {"text": "Chute d'appareils de navigation aérienne (CANA)",                                     "level": 1},
    {"text": "Dommages causés par les secours lors de l'intervention (eau d'extinction)",           "level": 1},
    {"text": "Extensions courantes (options contractuelles)", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Dégâts des eaux (DDO) : infiltrations, rupture de canalisations, inondation",         "level": 1},
    {"text": "Tempêtes, cyclones, grêle (TCG) — couverture événements climatiques",                 "level": 1},
    {"text": "Bris de glaces / Vol avec effraction (conditions particulières d'assurabilité)",      "level": 1},
    {"text": "Pertes d'exploitation consécutives à l'incendie (marge brute garantie)",              "level": 1},
    {"text": "RC Voisins et Recours des Voisins (RVRV) — dommages aux tiers",                      "level": 1},
    {"text": "Risques exclus courants", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Guerre, émeutes, actes de terrorisme (sauf extension spécifique)",                    "level": 1},
    {"text": "Faute intentionnelle ou dolosive de l'assuré (art. 12 Code CIMA)",                    "level": 1},
    {"text": "Dommages électriques (machines électriques) sans extension dédiée",                  "level": 1},
], label="MODULE 3 — Incendie", accent=RED)

slide_two_col(prs, "Incendie — Instruction du Dossier & Calcul de l'Indemnité",
    "Pièces constitutives du dossier sinistre", [
        "Déclaration écrite : date, heure, circonstances précises",
        "PV des sapeurs-pompiers (original ou copie certifiée)",
        "PV de police / gendarmerie (incendie criminel suspecté)",
        "Inventaire détaillé des biens sinistrés (mobilier + immobilier)",
        "Factures d'achat ou estimations des biens détruits / endommagés",
        "Photos et vidéos des dommages (préalables à tout déblaiement)",
        "Devis de remise en état des entreprises spécialisées",
        "Attestation syndic ou gestionnaire (si immeuble collectif)",
        "Rapport d'expertise incendie (expert mandaté par l'assureur)",
    ],
    "Calcul Indemnité — Règle Proportionnelle", [
        "Valeur assurée (VA) vs Valeur réelle (VR) au jour du sinistre",
        "Si VA ≥ VR → indemnité = sinistre × (1 — franchise)",
        "Si VA < VR → règle proportionnelle : I = S × (VA / VR)",
        "Exemple : VA=10M, VR=20M, S=4M → I = 4 × (10/20) = 2 M",
        "Application de la vétusté sur biens mobiliers et équipements",
        "Clause valeur à neuf : remboursement sans déduction vétusté",
        "Coassurance : répartition entre assureurs (chef de file)",
        "Pertes d'exploitation : calcul de la marge brute perdue",
        "Franchise absolue ou relative selon conditions particulières",
    ],
    c1=RED, c2=NAVY, label="MODULE 3 — Incendie")

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 4 — BRANCHE TRANSPORTS FACULTÉS
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "4", "BRANCHE TRANSPORTS\nFACULTÉS",
              "Facultés Terrestres  •  Facultés Aériennes  •  Facultés Maritimes")

slide_content(prs, "Transports — Notions Fondamentales, Intervenants & Documents", [
    {"text": "Définitions clés", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Facultés : marchandises transportées (par opposition à 'corps' = véhicule/navire/aéronef)", "level": 1},
    {"text": "Chargeur / expéditeur : propriétaire de la marchandise (assuré principal)",               "level": 1},
    {"text": "Destinataire : peut bénéficier de l'assurance selon la clause 'à qui il appartiendra'",  "level": 1},
    {"text": "Transporteur : tenu à une obligation de résultat — recours possible",                    "level": 1},
    {"text": "Documents de transport essentiels par mode", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Lettre de Voiture (LDV) / CMR international — transport routier",                        "level": 1},
    {"text": "Lettre de Transport Aérien (LTA / Air Waybill) — fret aérien",                           "level": 1},
    {"text": "Connaissement maritime (Bill of Lading) — titre de propriété de la marchandise",         "level": 1},
    {"text": "Certificat d'assurance transport — preuve de couverture présentable aux tiers",          "level": 1},
    {"text": "Garanties standard — Clauses Instituts de Londres (A/B/C)", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Clause A (All Risks) : tous risques sauf exclusions expressément listées",               "level": 1},
    {"text": "Clause B : avaries particulières + risques listés (naufrage, échouage, incendie, etc.)", "level": 1},
    {"text": "Clause C : risques majeurs uniquement — franchise d'avarie particulière (FAP sauf)",     "level": 1},
], label="MODULE 4 — Transports", accent=TEAL)

slide_two_col(prs, "Facultés Terrestres & Aériennes — Instruction des Sinistres",
    "Facultés Terrestres (route / rail)", [
        "Réserves sur le bon de livraison à la réception (OBLIGATOIRE)",
        "Confirmation écrite des réserves au transporteur sous 3 jours",
        "Constat contradictoire avec le transporteur ou son représentant",
        "Photos détaillées des dommages avant tout déplacement",
        "Documents : LDV/CMR, bon de livraison, facture commerciale",
        "Expertise des marchandises (qualitative + quantitative)",
        "Recours CMR : plafond 8,33 DTS/kg de poids brut",
        "Prescription : 1 an (national) — 1 an (CMR international)",
        "Avaries communes rares en transport terrestre",
    ],
    "Facultés Aériennes (fret aérien)", [
        "Déclaration à la compagnie aérienne dans les 24h suivant réception",
        "LTA (Air Waybill) : document de référence pour tous recours",
        "Convention de Montréal 1999 (remplace Convention de Varsovie)",
        "Limite responsabilité transporteur : 19 DTS/kg",
        "Expertise marchandises en présence du représentant transporteur",
        "Délais de protestation : 14 jours (avarie), 21 jours (retard)",
        "Prescription 2 ans à compter de la livraison",
        "Sinistres typiques : chocs, mouillures, manquants, vol en soute",
    ],
    c1=TEAL, c2=BLUE, label="MODULE 4 — Transports")

slide_content(prs, "Facultés Maritimes — Procédure & Avarie Commune", [
    {"text": "Actions immédiates à l'arrivée du navire", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Émission de réserves sur le Bon à Délivrer (BAD) lors de la prise en charge",           "level": 1},
    {"text": "Notification écrite au transporteur dans les 3 jours (avarie apparente)",               "level": 1},
    {"text": "Mandatement immédiat d'un commissaire d'avaries (expertise contradictoire)",            "level": 1},
    {"text": "Documents indispensables au dossier sinistre", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Connaissement (Bill of Lading) original — titre de propriété de la marchandise",        "level": 1},
    {"text": "Facture commerciale + liste de colisage (packing list) détaillée",                      "level": 1},
    {"text": "Certificats (poids, origine, phytosanitaire, qualité) selon nature marchandise",        "level": 1},
    {"text": "Rapport du commissaire d'avaries + protestation de mer du capitaine",                   "level": 1},
    {"text": "Déclaration d'avaries à la douane si nécessaire",                                      "level": 1},
    {"text": "Avarie Commune — Règles d'York-Anvers", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Acte volontaire et raisonnable pour sauver l'expédition en péril commun",              "level": 1},
    {"text": "Contribution proportionnelle de tous les intéressés à la valeur de l'expédition",      "level": 1},
    {"text": "Expertise spécialisée (average adjuster) — délais longs : 6 à 24 mois",               "level": 1},
], label="MODULE 4 — Transports", accent=TEAL)

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 5 — BRANCHE SANTÉ
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "5", "BRANCHE SANTÉ",
              "Frais médicaux  •  Hospitalisation  •  Maternité  •  Tiers Payant")

slide_content(prs, "Santé — Garanties, Mécanismes de Remboursement & Documents", [
    {"text": "Types de garanties santé", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Frais médicaux : consultations, médicaments, analyses, imagerie médicale",           "level": 1},
    {"text": "Hospitalisation : frais de séjour, honoraires chirurgicaux, forfait journalier",     "level": 1},
    {"text": "Maternité : frais d'accouchement, suivi prénatal et postnatal",                      "level": 1},
    {"text": "Soins dentaires et optique selon tableaux de garanties contractuels",                "level": 1},
    {"text": "Évacuation sanitaire et rapatriement médical international",                         "level": 1},
    {"text": "Mécanismes de remboursement", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Tiers payant : prise en charge directe assureur → prestataire (sans avance frais)",  "level": 1},
    {"text": "Remboursement a posteriori : l'assuré avance les frais puis se fait rembourser",     "level": 1},
    {"text": "Réseau de soins agréés vs hors réseau : différences de taux de remboursement",      "level": 1},
    {"text": "Ticket modérateur : part restant à charge de l'assuré (% ou montant fixe)",         "level": 1},
    {"text": "Documents standard pour la liquidation", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Ordonnance médicale originale + feuille de soins signée du médecin",                "level": 1},
    {"text": "Bulletin de Prise en Charge (BPC) pour hospitalisation programmée",                 "level": 1},
    {"text": "Facture détaillée du prestataire avec codification nationale des actes",             "level": 1},
], label="MODULE 5 — Santé", accent=GREEN)

slide_two_col(prs, "Santé — Gestion des Dossiers Hospitalisation & Contrôle",
    "Prise en Charge Préalable (PEC)", [
        "Demande de PEC déposée AVANT l'admission (sauf urgence médicale)",
        "Documents : devis d'hospitalisation signé du médecin traitant",
        "Délai de réponse assureur : 48h (hors urgence) / immédiat (urgence)",
        "Lettre de garantie envoyée directement à l'établissement",
        "Mention explicite des plafonds et exclusions pris en charge",
        "Suivi en cas de prolongation de séjour (renouvellement BPC)",
        "Régulation médicale si doute sur la pertinence ou durée des soins",
        "Protocole urgence : PEC rétroactive sous 72h après admission",
        "Gestion des hospitalisations à l'étranger : réseau international",
    ],
    "Contrôle & Liquidation des Factures", [
        "Vérification conformité facture / devis initial (écarts signalés)",
        "Contrôle médical par le médecin-conseil assureur si besoin",
        "Codification CCAM / Nomenclature nationale des actes médicaux",
        "Application des tableaux de garanties et plafonds contractuels",
        "Déduction des organismes sociaux (Sécurité Sociale, mutuelles)",
        "Calcul du ticket modérateur à la charge de l'assuré",
        "Émission du décompte de remboursement détaillé",
        "Délai de règlement : 10 jours après réception du dossier complet",
        "Détection des anomalies : doublons, tarifs hors nomenclature",
    ],
    c1=GREEN, c2=TEAL, label="MODULE 5 — Santé")

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 6 — INDIVIDUEL ACCIDENT
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "6", "INDIVIDUEL ACCIDENT",
              "Décès  •  PTIA  •  IPP / IPT  •  ITT  •  Frais médicaux accident")

slide_content(prs, "Individuel Accident — Garanties, Définitions & Exclusions", [
    {"text": "Garanties principales de la branche IA", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Décès accidentel : capital versé aux bénéficiaires désignés au contrat",               "level": 1},
    {"text": "PTIA (Perte Totale et Irréversible d'Autonomie) : 100% du capital souscrit",           "level": 1},
    {"text": "IPP (Invalidité Permanente Partielle) : % du capital selon barème médical",            "level": 1},
    {"text": "IPT (Invalidité Permanente Totale) : taux ≥ 66% d'incapacité fonctionnelle",          "level": 1},
    {"text": "ITT (Incapacité Temporaire Totale) : indemnité journalière pendant l'arrêt de travail","level": 1},
    {"text": "Frais médicaux et chirurgicaux consécutifs à un accident corporel",                    "level": 1},
    {"text": "Définitions contractuelles essentielles", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Accident : atteinte corporelle soudaine, involontaire, causée par un événement extérieur", "level": 1},
    {"text": "Distinction accident / maladie : enjeu fondamental pour le déclenchement de la garantie", "level": 1},
    {"text": "Consolidation : date à partir de laquelle l'état de santé est stabilisé",             "level": 1},
    {"text": "Exclusions standard : sports extrêmes, guerre, état alcoolique/substances, fait intentionnel", "level": 1},
], label="MODULE 6 — Individuel Accident", accent=PURPLE)

slide_content(prs, "Individuel Accident — Instruction, Expertise Médicale & Calcul", [
    {"text": "Pièces justificatives selon type de sinistre", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Décès : certificat de décès + rapport médical + actes de naissance bénéficiaires",    "level": 1},
    {"text": "IPP/IPT : certificat médical initial + rapport de consolidation du médecin traitant", "level": 1},
    {"text": "ITT : arrêts de travail médicaux renouvelés + bulletins de salaire",                  "level": 1},
    {"text": "PTIA : décision MDPH ou équivalent national + rapport médical spécialisé",            "level": 1},
    {"text": "Évaluation médicale du taux d'invalidité", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "Médecin-conseil désigné par l'assureur (examen médical contradictoire obligatoire)",  "level": 1},
    {"text": "Référence : Guide Barème Européen d'Évaluation des Atteintes à l'Intégrité Physique", "level": 1},
    {"text": "Taux fonctionnel (AIPP) + incidence professionnelle éventuelle",                     "level": 1},
    {"text": "Contre-expertise à la demande de l'assuré (droit CIMA)",                            "level": 1},
    {"text": "Formules de calcul des indemnités", "level": 0, "bold": True, "color": NAVY, "size": 18},
    {"text": "IPP : Capital souscrit × Taux d'invalidité (ex : 10 000 000 FCFA × 30% = 3 000 000 FCFA)", "level": 1},
    {"text": "ITT : Indemnité journalière × (jours d'arrêt − franchise temporelle de 3 jours)",    "level": 1},
    {"text": "Cumul de garanties possible : IA + frais médicaux + prévoyance employeur",           "level": 1},
], label="MODULE 6 — Individuel Accident", accent=PURPLE)

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 7 — CLASSEMENT & ARCHIVAGE
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "7", "CLASSEMENT &\nARCHIVAGE DES DOSSIERS",
              "Organisation  •  Conservation  •  Conformité réglementaire  •  GED")

slide_two_col(prs, "Classement & Archivage — Règles et Bonnes Pratiques",
    "Classement physique (dossier papier)", [
        "Chemise de couleur par branche (code chromatique normalisé)",
        "Numérotation unique : branche / année / numéro séquentiel",
        "Bordereau des pièces daté et signé en début de chemise",
        "Ordre chronologique strict des pièces dans le dossier",
        "Séparation nette : dossiers actifs / clôturés / contentieux",
        "Local d'archives sécurisé (accès restreint, protection incendie)",
        "Inventaire trimestriel des dossiers archivés",
        "Délai légal de conservation : 10 ans après clôture (CIMA art. 28)",
        "Destruction sécurisée : broyage certifié après délai légal",
    ],
    "Archivage numérique (GED)", [
        "Scan de tous les documents originaux à réception (qualité contrôlée)",
        "Indexation normalisée : branche, date, assuré, n° police, n° sinistre",
        "Convention de nommage interne des fichiers (charte documentaire)",
        "Contrôle qualité des scans : lisibilité, complétude, orientation",
        "Sauvegarde redondante : serveur local + cloud sécurisé (chiffré)",
        "Accès par profil utilisateur avec traçabilité des consultations",
        "Plan de reprise sur incident (PRA/PCA) validé",
        "Purge sécurisée après délai légal (effacement certifié RGPD)",
        "Audit annuel de conformité de l'archivage numérique",
    ],
    c1=NAVY, c2=BLUE, label="MODULE 7 — Classement & Archivage")

# ══════════════════════════════════════════════════════════════════════════════
# MODULE 8 — ÉVALUATION & CAS PRATIQUES
# ══════════════════════════════════════════════════════════════════════════════
slide_section(prs, "8", "ÉVALUATION &\nCAS PRATIQUES",
              "Mise en application — QCM — Jeux de rôle — Dossiers simulés")

slide_evaluation(prs)

slide_table(prs, "Récapitulatif — Délais Clés à Maîtriser par Branche",
    ["Branche", "Acte / Étape", "Délai CIMA / Convention", "Référence"],
    [
        ["Toutes branches",      "Accusé réception déclaration",     "24 à 48 heures",         "Art. 15 CIMA"],
        ["Toutes branches",      "Proposition d'indemnisation",      "30 jours après accord",  "Art. 16 CIMA"],
        ["Toutes branches",      "Paiement effectif de l'indemnité", "10 jours après quittance","Art. 18 CIMA"],
        ["Automobile",           "Mandatement expert",               "5 jours ouvrables",      "Convention expertise"],
        ["Automobile",           "Rapport d'expertise",              "15 jours ouvrables",     "Convention expertise"],
        ["Transports maritimes", "Notification avaries",             "3 jours après livraison","Convention Bruxelles"],
        ["Transports aériens",   "Protestation avaries",             "14 jours",               "Convention Montréal"],
        ["Santé",                "Réponse à la demande de PEC",      "48 heures (hors urgence)","Conditions générales"],
        ["IA — ITT",             "Franchise temporelle ITT",         "3 premiers jours",       "Conditions particulières"],
    ],
    label="MODULE 8 — Évaluation")

# ══════════════════════════════════════════════════════════════════════════════
# CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
slide_conclusion(prs)

# ══════════════════════════════════════════════════════════════════════════════
# SAUVEGARDE
# ══════════════════════════════════════════════════════════════════════════════
output = "Formation_Gestion_Sinistres_Toutes_Branches.pptx"
prs.save(output)
print(f"\nOK  Presentation generee : {output}")
print(f"    Nombre de diapositives : {len(prs.slides)}")
