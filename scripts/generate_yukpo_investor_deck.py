# -*- coding: utf-8 -*-
"""
Génère Yukpo_Investisseur_Presentation.pptx — dossier investisseur (200 M FCFA).

Langage volontairement accessible. Aucune référence aux chemins de fichiers internes
dans les diapos (usage externe : investisseurs).
"""
from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

# --- Typographie (lisibilité & rendu professionnel) ---
FS_H1 = 28
FS_SUBTITLE = 12.5
FS_BODY = 12
FS_CARD_TITLE = 12.5
FS_CARD_BULLET = 10.8
FS_TABLE_HDR = 11
FS_TABLE_BODY = 10.8
FS_NOTE = 10.5
FS_FOOTER = 9.5
FS_STAMP = 10
FS_KPI_BIG = 32
FS_KPI_LAB = 12.5
FS_KPI_SUB = 10
FS_COVER_KPI = 36

# --- Besoin global (aligné docs/investisseurs/BESOIN_FINANCEMENT_200M_FCFA.md) ---
B = {
    "total_m": 200,
    "marketing_m": 118,
    "personnel_m": 42,
    "tech_m": 18,
    "legal_m": 8,
    "divers_m": 7,
    "reserve_m": 7,
    "equity_cap_pct": 15,
}

assert sum(
    [
        B["marketing_m"],
        B["personnel_m"],
        B["tech_m"],
        B["legal_m"],
        B["divers_m"],
        B["reserve_m"],
    ]
) == B["total_m"]

# Répartition des charges sur 2 ans (mai 2026 – avril 2028) — cohérente avec l’enveloppe 200 M
CHARGES_ANNEE1_M = 118  # mai 2026 – avril 2027
CHARGES_ANNEE2_M = 82  # mai 2027 – avril 2028

# 8 trimestres (juil. 2026 → juin 2028), strictement alignés avec les bilans annuels :
# Année 1 (T3 26 → T2 27) : 118 M de charges, 62 M de revenus
# Année 2 (T3 27 → T2 28) : 82 M de charges, 106 M de revenus
PRELANCER_M = 0
TRIM_CHARGES_M = [28, 30, 30, 30, 21, 20, 20, 21]  # total 200 ; A1=118 ; A2=82

# Revenus (M FCFA / trimestre) — montée progressive, cohérente avec le bilan annuel
TRIM_REVENUS_M = [10, 14, 18, 20, 22, 25, 28, 31]  # total 168 ; A1=62 ; A2=106

FOUNDER_EXP_Y = 13

# --- Huit segments : 2 diapos chacun (besoins & solutions · impact population & économie) ---
MODULES_DEEP: list[dict[str, list[str] | str]] = [
    {
        "title": "E-commerce & vitrines marchands",
        "tag": "Commerce local",
        "besoins": [
            "Des milliers de boutiques physiques n’ont aucune vitrine fiable en ligne : les clients ne peuvent ni comparer ni commander sans passer par un intermédiaire humain.",
            "Créer un site ou une boutique « pro » coûte souvent plusieurs centaines de milliers de FCFA et exige des compétences techniques absentes chez la majorité des TPE.",
            "Les réseaux sociaux servent de catalogue improvisé : pas de stock clair, pas de livraison intégrée, pas de paiement traçable.",
            "Résultat : chiffre d’affaires capté par les grandes surfaces ou l’informel, alors que le stock et le savoir-faire existent déjà localement.",
        ],
        "solutions": [
            "Parcours guidé : photos, prix, disponibilité, promotions — sans compétence web préalable.",
            "Recherche par mots, image ou voix pour coller aux habitudes réelles des utilisateurs.",
            "Lien direct avec la livraison et le paiement mobile : une vente ne se perd plus entre deux conversations WhatsApp.",
            "Tableaux simples pour le marchand : commandes du jour, statut des paiements, clients à fidéliser.",
        ],
        "revolution": [
            "Les vendeurs de quartier gagnent la même visibilité qu’un site coûteux, sans quitter leur logique de travail habituelle.",
            "Les familles achètent « chez eux » en ayant confiance : moins de déplacements inutiles, plus d’argent qui reste dans le tissu local.",
            "Une partie du commerce informel peut enfin être comptabilisée et sécurisée, ce qui prépare l’inclusion financière.",
        ],
        "valeur": [
            "Commissions sur transactions et options de mise en avant : revenu direct lié au volume réel de ventes.",
            "Plus de marchands actifs = plus de récurrence d’usage de toute l’application (effet super-app).",
            "Réduction du coût d’acquisition marchand au fil du temps grâce au bouche-à-oreille digital.",
        ],
    },
    {
        "title": "Livraison",
        "tag": "Dernier kilomètre",
        "besoins": [
            "Les livraisons sont négociées au téléphone, prix variables, aucune traçabilité en cas de litige.",
            "Les coursiers manquent d’outil commun : charge irrégulière, revenus imprévisibles, risques sur la route non couverts.",
            "Les commerçants perdent des ventes faute de livreur disponible au bon moment.",
            "Les quartiers mal desservis sont pénalisés alors que la demande existe.",
        ],
        "solutions": [
            "Mise en relation automatique commande ↔ livreur avec estimation de temps et suivi GPS.",
            "Portefeuille intégré : moins de cash à transporter, historique clair pour le marchand et le coursier.",
            "Tarification plus lisible (forfaits, distance, urgence) pour éviter les malentendus.",
            "Files d’attente et zones chaudes pour équilibrer l’offre de livreurs sur la journée.",
        ],
        "revolution": [
            "Le travail de livreur devient une activité visible et rémunérée dans un cadre identifiable, pas seulement « au jour le jour ».",
            "Les ménages commandent davantage en confiance : la livraison n’est plus un parcours du combattant.",
            "Les petits commerces peuvent concurrencer les grandes surfaces sur la rapidité, là où elles ne vont pas.",
        ],
        "valeur": [
            "Prélèvement sur livraisons et abonnements livreurs : revenu récurrent lié au volume.",
            "Densité locale = baisse du coût marginal par course (effet réseau).",
            "La livraison tire les autres segments (commerce, restauration, santé) dans la même app.",
        ],
    },
    {
        "title": "Transport & mobilité",
        "tag": "Déplacements",
        "besoins": [
            "Billets interurbains et trajets urbains passent encore par files d’attente, groupes WhatsApp ou intermédiaires peu transparents.",
            "Les voyageurs manquent d’information centralisée sur horaires, prix et disponibilité en temps réel.",
            "Les chauffeurs et opérateurs subissent des commissions opaques et une faible fidélisation des clients.",
            "Les étudiants et travailleurs perdent un temps précieux à chercher un départ fiable.",
        ],
        "solutions": [
            "Billetterie et options de trajet regroupées avec paiement et reçu numérique (QR).",
            "Carte et filtres : partir « maintenant » ou à une heure donnée, avec comparaison simple.",
            "Espace dédié aux partenaires transport pour publier leurs offres dans les règles.",
            "Historique des trajets pour le remboursement professionnel ou familial.",
        ],
        "revolution": [
            "Réduction du stress et des arnaques sur la route : le voyageur sait à quoi s’en tenir avant de payer.",
            "Meilleure utilisation des sièges et des véhicules : moins de départs à vide, plus de revenus pour les opérateurs locaux.",
            "Les familles peuvent planifier santé, études et travail sans dépendre uniquement du bouche-à-oreille.",
        ],
        "valeur": [
            "Commissions sur ventes de billets et partenariats avec transporteurs.",
            "Usage récurrent (scolaire, travail) qui ancre Yukpo dans un geste quotidien.",
            "Données d’affluence utiles pour ajuster l’offre marketing et les zones prioritaires.",
        ],
    },
    {
        "title": "Navigation intelligente (NavigationScreen)",
        "tag": "Carte & itinéraires",
        "besoins": [
            "Les conducteurs et piétons jonglent entre plusieurs apps ou l’absence totale d’itinéraire fiable hors des axes principaux ; le trafic et les dangers locaux ne sont pas anticipés.",
            "Les usagers ont besoin d’un guidage clair en français (et voix) sans quitter l’écosystème Yukpo (commerce, livraison, rendez-vous).",
            "Les alertes « terrain » (travaux, zones sensibles, communauté) sont aujourd’hui dispersées sur des groupes ou des bouches-à-oreille.",
            "Peu d’outils combinent carte, statistiques d’activité et assistance IA dans une seule expérience mobile adaptée au contexte africain.",
        ],
        "solutions": [
            "Carte temps réel (Google Maps), itinéraires avec trafic, polylignes, étapes détaillées et plusieurs options de trajet.",
            "Guidage vocal (TTS) multilingue, instructions étape par étape, adaptation à la langue choisie dans l’app.",
            "Points d’intérêt le long du trajet, estimation de coûts, tarification dynamique et micro-paiement / crédits pour les fonctions avancées.",
            "Alertes communautaires et checkpoints (y compris tâche en arrière-plan), chat intelligent contextuel (NavigationScreen), marche libre & suivi d’activité (distance, durée, calories, lieux les plus visités).",
            "Sélecteur de lieux, partage de trajet, commentaires sur points de passage — le tout relié au compte Yukpo unique.",
        ],
        "revolution": [
            "Une seule app pour « se déplacer + acheter + se faire livrer » : moins de friction, plus de sécurité perçue.",
            "Les usagers gagnent du temps et de la clarté ; les petits commerces deviennent trouvables depuis la carte sans changer d’application.",
            "La voix et les alertes réduisent la charge cognitive en circulation (important pour motos, taxi, piétons).",
        ],
        "valeur": [
            "Revenus sur crédits, options premium, abonnement coaching lié à la navigation (leviers déjà câblés dans le produit).",
            "Forte fréquence d’usage = rétention et données de flux utiles pour cibler offres commerciales et pub géolocalisée.",
            "Différenciation forte vs simple marketplace : Yukpo devient un compagnon de mobilité au quotidien.",
        ],
    },
    {
        "title": "Santé & accès aux soins",
        "tag": "Orientation",
        "besoins": [
            "En cas d’urgence modérée ou de besoin courant, beaucoup de gens ne savent pas où aller en premier : pharmacie, clinique, spécialiste.",
            "Les créneaux et les tarifs sont souvent obtenus par appels répétés ; pas de vue d’ensemble sur la proximité réelle.",
            "Les professionnels de santé manquent de visibilité au-delà de leur clientèle habituelle.",
            "La prévention et le suivi (ordonnances, rappels) restent peu outillés pour le grand public.",
        ],
        "solutions": [
            "Annuaire géolocalisé avec types de structures et spécialités (dans le respect du cadre légal national).",
            "Prise de contact et prise de rendez-vous simplifiée lorsque le partenaire l’autorise.",
            "Rappels et documents : historique des demandes côté utilisateur pour faciliter le parcours de soins.",
            "Espace « santé bien-être » relié aux autres besoins (transport, livraison de médicaments si partenaires).",
        ],
        "revolution": [
            "Moins de errances et de coûts cachés avant d’être pris en charge.",
            "Les structures de proximité gagnent en affluence légitime plutôt qu’en concurrence opaque.",
            "Les familles gagnent du temps — ressource rare pour les parents et les aidants.",
        ],
        "valeur": [
            "Abonnements ou mises en avant pour les établissements et laboratoires partenaires.",
            "Fort volume de recherches : monétisation ciblée sans surcharger l’usager final.",
            "Image de confiance pour toute la marque Yukpo.",
        ],
    },
    {
        "title": "Immobilier",
        "tag": "Logement",
        "besoins": [
            "Loyers et colocations circulent sur des canaux non structurés : arnaques, doubles locations, photos trompeuses.",
            "Les étudiants et jeunes actifs paient souvent trop cher faute de comparateur fiable sur un même quartier.",
            "Les bailleurs perdent du temps à filtrer des candidats non sérieux.",
            "Les agences locales manquent d’outil simple pour digitaliser leur stock sans gros investissement.",
        ],
        "solutions": [
            "Annonces géolocalisées avec critères (loyer, chambres, meublé, proximité écoles / transports).",
            "Alertes lorsque le budget et le quartier correspondent.",
            "Messagerie intégrée et signalement pour réduire les abus.",
            "Parcours « bailleur » et « locataire » séparés mais reliés au même compte Yukpo.",
        ],
        "revolution": [
            "Plus de transparence sur le marché locatif : le locataire compare avant de se déplacer.",
            "Réduction des frais perdus dans des dossiers qui n’aboutissent pas.",
            "Meilleure rotation des logements disponibles, ce qui limite la vacance inutile.",
        ],
        "valeur": [
            "Frais de publication premium et partenariats avec agences.",
            "Montants moyens élevés par transaction : faible volume, forte contribution au revenu.",
            "Complément naturel au transport et à la vie pratique (déménagement, ménage).",
        ],
    },
    {
        "title": "Vie pratique",
        "tag": "Services du quotidien",
        "besoins": [
            "Réparations, cours à domicile, ménage, événements : tout repose sur le réseau personnel, donc long à trouver en urgence.",
            "Les prestataires indépendants manquent de visibilité stable ; leur chiffre d’affaires est irrégulier.",
            "Peu d’avis vérifiables : l’utilisateur hésite à faire entrer un inconnu chez lui.",
            "Les prix sont souvent négociés au téléphone sans référence de marché.",
        ],
        "solutions": [
            "Fiches prestataires avec spécialités, zone d’intervention, créneaux et tarifs indicatifs.",
            "Avis et historique de missions pour rassurer l’utilisateur.",
            "Liaison avec le paiement et, si besoin, la livraison de matériel.",
            "Suggestions croisées depuis le commerce (ex. achat d’électroménager → installation).",
        ],
        "revolution": [
            "Les services du quotidien deviennent aussi simples à commander qu’un plat ou un taxi.",
            "Les travailleurs indépendants gagnent une file d’attente de clients sans dépendre uniquement du bouche-à-oreille.",
            "Les ménages féminisés ou à temps très contraint gagnent des heures de vie chaque semaine.",
        ],
        "valeur": [
            "Commission sur prestations : revenu sur un très grand nombre de micro-transactions.",
            "Complète le panier e-commerce et renforce la rétention hebdomadaire.",
            "Données utiles pour cibler le marketing local (quartiers, types de ménages).",
        ],
    },
    {
        "title": "Bourse du livre scolaire",
        "tag": "Éducation",
        "besoins": [
            "La rentrée représente un choc de dépenses : manuels neufs coûteux, occasion difficile à trouver au bon niveau.",
            "Les familles revendent trop tard ou trop cher faute de lieu structuré pour l’occasion qualifiée.",
            "Les enseignants et parents manquent parfois d’une vue claire des programmes requis par niveau.",
            "Le gaspillage de livres encore utilisables est énorme d’une année sur l’autre.",
        ],
        "solutions": [
            "Matching par classe, programme et état du livre (neuf / très bon / bon).",
            "Messagerie intégrée et points de retrait possibles via partenaires (librairies, écoles).",
            "Rappels saisonniers et contenus « pré-rentrée » pour capter la demande au bon moment.",
            "Lien avec l’orientation scolaire et les frais connexes dans Yukpo.",
        ],
        "revolution": [
            "Réduction du coût réel de la scolarité pour les ménages modestes.",
            "Circulation du savoir sous forme de livre : moins de déchet, plus d’équité.",
            "Image forte pour Yukpo : une app utile à la vie des enfants, pas seulement commerciale.",
        ],
        "valeur": [
            "Commission sur transactions ; pics de revenus prévisibles chaque rentrée.",
            "Très fort partage social organique (parents, enseignants).",
            "Données sur les niveaux scolaires utiles au ciblage d’autres services (transport scolaire, assurances).",
        ],
    },
    {
        "title": "Supermarchés & grande distribution",
        "tag": "Retail moderne",
        "besoins": [
            "Les supermarchés ont du stock et des promotions, mais l’information utile n’arrive pas assez vite aux familles et aux quartiers où le besoin est le plus fort.",
            "Les ménages comparent mal les prix réels des paniers essentiels et perdent du temps en déplacements pour vérifier disponibilité et promotions.",
            "Les enseignes manquent souvent d’un canal unifié pour publier les offres, traiter les commandes locales, et orchestrer la livraison de proximité.",
            "Dans les périodes critiques (rentrée, fêtes, inflation), l’absence de visibilité sur les produits de base accentue la pression sur le pouvoir d’achat.",
        ],
        "solutions": [
            "Catalogue supermarché dynamique dans Yukpo : prix, promotions, disponibilité et comparaison de paniers en quelques clics.",
            "Commande assistée (texte, image, voix), retrait/livraison programmés et paiement mobile sécurisé dans un seul parcours.",
            "Back-office léger pour les enseignes : campagnes géolocalisées, mesure de conversion, suivi des ruptures sur produits sensibles.",
            "Parcours multi-segments : courses + livraison + services ménagers + mobilité, pour une expérience complète côté ménage.",
        ],
        "revolution": [
            "Le ménage reprend la main sur son budget alimentaire : plus de transparence prix, moins de déplacements inutiles, et décisions d’achat mieux informées.",
            "Les supermarchés deviennent accessibles au-delà de leur zone habituelle, avec une distribution plus inclusive dans des quartiers moins servis.",
            "Le segment produit un impact social visible : meilleure accessibilité des biens essentiels et baisse de la friction quotidienne.",
        ],
        "valeur": [
            "Commissions sur commandes, mises en avant promotionnelles, et abonnements enseignes pour outils de pilotage.",
            "Ticket moyen élevé et récurrence hebdomadaire : segment stabilisateur des revenus.",
            "Effet d’entraînement sur livraison, paiements et vie pratique, avec hausse de la rétention globale.",
        ],
    },
    {
        "title": "Assurance & protection",
        "tag": "Couverture",
        "besoins": [
            "Peu de comparabilité claire entre offres ; jargon incompréhensible pour une grande partie du public.",
            "Les TPE et auto-entrepreneurs sous-assurent leurs activités faute de parcours simple.",
            "En cas de sinistre, le suivi est souvent lourd : l’assuré ne sait pas où télécharger ou relancer.",
            "Les courtiers locaux manquent d’un canal digital moderne pour toucher les jeunes actifs.",
        ],
        "solutions": [
            "Vitrine de produits partenaires (auto, santé, habitation, activité) avec langage simplifié.",
            "Parcours de souscription guidé et rappels d’échéances dans l’app.",
            "Espace documents pour télécharger attestation ou contacter le partenaire.",
            "Lien naturel avec l’immobilier, le transport et le commerce (assurance marchand, RC pro).",
        ],
        "revolution": [
            "Démocratisation de la protection : moins de familles et de micro-entreprises « à découvert ».",
            "Comportement plus responsable (véhicule, logement) quand la couverture est accessible.",
            "Réduction de l’angoisse liée aux imprévus financiers pour les classes moyennes basses.",
        ],
        "valeur": [
            "Commissions distributeur sur primes : montants élevés par contrat.",
            "Une fois la confiance installée, fort potentiel de renouvellement annuel.",
            "Positionne Yukpo comme intermédiaire de confiance au-delà du commerce pur.",
        ],
    },
]

_LOGO_PATH: Path | None = None

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.45)
MARGIN_R = Inches(0.38)
TABLE_LEFT = Inches(0.48)
LOGO_H = Inches(0.42)
LOGO_MARGIN = Inches(0.2)
TABLE_W = Inches(10) - TABLE_LEFT - MARGIN_R - LOGO_H - LOGO_MARGIN
TITLE_W = TABLE_W
TITLE_TOP_IN = 0.22
TITLE_TOP = Inches(TITLE_TOP_IN)
SAFE_CONTENT_BOTTOM_IN = 6.42
FOOTER_BAND_TOP_IN = 6.5

NAVY = RGBColor(0x12, 0x2A, 0x4A)
ACCENT = RGBColor(0x1E, 0x5A, 0xB8)
ACCENT_LIGHT = RGBColor(0xE8, 0xEF, 0xFB)
ACCENT_HEADER = RGBColor(0x15, 0x3D, 0x7A)
GOLD = RGBColor(0xC9, 0x8A, 0x2E)
MUTED = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BOX_BG = RGBColor(0xF0, 0xF4, 0xFC)
CARD_LINE = RGBColor(0xB8, 0xC9, 0xE8)
CARD_HEAD_BG = RGBColor(0x15, 0x3D, 0x7A)


def resolve_logo_path() -> Path | None:
    root = Path(__file__).resolve().parent.parent
    c = root / "mobile" / "assets" / "icon.png"
    return c if c.is_file() else None


def apply_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFE)


def add_slide_top_bar(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0.0, 0.0, float(SLIDE_W), float(Inches(0.085)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()


def add_accent_sidebar(slide) -> None:
    """Bandeau vertical fort + léger dégradé visuel (impact)."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, float(MARGIN_L), float(Inches(0.88)), float(Inches(0.09)), float(Inches(5.55)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()
    sh2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        float(MARGIN_L) + float(Inches(0.09)),
        float(Inches(0.88)),
        float(Inches(0.03)),
        float(Inches(5.55)),
    )
    sh2.fill.solid()
    sh2.fill.fore_color.rgb = ACCENT_LIGHT
    sh2.line.fill.background()


def add_logo_br(slide) -> None:
    if _LOGO_PATH is None or not _LOGO_PATH.is_file():
        return
    left = float(SLIDE_W) - float(LOGO_H) - float(MARGIN_R)
    top = float(SLIDE_H) - float(LOGO_H) - float(MARGIN_R)
    slide.shapes.add_picture(str(_LOGO_PATH), left, top, height=LOGO_H)


def footer(slide, text: str | None = None) -> None:
    t = text or "Confidentiel — Yukpo"
    box = slide.shapes.add_textbox(
        float(Inches(1.15)),
        float(Inches(FOOTER_BAND_TOP_IN)),
        float(Inches(8.4)),
        float(Inches(0.32)),
    )
    p = box.text_frame.paragraphs[0]
    p.text = t
    p.font.size = Pt(FS_FOOTER)
    p.font.color.rgb = MUTED


def stamp(prs: Presentation) -> None:
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        box = slide.shapes.add_textbox(
            float(MARGIN_L),
            float(Inches(FOOTER_BAND_TOP_IN)),
            float(Inches(1.0)),
            float(Inches(0.28)),
        )
        p = box.text_frame.paragraphs[0]
        p.text = f"{i + 1} / {total}"
        p.font.size = Pt(FS_STAMP)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def max_body_height_inches(y_top: float) -> float:
    return max(0.35, SAFE_CONTENT_BOTTOM_IN - y_top)


def _estimate_wrapped_bullet_lines(lines: list[str], chars_per_line: float = 46.0) -> float:
    """Estime le nombre de lignes visibles (retours à la ligne) pour des puces étroites."""
    s = 0.0
    for t in lines:
        s += max(1.0, len(t) / chars_per_line)
    return s


def _module_card_height_in(besoins: list[str], solutions: list[str], y0: float) -> float:
    """Hauteur des cartes côte à côte : suit le contenu, sans remplir artificiellement la diapo."""
    n_left = _estimate_wrapped_bullet_lines(besoins)
    n_right = _estimate_wrapped_bullet_lines(solutions)
    n = max(n_left, n_right)
    body_h = 0.16 + 0.182 * max(4.0, min(n, 18.0))
    h = 0.52 + body_h
    cap = max_body_height_inches(y0) - 0.04
    return min(max(h, 2.28), cap)


def _module_card_height_in_part2(revolution: list[str], valeur: list[str], y0: float) -> float:
    n = max(_estimate_wrapped_bullet_lines(revolution), _estimate_wrapped_bullet_lines(valeur))
    body_h = 0.16 + 0.182 * max(4.0, min(n, 18.0))
    h = 0.52 + body_h
    cap = max_body_height_inches(y0) - 0.04
    return min(max(h, 2.28), cap)


def apply_bar_chart_style(chart) -> None:
    """Barres horizontales : couleur marque, labels lisibles, moins de vide autour du plot."""
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT
    dl = plot.data_labels
    dl.font.size = Pt(10.8)
    dl.font.bold = True
    dl.font.color.rgb = NAVY
    try:
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except (AttributeError, ValueError):
        pass
    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD8, 0xE2, 0xEF)
        chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    except (AttributeError, ValueError):
        pass
    try:
        chart.category_axis.tick_labels.font.size = Pt(10.5)
        chart.category_axis.tick_labels.font.color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
    except (AttributeError, ValueError):
        pass
    try:
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.color.rgb = MUTED
    except (AttributeError, ValueError):
        pass
    gw = getattr(plot, "gap_width", None)
    if gw is not None:
        try:
            plot.gap_width = 38
        except (AttributeError, ValueError):
            pass


def title_block(slide, title: str, subtitle: str | None = None) -> float:
    title_h_in = 1.18 if subtitle else 0.82
    h = Inches(title_h_in)
    box = slide.shapes.add_textbox(TABLE_LEFT, TITLE_TOP, TITLE_W, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_bottom = Pt(4)
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(FS_H1)
    p0.font.bold = True
    p0.font.color.rgb = NAVY
    if subtitle:
        p1 = tf.add_paragraph()
        p1.text = subtitle
        p1.font.size = Pt(FS_SUBTITLE)
        p1.font.color.rgb = MUTED
        p1.space_before = Pt(5)
    return TITLE_TOP_IN + title_h_in + 0.1


def slide_blank(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    add_slide_top_bar(slide)
    return slide


def add_kpi_row(
    slide,
    y_in: float,
    items: list[tuple[str, str, str]],
    foot: str | None = None,
) -> None:
    n = len(items)
    gap = Inches(0.12)
    total_w = float(TABLE_W) - float(gap) * (n - 1)
    w = total_w / n
    x0 = float(TABLE_LEFT)
    y = Inches(y_in)
    h = Inches(2.25)
    for i, (big, lab, sub) in enumerate(items):
        x = x0 + i * (w + float(gap))
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = BOX_BG
        sh.line.color.rgb = CARD_LINE
        sh.line.width = Pt(1)
        tf = sh.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.text = big
        p0.font.size = Pt(FS_KPI_BIG + 2 if len(big) < 14 else max(FS_KPI_BIG - 4, 26))
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        p0.alignment = PP_ALIGN.CENTER
        p1 = tf.add_paragraph()
        p1.text = lab
        p1.font.size = Pt(FS_KPI_LAB)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(8)
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(FS_KPI_SUB)
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(5)
    if foot:
        bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(y_in + 2.38), TABLE_W, Inches(0.42))
        bx.text_frame.paragraphs[0].text = foot
        bx.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
        bx.text_frame.paragraphs[0].font.color.rgb = MUTED
        bx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_two_kpi_wide(slide, y_in: float, left: tuple[str, str, str], right: tuple[str, str, str]) -> None:
    w = (float(TABLE_W) - float(Inches(0.2))) / 2
    x0 = float(TABLE_LEFT)
    y = Inches(y_in)
    h = Inches(2.4)
    for i, (big, lab, sub) in enumerate((left, right)):
        x = x0 + i * (w + float(Inches(0.2)))
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = NAVY
        sh.line.fill.background()
        tf = sh.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.text = big
        p0.font.size = Pt(FS_COVER_KPI + 2)
        p0.font.bold = True
        p0.font.color.rgb = WHITE
        p0.alignment = PP_ALIGN.CENTER
        p1 = tf.add_paragraph()
        p1.text = lab
        p1.font.size = Pt(FS_KPI_LAB)
        p1.font.color.rgb = RGBColor(0xDD, 0xE5, 0xF5)
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(10)
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(FS_KPI_SUB)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)


def add_chart_column(
    slide,
    title: str,
    categories: list[str],
    series_name: str,
    values: list[float],
    y_sub: str | None,
    unit_note: str,
) -> None:
    yb = title_block(slide, title, y_sub)
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    n = max(len(categories), 1)
    # Barres horizontales : hauteur liée au nombre de catégories (moins d’espace vide)
    cy_in = min(4.25, max(2.65, 0.38 * n + 1.35))
    x = float(TABLE_LEFT)
    y = float(Inches(yb + 0.05))
    cx = float(TABLE_W)
    cy = float(Inches(cy_in))
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    apply_bar_chart_style(chart)
    note_y = float(yb) + 0.05 + cy_in + 0.06
    tb = slide.shapes.add_textbox(TABLE_LEFT, Inches(note_y), TABLE_W, Inches(0.42))
    tb.text_frame.paragraphs[0].text = unit_note
    tb.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    tb.text_frame.paragraphs[0].font.color.rgb = MUTED


def add_slide_market_informal(prs: Presentation) -> None:
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    yb = title_block(
        slide,
        "Marché informel & besoin réel",
        "Une opportunité de structuration, pas seulement « du digital »",
    )
    bx = slide.shapes.add_textbox(
        TABLE_LEFT,
        Inches(yb + 0.1),
        TABLE_W,
        Inches(max_body_height_inches(yb + 0.1)),
    )
    tf = bx.text_frame
    tf.word_wrap = True
    lines = [
        "Une part majeure de l’activité économique repose sur des échanges peu ou pas digitalisés : visibilité faible, prix et délais peu comparables.",
        "Les ménages jonglent entre plusieurs canaux (boutique, téléphone, réseaux sociaux) pour se loger, se soigner, transporter des biens ou des personnes.",
        "La demande existe : ce qui manque souvent, c’est une couche simple, en français, mobile, qui relie l’offre locale à la demande sans imposer un coût d’entrée prohibitif aux petits acteurs.",
        "Yukpo vise à capter cette demande en regroupant des usages du quotidien dans une seule application, avec un modèle où chaque segment peut contribuer aux revenus.",
    ]
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = "▸  " + line
        p.font.size = Pt(FS_BODY)
        p.space_after = Pt(10)
    footer(slide)
    add_logo_br(slide)


def _fill_card_with_bullets(tf, lines: list[str], bullet: str = "•") -> None:
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"{bullet}  {line}"
        p.font.size = Pt(FS_CARD_BULLET)
        p.font.color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
        p.space_after = Pt(6)
        p.line_spacing = 1.14


def _short_text(s: str, limit: int = 240) -> str:
    txt = " ".join(s.split())
    if not txt:
        return txt
    if len(txt) <= limit:
        return txt if txt[-1] in ".!?" else f"{txt}."
    cut = txt[:limit].rstrip()
    if " " in cut:
        cut = cut.rsplit(" ", 1)[0]
    # force une fin de phrase propre au lieu d'une coupe brute
    return cut + "."


def _take_complete_sentences(lines: list[str], max_chars: int = 240, max_sentences: int = 2) -> str:
    """Construit un résumé court avec phrases complètes uniquement."""
    if not lines:
        return ""
    buf = " ".join(" ".join(l.split()) for l in lines if l.strip())
    if not buf:
        return ""
    parts = re.split(r"(?<=[.!?])\s+", buf)
    picked: list[str] = []
    size = 0
    for sent in parts:
        st = sent.strip()
        if not st:
            continue
        if st[-1] not in ".!?":
            st += "."
        add_len = len(st) + (1 if picked else 0)
        if picked and size + add_len > max_chars:
            break
        picked.append(st)
        size += add_len
        if len(picked) >= max_sentences:
            break
    if picked:
        return " ".join(picked)
    return _short_text(buf, max_chars)


def _segment_flow_blocks(m: dict) -> list[tuple[str, str, RGBColor]]:
    besoins = list(m["besoins"])
    solutions = list(m["solutions"])
    probleme = _take_complete_sentences(besoins[:2], 250, 2) if besoins else ""
    besoin_reel = _take_complete_sentences(besoins[2:4] if len(besoins) > 2 else besoins, 250, 2)
    solution = _take_complete_sentences(solutions[:2], 250, 2) if solutions else ""
    return [
        ("Problématique terrain", probleme, ACCENT_HEADER),
        ("Besoin utilisateur réel", besoin_reel, GOLD),
        ("Solution digitale Yukpo", solution, ACCENT),
    ]


def _draw_segment_flow(slide, x: float, y: float, w: float, h: float, title: str, body: str, color: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.4)
    hb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, float(Inches(0.48)))
    hb.fill.solid()
    hb.fill.fore_color.rgb = color
    hb.line.fill.background()
    ht = slide.shapes.add_textbox(x + float(Inches(0.08)), y + float(Inches(0.06)), w - float(Inches(0.16)), float(Inches(0.34)))
    ht.text_frame.paragraphs[0].text = title
    ht.text_frame.paragraphs[0].font.bold = True
    ht.text_frame.paragraphs[0].font.size = Pt(FS_CARD_TITLE)
    ht.text_frame.paragraphs[0].font.color.rgb = WHITE

    bd = slide.shapes.add_textbox(x + float(Inches(0.1)), y + float(Inches(0.56)), w - float(Inches(0.2)), h - float(Inches(0.64)))
    tf = bd.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(FS_CARD_BULLET + 0.4)
    p.font.color.rgb = RGBColor(0x2A, 0x2A, 0x2A)
    p.line_spacing = 1.16


def _segment_kpi_line(title: str) -> str:
    kpi_by_title = {
        "E-commerce & vitrines marchands": "KPI cible : +30 à +45 % de commandes digitalisées chez les marchands actifs en 12 mois.",
        "Livraison": "KPI cible : délai moyen de livraison réduit de 20 à 35 % dans les zones densifiées.",
        "Transport & mobilité": "KPI cible : +25 % de trajets planifiés via app et baisse des trajets à vide partenaires.",
        "Navigation intelligente (NavigationScreen)": "KPI cible : +15 à +25 % de rétention mensuelle via usage navigation récurrent.",
        "Santé & accès aux soins": "KPI cible : réduction du temps d’orientation vers un point de soin pertinent (suivi en minutes).",
        "Immobilier": "KPI cible : baisse des visites non qualifiées et amélioration du taux contact-annonce.",
        "Vie pratique": "KPI cible : hausse de la fréquence hebdomadaire des services du quotidien par ménage actif.",
        "Bourse du livre scolaire": "KPI cible : économie mesurable sur le panier rentrée par famille et forte acquisition saisonnière.",
        "Supermarchés & grande distribution": "KPI cible : hausse du panier moyen et amélioration de la transparence prix sur produits essentiels.",
        "Assurance & protection": "KPI cible : progression du taux de souscription/renouvellement sur offres partenaires simplifiées.",
    }
    return kpi_by_title.get(title, "KPI cible : progression conjointe usage, conversion et rétention sur le segment.")


def _segment_investor_story(m: dict) -> str:
    besoins = list(m["besoins"])
    solutions = list(m["solutions"])
    rev = list(m["revolution"])
    douleur = _short_text(besoins[0] if besoins else "Besoin terrain insuffisamment couvert.", 165)
    transfo = _short_text(solutions[0] if solutions else "Réponse digitale Yukpo sur le segment.", 165)
    impact = _short_text(rev[0] if rev else "Impact social et opérationnel positif.", 165)
    kpi = _segment_kpi_line(str(m.get("title", "")))
    return f"Douleur marché : {douleur}  |  Transformation Yukpo : {transfo}  |  Impact chiffrable : {impact} {kpi}"


def _segment_pitch_10s(title: str) -> str:
    pitches = {
        "E-commerce & vitrines marchands": "Pitch 10s : on transforme les boutiques invisibles en vitrines monétisables, avec plus de ventes locales traçables.",
        "Livraison": "Pitch 10s : on rend la livraison fiable et mesurable, donc plus de commandes finalisées et plus de revenus pour les coursiers.",
        "Transport & mobilité": "Pitch 10s : on réduit la friction du déplacement quotidien et on convertit ce flux en usage récurrent de la super-app.",
        "Navigation intelligente (NavigationScreen)": "Pitch 10s : on capte un usage quotidien à forte rétention, puis on le convertit en commerce et services géolocalisés.",
        "Santé & accès aux soins": "Pitch 10s : on raccourcit le chemin vers le bon soin, ce qui crée confiance, récurrence et valeur de marque durable.",
        "Immobilier": "Pitch 10s : on fluidifie un marché opaque et on monétise la mise en relation qualifiée sur des tickets élevés.",
        "Vie pratique": "Pitch 10s : on digitalise les services du quotidien, augmentant la fréquence d’usage et la valeur hebdomadaire par ménage.",
        "Bourse du livre scolaire": "Pitch 10s : on réduit le coût de la rentrée pour les familles et on active un puissant moteur d’acquisition saisonnier.",
        "Supermarchés & grande distribution": "Pitch 10s : on connecte le panier essentiel au digital local, avec impact pouvoir d’achat et revenus récurrents.",
        "Assurance & protection": "Pitch 10s : on simplifie la protection des ménages/TPE et on crée un revenu à renouvellement annuel.",
    }
    return pitches.get(title, "Pitch 10s : Yukpo résout un besoin quotidien réel et le convertit en revenu récurrent mesurable.")


def add_slide_module_part1(prs: Presentation, m: dict) -> None:
    """Segment — diapo 1 : besoins du terrain + solutions Yukpo."""
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    sub = f"{m['tag']} · Chaîne visuelle : problème → besoin réel → solution Yukpo"
    yb = title_block(slide, str(m["title"]), str(sub))
    blocks = _segment_flow_blocks(m)
    gap = float(Inches(0.12))
    w3 = (float(TABLE_W) - 2 * gap) / 3.0
    x0 = float(TABLE_LEFT)
    y0 = float(yb) + 0.1
    h = min(max_body_height_inches(y0) - 0.06, 3.9)

    for idx, (head, body, color) in enumerate(blocks):
        x = x0 + idx * (w3 + gap)
        _draw_segment_flow(slide, x, float(Inches(y0)), w3, float(Inches(h)), head, body, color)
        # Badge de séquence pour un storytelling plus "deck investisseur"
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + float(Inches(0.02)),
            float(Inches(y0)) - float(Inches(0.16)),
            float(Inches(0.24)),
            float(Inches(0.24)),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = str(idx + 1)
        bp.font.bold = True
        bp.font.size = Pt(10)
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER
        if idx < len(blocks) - 1:
            ax = x + w3 + float(Inches(0.02))
            arrow = slide.shapes.add_textbox(ax, float(Inches(y0 + 1.52)), float(Inches(0.08)), float(Inches(0.4)))
            ap = arrow.text_frame.paragraphs[0]
            ap.text = "→"
            ap.font.size = Pt(18)
            ap.font.bold = True
            ap.font.color.rgb = ACCENT

    nb = slide.shapes.add_textbox(TABLE_LEFT, Inches(y0 + h + 0.06), TABLE_W, Inches(0.48))
    np = nb.text_frame.paragraphs[0]
    np.text = (
        "Lecture investisseur : chaque segment part d’un besoin social réel, le transforme en usage mesurable, "
        "puis active plusieurs solutions digitales Yukpo avec impact direct sur ménages, acteurs locaux et revenus."
    )
    np.font.size = Pt(FS_NOTE - 0.5)
    np.font.color.rgb = MUTED
    footer(slide)
    add_logo_br(slide)


def add_slide_module_part2(prs: Presentation, m: dict) -> None:
    """Segment — diapo 2 : impact population + valeur économique."""
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    sub = f"{m['tag']} · Impact social & modèle économique"
    yb = title_block(slide, str(m["title"]), str(sub))
    gap = float(Inches(0.14))
    w2 = (float(TABLE_W) - gap) / 2.0
    x0 = float(TABLE_LEFT)
    y0 = float(yb) + 0.1
    h = min(_module_card_height_in_part2(list(m["revolution"]), list(m["valeur"]), y0), 3.58)
    blocks: list[tuple[str, list[str], RGBColor]] = [
        ("Impact pour la population", list(m["revolution"]), GOLD),
        ("Valeur économique & leviers de revenus", list(m["valeur"]), NAVY),
    ]
    for k, (head, lines, head_c) in enumerate(blocks):
        x = x0 + k * (w2 + gap)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y0), w2, Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = WHITE
        sh.line.color.rgb = ACCENT
        sh.line.width = Pt(1.15)
        hb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(y0), w2, Inches(0.42))
        hb.fill.solid()
        hb.fill.fore_color.rgb = head_c
        hb.line.fill.background()
        ht = slide.shapes.add_textbox(x + float(Inches(0.08)), Inches(y0 + 0.06), w2 - float(Inches(0.16)), Inches(0.34))
        ht.text_frame.paragraphs[0].text = head
        ht.text_frame.paragraphs[0].font.bold = True
        ht.text_frame.paragraphs[0].font.size = Pt(FS_CARD_TITLE)
        ht.text_frame.paragraphs[0].font.color.rgb = WHITE
        ht.text_frame.word_wrap = True
        body = slide.shapes.add_textbox(x + float(Inches(0.12)), Inches(y0 + 0.48), w2 - float(Inches(0.24)), Inches(h - 0.52))
        tf = body.text_frame
        tf.word_wrap = True
        _fill_card_with_bullets(tf, lines)

    story = slide.shapes.add_textbox(TABLE_LEFT, Inches(y0 + h + 0.06), TABLE_W, Inches(0.72))
    stf = story.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = _segment_investor_story(m)
    sp.font.size = Pt(FS_NOTE - 0.5)
    sp.font.bold = True
    sp.font.color.rgb = NAVY

    pitch = slide.shapes.add_textbox(TABLE_LEFT, Inches(y0 + h + 0.62), TABLE_W, Inches(0.34))
    pp = pitch.text_frame.paragraphs[0]
    pp.text = _segment_pitch_10s(str(m.get("title", "")))
    pp.font.size = Pt(FS_NOTE - 0.8)
    pp.font.color.rgb = RGBColor(0x1B, 0x4D, 0x8F)
    pp.font.bold = True
    footer(slide)
    add_logo_br(slide)


def add_slide_financing_overview(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Plan de financement — vue d’ensemble",
        f"{B['total_m']} millions FCFA sur 20 mois · lancement effectif juillet 2026",
    )
    bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(yb + 0.1), TABLE_W, Inches(4.9))
    tf = bx.text_frame
    tf.word_wrap = True
    blocks = [
        "Les fonds servent avant tout à faire connaître l’application et à recruter les bons profils pour activer chaque segment (commerce, livraison, santé, etc.).",
        "Le budget marketing et le budget « équipe » sont volontairement élevés : ce sont les deux leviers qui transforment un produit déjà construit en usage réel sur le terrain.",
        "Technologie : budget qui augmente avec la charge (hébergement, sécurité, évolutions) — voir répartition sur deux ans.",
        "Calendrier : dépenses à partir de mai 2026 ; application en production à partir de juillet 2026 ; suivi budgétaire sur 20 mois (jusqu’à fin 2027).",
        f"Pour le capital : toute entrée d’investisseur reste discutée avec un plafond de {B['equity_cap_pct']} % pour protéger le contrôle du projet.",
    ]
    for i, line in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸  " + line
        p.font.size = Pt(FS_BODY)
        p.space_after = Pt(11)
    footer(slide)
    add_logo_br(slide)


def add_slide_financing_options(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Comment financer les 200 millions ?",
        "Plusieurs combinaisons possibles — à valider avec votre banque et vos conseils",
    )
    bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(yb + 0.12), TABLE_W, Inches(max_body_height_inches(yb + 0.12)))
    tf = bx.text_frame
    tf.word_wrap = True
    items = [
        (
            "Prêt ou ligne de crédit",
            "Vous empruntez une partie de la somme et vous remboursez selon un calendrier convenu. Avantage : pas de partage du capital si le financement est entièrement en dette.",
        ),
        (
            "Caution ou garantie bancaire",
            "Une banque ou un tiers peut garantir votre emprunt pour réduire le risque perçu. À prévoir : coût de la garantie et conditions précises.",
        ),
        (
            "Investisseur au capital (minoritaire)",
            f"Échange d’argent contre une part limitée des actions — dans ce dossier, plafonnée à {B['equity_cap_pct']} % pour encadrer la dilution.",
        ),
        (
            "Mélange des trois",
            "Exemple : une partie en prêt, une petite part au capital, une garantie sur le reste — à modéliser avec un professionnel.",
        ),
    ]
    for idx, (h, body) in enumerate(items):
        ph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        ph.text = h
        ph.font.bold = True
        ph.font.size = Pt(FS_BODY + 1)
        ph.font.color.rgb = NAVY
        ph.space_before = Pt(0 if idx == 0 else 14)
        pb = tf.add_paragraph()
        pb.text = body
        pb.font.size = Pt(FS_BODY)
        pb.space_after = Pt(4)
    footer(slide)
    add_logo_br(slide)


def add_slide_charges_two_years(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Répartition des charges sur deux ans",
        "Mai 2026 – avril 2028 (cohérent avec l’enveloppe 200 M)",
    )
    # Tech & personnel augmentent en 2ᵉ année (montée en charge, outils, équipe élargie) ; marketing relativement plus fort en 1ʳᵉ année
    rows = [
        ("Marketing & acquisition", "78 M", "40 M"),
        ("Personnel & rémunérations", "18 M", "24 M"),
        ("Technologie & infrastructure", "7 M", "11 M"),
        ("Juridique & conformité", "4 M", "4 M"),
        ("Frais opérationnels & réserve", "11 M", "3 M"),
        ("TOTAL charges", f"{CHARGES_ANNEE1_M} M", f"{CHARGES_ANNEE2_M} M"),
    ]
    top = Inches(yb + 0.08)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.44 * nrows))
    table = shape.table
    for j, w in enumerate((Inches(3.5), Inches(2.35), Inches(2.35))):
        table.columns[j].width = w
    hdr = ("Poste de charge", "1ʳᵉ année (mai 26 – avr. 27)", "2ᵉ année (mai 27 – avr. 28)")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY)
            if i == nrows - 1:
                c.text_frame.paragraphs[0].font.bold = True
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF5)
            elif i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    note = slide.shapes.add_textbox(TABLE_LEFT, Inches(5.55), TABLE_W, Inches(0.75))
    note.text_frame.paragraphs[0].text = (
        "Logique : gros investissement acquisition en 1ʳᵉ année ; en 2ᵉ année, hausse volontaire des coûts tech (infra, sécurité, évolutions) "
        "et des salaires (renforts terrain & support) tandis que le marketing peut être optimisé grâce à la notoriété."
    )
    note.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    note.text_frame.paragraphs[0].font.color.rgb = MUTED
    footer(slide)
    add_logo_br(slide)


def _add_quarterly_table_at(
    slide,
    y_top: float,
    labels: list[str],
    rev: list[float],
    chg: list[float],
) -> None:
    solde = [round(rev[i] - chg[i], 1) for i in range(len(labels))]
    top = Inches(y_top)
    ncol = len(labels) + 1
    nrows = 4
    shape = slide.shapes.add_table(nrows, ncol, TABLE_LEFT, top, TABLE_W, Inches(0.42 * nrows))
    table = shape.table
    hdr = ["Indicateur (M)", *labels]
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    row_data = [
        ("Revenus", [str(x) for x in rev]),
        ("Charges", [str(x) for x in chg]),
        ("Résultat", [str(s) for s in solde]),
    ]
    for ri, (label, vals) in enumerate(row_data, start=1):
        table.cell(ri, 0).text = ""
        table.cell(ri, 0).text_frame.paragraphs[0].text = label
        table.cell(ri, 0).text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY)
        table.cell(ri, 0).text_frame.paragraphs[0].font.bold = True
        for j, v in enumerate(vals, start=1):
            table.cell(ri, j).text = ""
            table.cell(ri, j).text_frame.paragraphs[0].text = v
            table.cell(ri, j).text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY)


def add_slide_simulation_quarters_1(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Simulation trimestrielle (1/2)",
        "Juil. 2026 – juin 2027 · millions FCFA",
    )
    _add_quarterly_table_at(
        slide,
        yb + 0.08,
        ["T3 26", "T4 26", "T1 27", "T2 27"],
        TRIM_REVENUS_M[:4],
        TRIM_CHARGES_M[:4],
    )
    note = slide.shapes.add_textbox(TABLE_LEFT, Inches(5.1), TABLE_W, Inches(0.95))
    note.text_frame.paragraphs[0].text = (
        "Base de calcul revenus : volume de transactions × panier moyen × taux de commission net + abonnements marchands. "
        "Période alignée avec le bilan annuel : T3 2026 à T2 2027."
    )
    note.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    note.text_frame.paragraphs[0].font.color.rgb = MUTED
    footer(slide)
    add_logo_br(slide)


def add_slide_simulation_quarters_2(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Simulation trimestrielle (2/2)",
        "Juil. 2027 – juin 2028 · accélération des revenus",
    )
    _add_quarterly_table_at(
        slide,
        yb + 0.08,
        ["T3 27", "T4 27", "T1 28", "T2 28"],
        TRIM_REVENUS_M[4:8],
        TRIM_CHARGES_M[4:8],
    )
    cum_rev = sum(TRIM_REVENUS_M)
    cum_chg = sum(TRIM_CHARGES_M)
    note = slide.shapes.add_textbox(TABLE_LEFT, Inches(5.0), TABLE_W, Inches(1.05))
    note.text_frame.paragraphs[0].text = (
        f"Sur les 8 trimestres : revenus cumulés ≈ {cum_rev} M ; charges d’activité ≈ {cum_chg} M "
        "(cohérent avec les deux bilans annuels). "
        "Hypothèse de revenu conservatrice : commissions moyennes 8–12 % selon segment, panier progressif avec densité utilisateurs."
    )
    note.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    note.text_frame.paragraphs[0].font.color.rgb = MUTED
    footer(slide)
    add_logo_br(slide)


def add_slide_acquisition_croissance(prs: Presentation) -> None:
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    yb = title_block(
        slide,
        "Acquisition & traction",
        "Objectifs indicatifs — pilotage mensuel recommandé",
    )
    rows = [
        ("Utilisateurs actifs cumulés (milliers)", "≈ 28", "≈ 95"),
        ("Commerçants / vendeurs actifs", "≈ 2 200", "≈ 8 500"),
        ("Commandes & prestations / mois (milliers, fin période)", "≈ 38", "≈ 145"),
        ("Revenus annuels agrégés (M FCFA)", "≈ 62", "≈ 106"),
    ]
    top = Inches(yb + 0.1)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.46 * nrows))
    table = shape.table
    hdr = ("Indicateur", "Fin 1ʳᵉ année d’exploitation", "Fin 2ᵉ année d’exploitation")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY + 0.5)
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(5.35), TABLE_W, Inches(0.85))
    bx.text_frame.paragraphs[0].text = (
        "La densité d’utilisateurs et de marchands fait croître les revenus plus vite que les coûts fixes — "
        "levier central pour dépasser le seuil de rentabilité en 2ᵉ année. "
        "Base : revenus = (transactions/mois × panier moyen × commission nette) + abonnements et mises en avant."
    )
    bx.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    bx.text_frame.paragraphs[0].font.color.rgb = MUTED
    footer(slide)
    add_logo_br(slide)


def add_slide_revenus_base_calcul(prs: Presentation) -> None:
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    yb = title_block(
        slide,
        "Comment les revenus sont évalués",
        "Méthode de calcul utilisée pour les simulations trimestrielles et le bilan annuel",
    )
    bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(yb + 0.08), TABLE_W, Inches(4.95))
    tf = bx.text_frame
    tf.word_wrap = True
    lines = [
        "Formule principale : revenus = (transactions validées × panier moyen × commission nette) + abonnements partenaires + mises en avant.",
        "Commissions nettes retenues : intervalle prudent de 8 % à 12 % selon le segment (commerce, livraison, transport, immobilier, assurance).",
        "Montée en charge : progression du volume trimestriel via l’augmentation des utilisateurs actifs, marchands actifs et commandes mensuelles.",
        "Cohérence stricte : somme des 4 trimestres année 1 = 62 M et année 2 = 106 M, identique au bilan annuel.",
        "Hypothèse conservatrice : pas de revenus exceptionnels intégrés, uniquement revenus d’exploitation récurrents.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸  " + line
        p.font.size = Pt(FS_BODY)
        p.space_after = Pt(10)
    footer(slide)
    add_logo_br(slide)


def add_slide_bilan_rentabilite(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Bilan par exercice — rentabilité en 2ᵉ année",
        "Hypothèses de synthèse (mai 2026 – avril 2028)",
    )
    rows = [
        ("Revenus agrégés", "62 M", "106 M"),
        ("Charges totales", f"{CHARGES_ANNEE1_M} M", f"{CHARGES_ANNEE2_M} M"),
        ("Résultat (revenus – charges)", "−56 M", "+24 M"),
    ]
    top = Inches(yb + 0.12)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.52 * nrows))
    table = shape.table
    hdr = ("", "1ʳᵉ année (mai 26 – avr. 27)", "2ᵉ année (mai 27 – avr. 28)")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_BODY)
            if j == 0:
                c.text_frame.paragraphs[0].font.bold = True
            if i == nrows - 1 and j > 0:
                c.text_frame.paragraphs[0].font.bold = True
                c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x0F, 0x5A, 0x2E)
                if val.startswith("+"):
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    note = slide.shapes.add_textbox(TABLE_LEFT, Inches(4.85), TABLE_W, Inches(1.35))
    note.text_frame.paragraphs[0].text = (
        "À la fin de la 2ᵉ année complète, les revenus dépassent les charges : l’entreprise peut être "
        "rentable sur cette base (hors nouveaux tours de financement de croissance)."
    )
    note.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    note.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p2 = note.text_frame.add_paragraph()
    p2.text = "Ordres de grandeur pour discussion — à confirmer avec la comptabilité de gestion."
    p2.font.size = Pt(FS_NOTE - 0.5)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(8)
    footer(slide)
    add_logo_br(slide)


def add_slide_marketing_depenses_lignes(prs: Presentation) -> None:
    """Lignes de dépenses marketing (sur le poste 118 M) — poches budgétaires."""
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Déploiement marketing — lignes de dépenses",
        f"Poste acquisition & notoriété : {B['marketing_m']} M FCFA (enveloppe 200 M)",
    )
    mk = B["marketing_m"]
    rows = [
        ("Acquisition digitale (social, UAC, ASO, remarketing)", f"≈ {round(mk * 0.32)} M", "≈ 32 %"),
        ("Terrain & ambassadeurs (villes pilotes, micro-influence locale)", f"≈ {round(mk * 0.24)} M", "≈ 24 %"),
        ("Contenus & création (vidéos, studios légers, UGC)", f"≈ {round(mk * 0.13)} M", "≈ 13 %"),
        ("Partenariats médias, écoles, radios, associations", f"≈ {round(mk * 0.10)} M", "≈ 10 %"),
        ("Événements & saisons fortes (rentrée, fêtes, salons)", f"≈ {round(mk * 0.09)} M", "≈ 9 %"),
        ("Outils, CRM acquisition, mesure & analytics", f"≈ {round(mk * 0.07)} M", "≈ 7 %"),
        ("Marge de manœuvre / tests A-B", f"≈ {round(mk * 0.05)} M", "≈ 5 %"),
    ]
    top = Inches(yb + 0.08)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.4 * nrows))
    table = shape.table
    hdr = ("Ligne de dépense", "Montant indicatif", "Part du poste marketing")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY)
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    note = slide.shapes.add_textbox(TABLE_LEFT, Inches(5.75), TABLE_W, Inches(0.65))
    note.text_frame.paragraphs[0].text = (
        "Les parts sont indicatives : elles permettent d’aligner équipes et investisseurs sur les « poches » "
        "où l’argent part réellement (digital vs terrain vs partenariats)."
    )
    note.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    note.text_frame.paragraphs[0].font.color.rgb = MUTED
    footer(slide)
    add_logo_br(slide)


def add_slide_marketing_poches_consommation(prs: Presentation) -> None:
    """Poches de consommation & impact sur la croissance de l’app."""
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    yb = title_block(
        slide,
        "Poches de consommation & levier de croissance",
        "Où dépense l’utilisateur — comment ça fait bouger les indicateurs",
    )
    bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(yb + 0.08), TABLE_W, Inches(5.15))
    tf = bx.text_frame
    tf.word_wrap = True
    blocks = [
        (
            "Jeunes actifs & étudiants",
            "Livraison rapide, mobilité, navigation, bourse du livre, vie pratique. Impact : forte fréquence d’ouverture, partage social, acquisition virale peu chère.",
        ),
        (
            "Ménages & familles",
            "Courses, santé, logement, assurance, école. Impact : panier plus élevé par session, récurrence hebdomadaire, confiance long terme.",
        ),
        (
            "Commerçants & TPE",
            "E-commerce, visibilité, livraison. Impact : liquidité sur la plateforme, commissions, preuve de traction pour les partenaires institutionnels.",
        ),
        (
            "Usagers de la mobilité quotidienne",
            "Navigation intelligente, transport, trajets. Impact : habitude quotidienne (rétention), données de flux pour ciblage publicitaire et offres locales.",
        ),
        (
            "Saison scolaire & rentrée",
            "Bourse du livre, orientation. Impact : pics de téléchargement prévisibles, image « utile socialement », fort bouche-à-oreille parental.",
        ),
    ]
    for idx, (h, body) in enumerate(blocks):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = h
        p.font.bold = True
        p.font.size = Pt(FS_BODY)
        p.font.color.rgb = NAVY
        p.space_before = Pt(12 if idx else 0)
        pb = tf.add_paragraph()
        pb.text = body
        pb.font.size = Pt(FS_NOTE)
        pb.space_after = Pt(4)
    footer(slide)
    add_logo_br(slide)


def add_slide_deploiement_geographique(prs: Presentation) -> None:
    """Évolution Douala / Yaoundé / Bafoussam puis Côte d’Ivoire puis Sénégal."""
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Déploiement géographique — calendrier",
        "Cameroun d’abord · extension Afrique de l’Ouest ensuite",
    )
    rows = [
        (
            "T3 2026",
            "Douala & Yaoundé",
            "Lancement national priorisé sur les deux pôles économiques : densification commerçants, livraison, navigation, acquisition digitale + terrain.",
        ),
        (
            "T4 2026",
            "Côte d’Ivoire (Abidjan)",
            "Démarrage extension internationale avec noyau opérationnel local ; adaptation marketing et conformité avant montée en charge.",
        ),
        (
            "T1 – T2 2027",
            "Bafoussam + consolidation",
            "Ouverture de l’ouest camerounais et renforcement des zones ouvertes (Cameroun + Abidjan) : preuve de scalabilité multi-villes.",
        ),
        (
            "T3 – T4 2027",
            "Côte d’Ivoire (extension) + optimisation",
            "Extension hors Abidjan et optimisation unit economics sur les marchés déjà ouverts.",
        ),
        (
            "2028 (cible)",
            "Sénégal",
            "Après stabilisation Cameroun + Côte d’Ivoire : hub Dakar avec le même playbook produit/marketing adapté.",
        ),
    ]
    top = Inches(yb + 0.08)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.52 * nrows))
    table = shape.table
    for j, w in enumerate((Inches(2.05), Inches(1.85), Inches(4.45))):
        table.columns[j].width = w
    hdr = ("Période", "Zone", "Objectif opérationnel")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY)
            c.text_frame.word_wrap = True
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    note = slide.shapes.add_textbox(TABLE_LEFT, Inches(5.85), TABLE_W, Inches(0.55))
    note.text_frame.paragraphs[0].text = (
        "L’ordre Cameroun → Côte d’Ivoire → Sénégal maximise la proximité culturelle et linguistique avant élargissement UEMOA / régional."
    )
    note.text_frame.paragraphs[0].font.size = Pt(FS_NOTE)
    note.text_frame.paragraphs[0].font.color.rgb = MUTED
    footer(slide)
    add_logo_br(slide)


def add_slide_marketing_1(prs: Presentation) -> None:
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    yb = title_block(
        slide,
        "Stratégie marketing — principes",
        "Toucher chaque segment sans disperser le message",
    )
    bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(yb + 0.1), TABLE_W, Inches(4.95))
    tf = bx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(
        [
            "Un seul nom de marque (Yukpo), plusieurs « portes d’entrée » : commerce, livraison, santé, logement, etc.",
            "Contenus courts et concrets (vidéos, témoignages marchands, avant/après) plutôt que jargon technologique.",
            "Acquisition payante ciblée par ville ou par quartier pour densifier l’offre avant d’élargir.",
            "Partenariats : écoles, associations, commerçants influents, radios locales — pour emprunter la confiance déjà installée.",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸  " + line
        p.font.size = Pt(FS_BODY + 0.5)
        p.space_after = Pt(12)
    footer(slide)
    add_logo_br(slide)


def add_slide_marketing_2(prs: Presentation) -> None:
    slide = slide_blank(prs)
    add_accent_sidebar(slide)
    yb = title_block(
        slide,
        "Déploiement marketing — canaux",
        "Digital + terrain + animation continue",
    )
    bx = slide.shapes.add_textbox(TABLE_LEFT, Inches(yb + 0.1), TABLE_W, Inches(4.95))
    tf = bx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(
        [
            "Réseaux sociaux et moteurs de recherche : campagnes géolocalisées, créatives renouvelées tous les mois.",
            "Terrain : équipes ou ambassadeurs pour inscrire les commerçants et expliquer l’intérêt du livre scolaire, de la livraison et des services.",
            "Événements thématiques (rentrée scolaire, salons, quartiers) pour activer plusieurs modules en une fois.",
            "Support client réactif : réduction du churn dès les premières semaines après téléchargement.",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸  " + line
        p.font.size = Pt(FS_BODY + 0.5)
        p.space_after = Pt(12)
    footer(slide)
    add_logo_br(slide)


def add_slide_marketing_3(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Calendrier marketing & modules",
        "Aligné sur le déploiement Douala / Yaoundé / Bafoussam puis international",
    )
    rows = [
        ("T2 2026", "Teasing bourse du livre + e-commerce, livraison, navigation", "Préparer la rentrée : pré-inscriptions parents/élèves, collecte de stock et partenaires écoles/libraires"),
        ("T3 2026", "Accent fort bourse du livre + orientation + campagnes terrain écoles", "Capteur principal rentrée de septembre ; forte acquisition organique et preuve d’impact social immédiat"),
        ("T4 2026 – T2 2027", "Transport, santé, supermarchés, vie pratique + renfort bourse du livre", "Consolidation Cameroun et démarrage Côte d’Ivoire avec hausse du panier moyen"),
        ("T3 – T4 2027", "Immobilier, assurance, optimisation multi-pays", "Segments à ticket élevé + amélioration de la marge et baisse du coût d’acquisition"),
    ]
    top = Inches(yb + 0.1)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.48 * nrows))
    table = shape.table
    for j, w in enumerate((Inches(2.2), Inches(3.8), Inches(2.35))):
        table.columns[j].width = w
    hdr = ("Période", "Modules mis en avant", "Objectif")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY)
            row_text = " ".join(row).lower()
            has_book_market = "bourse du livre" in row_text
            if has_book_market:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xFF, 0xF6, 0xDE)
                c.text_frame.paragraphs[0].font.bold = True
            elif i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    mk_note = slide.shapes.add_textbox(TABLE_LEFT, Inches(5.78), TABLE_W, Inches(0.55))
    mkp = mk_note.text_frame.paragraphs[0]
    mkp.text = "Priorité marketing 2026 : bourse du livre dès T2 puis campagne forte en T3 pour capter la rentrée de septembre."
    mkp.font.size = Pt(FS_NOTE)
    mkp.font.color.rgb = NAVY
    footer(slide)
    add_logo_br(slide)


def add_slide_competition(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(
        slide,
        "Concurrents identifiés & positionnement Yukpo",
        "Acteurs nommés : chacun couvre une partie du besoin, aucun ne relie tous les segments locaux",
    )
    rows = [
        ("Jumia (e-commerce)", "Catalogue et logistique e-commerce structurés", "Faible profondeur multi-services du quotidien local (santé, orientation, vie pratique, éducation)."),
        ("Glovo / Bolt Food (livraison)", "Exécution rapide de livraison en zones denses", "Couverture centrée restauration/livraison ; faible intégration super-app transversale pour ménages et TPE."),
        ("Yango / Heetch (mobilité)", "Trajets urbains et disponibilité transport", "Peu de passerelles natives vers commerce local, bourse du livre, immobilier et services ménagers."),
        ("Facebook Marketplace / WhatsApp", "Audience massive et adoption spontanée", "Transactions peu sécurisées, traçabilité faible, parcours paiement/livraison non unifié."),
        ("Yukpo (positionnement cible)", "Super-app locale : commerce + livraison + mobilité + éducation + services", "Exécution opérationnelle à tenir ville par ville ; avantage compétitif lié à l’intégration segmentaire."),
    ]
    top = Inches(yb + 0.1)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.39 * nrows))
    table = shape.table
    hdr = ("Type d’acteur", "Ce qu’il fait bien", "Limite pour l’usager")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY)
            c.text_frame.word_wrap = True
            if i == nrows - 1:
                c.text_frame.paragraphs[0].font.bold = True
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xE8, 0xEF, 0xFB)
            if i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    footer(slide)
    add_logo_br(slide)


def add_table_budget_detail(prs: Presentation) -> None:
    slide = slide_blank(prs)
    yb = title_block(slide, "Détail du besoin — 200 millions FCFA", "Ventilation par poste")
    rows = [
        ("Marketing & acquisition", f"{B['marketing_m']} M", f"{round(100 * B['marketing_m'] / B['total_m'])} %"),
        ("Personnel & rémunérations", f"{B['personnel_m']} M", f"{round(100 * B['personnel_m'] / B['total_m'])} %"),
        ("Technologie & infrastructure", f"{B['tech_m']} M", f"{round(100 * B['tech_m'] / B['total_m'])} %"),
        ("Juridique & conformité", f"{B['legal_m']} M", f"{round(100 * B['legal_m'] / B['total_m'])} %"),
        ("Frais opérationnels", f"{B['divers_m']} M", f"{round(100 * B['divers_m'] / B['total_m'])} %"),
        ("Réserve / imprévus", f"{B['reserve_m']} M", f"{round(100 * B['reserve_m'] / B['total_m'])} %"),
        ("TOTAL", f"{B['total_m']} M", "100 %"),
    ]
    nrows = len(rows) + 1
    top = Inches(yb + 0.08)
    shape = slide.shapes.add_table(nrows, 3, TABLE_LEFT, top, TABLE_W, Inches(0.48 * nrows))
    table = shape.table
    for j, w in enumerate((Inches(4.0), Inches(2.2), Inches(1.85))):
        table.columns[j].width = w
    hdr = ("Poste", "Montant", "Part")
    for j, h in enumerate(hdr):
        c = table.cell(0, j)
        c.text = ""
        c.text_frame.paragraphs[0].text = h
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_HDR + 0.5)
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, start=1):
        is_total = row[0].startswith("TOTAL")
        for j, val in enumerate(row):
            c = table.cell(i, j)
            c.text = ""
            c.text_frame.paragraphs[0].text = val
            c.text_frame.paragraphs[0].font.size = Pt(FS_TABLE_BODY + 0.5)
            if is_total:
                c.text_frame.paragraphs[0].font.bold = True
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF5)
            elif i % 2 == 0:
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)
    footer(slide)
    add_logo_br(slide)


def _boardroom_compact_text(text: str) -> str:
    s = " ".join((text or "").split())
    if not s:
        return s
    if s.startswith(("•", "▸")):
        s = s[1:].strip()
    # Garde prioritairement la 1re phrase, puis tronque si nécessaire.
    parts = re.split(r"(?<=[.!?])\s+", s)
    head = parts[0] if parts else s
    if len(head) > 135:
        head = head[:132].rsplit(" ", 1)[0] + "."
    if head and head[-1] not in ".!?":
        head += "."
    return head


def _slide_all_texts(slide) -> list[str]:
    out: list[str] = []
    for sh in slide.shapes:
        tf = getattr(sh, "text_frame", None)
        if tf is None:
            continue
        for p in tf.paragraphs:
            t = (p.text or "").strip()
            if t:
                out.append(t)
    return out


def _remove_slide(prs: Presentation, idx: int) -> None:
    slide_id = prs.slides._sldIdLst[idx]  # type: ignore[attr-defined]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[idx]  # type: ignore[attr-defined]


def create_boardroom_variant(src: Path, dst: Path) -> None:
    prs = Presentation(str(src))

    # Condensation boardroom : retire une partie des diapos détaillées/répétitives.
    remove_idxs: list[int] = []
    for i, slide in enumerate(prs.slides):
        all_text = " ".join(_slide_all_texts(slide)).lower()
        if "impact social" in all_text:
            remove_idxs.append(i)
            continue
        if "déploiement marketing — canaux" in all_text or "deploiement marketing" in all_text:
            remove_idxs.append(i)
            continue
    for idx in sorted(set(remove_idxs), reverse=True):
        _remove_slide(prs, idx)

    for slide in prs.slides:
        for shape in slide.shapes:
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            for p in tf.paragraphs:
                t = (p.text or "").strip()
                if not t:
                    continue
                if len(t) > 120 or t.startswith(("•", "▸")):
                    p.text = _boardroom_compact_text(t)
                # Style plus sobre "boardroom"
                if p.font.size is not None and p.font.size.pt <= 11:
                    p.font.size = Pt(max(10.5, p.font.size.pt))

    prs.save(dst)


def main() -> None:
    global _LOGO_PATH
    _LOGO_PATH = resolve_logo_path()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Couverture
    s = slide_blank(prs)
    add_accent_sidebar(s)
    title_block(
        s,
        "Yukpo",
        f"Besoin de financement : {B['total_m']} millions FCFA · déploiement mai 2026 – fin 2027",
    )
    add_two_kpi_wide(
        s,
        1.75,
        (f"{B['total_m']} M", "Enveloppe demandée", "Budget marketing & équipe renforcés"),
        ("Juillet 2026", "Mise en ligne effective", "18 mois de production suivis dans le plan"),
    )
    auth = s.shapes.add_textbox(Inches(0.85), Inches(4.85), Inches(8.4), Inches(1.2))
    t = auth.text_frame
    t.paragraphs[0].text = "Hernandez LELE — CEO, Yukpo Company"
    t.paragraphs[0].font.size = Pt(FS_BODY + 1)
    t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.color.rgb = NAVY
    t.paragraphs[0].alignment = PP_ALIGN.CENTER
    p2 = t.add_paragraph()
    p2.text = f"Entrée au capital possible jusqu’à {B['equity_cap_pct']} % (discussion encadrée)"
    p2.font.size = Pt(FS_NOTE)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    footer(s)
    add_logo_br(s)

    # Marché informel — une diapo
    add_slide_market_informal(prs)

    # Problème / opportunité (langage simple)
    s = slide_blank(prs)
    yt = title_block(s, "Pourquoi investir maintenant ?", "Une masse de besoins peu digitalisée")
    add_kpi_row(
        s,
        yt + 0.05,
        [
            ("85 % env.", "Peu visibles en ligne", "Commerces locaux"),
            ("100K – 1M", "Coût d’une boutique web classique", "Barrière pour les TPE"),
            ("2 à 5", "Applications en parallèle", "Fatigue pour les familles"),
        ],
    )
    footer(s)
    add_logo_br(s)

    # Réponse produit
    s = slide_blank(prs)
    yt = title_block(s, "La réponse Yukpo", "Une application unique pour plusieurs vies quotidiennes")
    add_kpi_row(
        s,
        yt + 0.05,
        [
            ("≈ 5 min", "Mise en ligne guidée", "Pour le marchand"),
            ("9 leviers", "E-commerce, navigation, santé", "Même compte utilisateur"),
            ("1 ville", "Densifier puis étendre", "Stratégie progressive"),
        ],
    )
    footer(s)
    add_logo_br(s)

    # Huit segments × 2 diapos (besoins/solutions + impact/valeur)
    for mod in MODULES_DEEP:
        add_slide_module_part1(prs, mod)
        add_slide_module_part2(prs, mod)

    add_slide_competition(prs)

    # Marketing x3
    add_slide_marketing_1(prs)
    add_slide_marketing_depenses_lignes(prs)
    add_slide_marketing_2(prs)
    add_slide_marketing_poches_consommation(prs)
    add_slide_marketing_3(prs)
    add_slide_deploiement_geographique(prs)

    # Budget
    s = slide_blank(prs)
    add_chart_column(
        s,
        "Répartition du besoin (millions FCFA)",
        [
            "Marketing",
            "Personnel",
            "Tech",
            "Juridique",
            "Frais",
            "Réserve",
        ],
        "Millions FCFA",
        [
            float(B["marketing_m"]),
            float(B["personnel_m"]),
            float(B["tech_m"]),
            float(B["legal_m"]),
            float(B["divers_m"]),
            float(B["reserve_m"]),
        ],
        "Les deux premiers postes concentrent 80 % de l’enveloppe",
        f"Total : {B['total_m']} millions FCFA — hypothèses arrondies pour la lecture",
    )
    footer(s)
    add_logo_br(s)

    add_table_budget_detail(prs)
    add_slide_charges_two_years(prs)
    add_slide_financing_overview(prs)
    add_slide_financing_options(prs)
    add_slide_acquisition_croissance(prs)
    add_slide_revenus_base_calcul(prs)
    add_slide_simulation_quarters_1(prs)
    add_slide_simulation_quarters_2(prs)
    add_slide_bilan_rentabilite(prs)

    # Équipe & profils (budget personnel : 42 M)
    s = slide_blank(prs)
    add_accent_sidebar(s)
    yb_eq = title_block(s, "Équipe & profils recherchés", "42 millions FCFA sur 200 — charges salariales croissantes en 2ᵉ année")
    bx = s.shapes.add_textbox(TABLE_LEFT, Inches(yb_eq + 0.08), TABLE_W, Inches(max_body_height_inches(yb_eq + 0.08)))
    t = bx.text_frame
    t.word_wrap = True
    t.paragraphs[0].text = f"Hernandez LELE — fondateur & CEO ({FOUNDER_EXP_Y} ans d’expérience, institutionnel & terrain)."
    t.paragraphs[0].font.size = Pt(FS_BODY + 1.5)
    t.paragraphs[0].font.bold = True
    t.paragraphs[0].font.color.rgb = NAVY
    profs = [
        "Responsable acquisition terrain & partenariats (marchands, écoles, radios).",
        "Chargé(e) marketing digital & contenus (réseaux, créatives, influence locale).",
        "Chef de projet opérations livraison / logistique (qualité de service, coursiers).",
        "Support client & onboarding (réduction du churn, formation utilisateurs).",
        "Commerciaux B2B à temps partagé ou mission (secteur santé, assurance, immobilier).",
        "Renfort technique à la montée en charge : supervision cloud, sécurité, astreinte (en lien avec la hausse du poste tech).",
    ]
    for i, line in enumerate(profs):
        p = t.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(FS_NOTE)
        p.space_after = Pt(5)
        p.space_before = Pt(10 if i == 0 else 0)
    footer(s)
    add_logo_br(s)

    s = slide_blank(prs)
    yb_z = title_block(s, "Périmètre géographique", "Cameroun d’abord — extension Afrique francophone")
    bx = s.shapes.add_textbox(TABLE_LEFT, Inches(yb_z + 0.4), TABLE_W, Inches(2.0))
    bx.text_frame.paragraphs[0].text = "Déploiement progressif : preuve locale avant duplication."
    bx.text_frame.paragraphs[0].font.size = Pt(19)
    bx.text_frame.paragraphs[0].font.bold = True
    bx.text_frame.paragraphs[0].font.color.rgb = NAVY
    footer(s)
    add_logo_br(s)

    s = slide_blank(prs)
    yb_t = title_block(s, "Technologie", "Prête à encaisser la montée en charge")
    bx = s.shapes.add_textbox(TABLE_LEFT, Inches(yb_t + 0.12), TABLE_W, Inches(max_body_height_inches(yb_t + 0.12)))
    for k, line in enumerate(
        [
            "Backend performant (Rust) — API et temps réel.",
            "Application mobile (React Native) — parcours unifiés pour tous les segments.",
            f"Budget global tech ({B['tech_m']} M sur {B['total_m']} M) : les décaissements augmentent avec les utilisateurs (cloud, monitoring, mises à jour).",
        ]
    ):
        p = bx.text_frame.paragraphs[0] if k == 0 else bx.text_frame.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(FS_BODY + 2)
        p.space_after = Pt(11)
    footer(s)
    add_logo_br(s)

    s = slide_blank(prs)
    add_accent_sidebar(s)
    yb_c = title_block(s, "Contact", None)
    bx = s.shapes.add_textbox(TABLE_LEFT, Inches(yb_c + 0.35), TABLE_W, Inches(3.0))
    for k, line in enumerate(
        ["Hernandez LELE", "lelehernandez2007@yahoo.fr", "+237 674 546 895", "yukpomnang.com"]
    ):
        p = bx.text_frame.paragraphs[0] if k == 0 else bx.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(17 if k == 0 else 15)
        p.font.bold = k == 0
        p.font.color.rgb = NAVY if k == 0 else RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(10)
    footer(s, "NDA possible — data room sur demande")
    add_logo_br(s)

    stamp(prs)
    root = Path(__file__).resolve().parent.parent
    base_out = root / "Yukpo_Investisseur_Presentation.pptx"
    out = root / "Yukpo_Investisseur_Presentation_Finale.pptx"
    prs.save(base_out)
    prs.save(out)
    boardroom_out = root / "Yukpo_Investisseur_Presentation_Boardroom.pptx"
    create_boardroom_variant(out, boardroom_out)
    print(f"Fichier créé : {base_out}")
    print(f"Fichier créé : {out}")
    print(f"Fichier créé : {boardroom_out}")


if __name__ == "__main__":
    main()
