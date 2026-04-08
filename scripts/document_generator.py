#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
YukpoIA Document Generator — v2
Reads a JSON outline from stdin, generates PPTX / DOCX / PDF, writes bytes to stdout.

Input JSON schema:
{
  "document_type": "pptx" | "docx" | "pdf",
  "title": "Titre du document",
  "subtitle": "Sous-titre optionnel",
  "author": "Auteur",
  "theme": "blue" | "green" | "purple" | "orange" | "dark" | "light",
  "logo_base64": "...",          // Image logo/bannière en base64 (optionnel)
  "logo_mime": "image/png",      // MIME du logo (optionnel, défaut image/png)

  // Pour PPTX
  "slides": [
    {
      "title": "Titre de la diapo",
      "layout": "title" | "content" | "two_column" | "kpi" | "table" | "image" | "blank",
      "content": "Texte principal",
      "bullets": ["Point 1", "Point 2"],
      "sub_bullets": [["Sub A", "Sub B"], []],   // Sous-points par bullet (optionnel)
      "left": "Colonne gauche (two_column)",
      "right": "Colonne droite (two_column)",
      "table": {"headers": ["Col1"], "rows": [["val1"]]},
      "kpis": [{"value": "95%", "label": "Satisfaction", "trend": "+5%"}],
      "notes": "Notes de présentation",
      "accent_color": "#FF5733",
      "image_base64": "...",     // Image pleine-diapo (layout=image)
      "image_caption": "Légende",
      "icon": "▸"               // Icone personnalisée pour les bullets
    }
  ],

  // Pour DOCX
  "sections": [
    {
      "heading": "Titre de section",
      "level": 1,
      "paragraphs": ["Paragraphe 1"],
      "bullets": ["Point 1"],
      "table": {"headers": ["Col1"], "rows": [["val1"]]},
      "image_base64": "...",     // Image dans la section (optionnel)
      "image_caption": "Légende"
    }
  ],

  // Pour PDF
  "pages": [
    {
      "title": "Titre de page",
      "content": "Corps de texte",
      "bullets": ["Point 1", "Point 2"],
      "table": {"headers": ["Col1"], "rows": [["val1"]]},
      "kpis": [{"value": "95%", "label": "Satisfaction"}],
      "image_base64": "...",
      "image_caption": "Légende",
      "layout": "content" | "cover" | "kpi" | "table"
    }
  ]
}
"""

from __future__ import annotations
import base64
import io
import json
import os
import sys
import tempfile
import traceback
from typing import Any

# ─── Thèmes ─────────────────────────────────────────────────────────────────

THEMES: dict[str, dict[str, tuple[int, int, int]]] = {
    "blue":   {"primary": (0, 84, 166),   "accent": (0, 176, 240),  "bg": (245, 248, 252), "text": (30, 30, 50),   "subtitle": (90, 130, 180),  "table_alt": (224, 235, 252)},
    "green":  {"primary": (27, 120, 53),  "accent": (46, 213, 115), "bg": (243, 252, 246), "text": (20, 50, 20),   "subtitle": (60, 140, 80),   "table_alt": (220, 248, 230)},
    "purple": {"primary": (102, 0, 153),  "accent": (180, 100, 255),"bg": (250, 245, 255), "text": (40, 10, 60),   "subtitle": (130, 70, 180),  "table_alt": (238, 224, 255)},
    "orange": {"primary": (200, 75, 0),   "accent": (255, 160, 0),  "bg": (255, 250, 240), "text": (60, 30, 0),    "subtitle": (180, 100, 30),  "table_alt": (255, 238, 212)},
    "dark":   {"primary": (20, 20, 45),   "accent": (0, 200, 255),  "bg": (22, 22, 42),    "text": (225, 225, 240),"subtitle": (155, 165, 200), "table_alt": (35, 35, 65)},
    "light":  {"primary": (45, 55, 130),  "accent": (255, 130, 0),  "bg": (255, 255, 255), "text": (30, 30, 50),   "subtitle": (100, 110, 160), "table_alt": (235, 236, 248)},
}
DEFAULT_THEME = "blue"


def get_theme(name: str) -> dict[str, tuple[int, int, int]]:
    return THEMES.get(name, THEMES[DEFAULT_THEME])


def hex_to_rgb(hex_str: str) -> tuple[int, int, int]:
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def decode_image_b64(b64_str: str) -> bytes | None:
    try:
        data = b64_str.strip()
        if "," in data:
            data = data.split(",", 1)[1]
        return base64.b64decode(data)
    except Exception:
        return None


def save_tmp_image(b64_str: str, suffix: str = ".png") -> str | None:
    """Decode base64 image and save to a temp file; return path."""
    img_bytes = decode_image_b64(b64_str)
    if not img_bytes:
        return None
    fd, path = tempfile.mkstemp(suffix=suffix)
    try:
        os.write(fd, img_bytes)
    finally:
        os.close(fd)
    return path


def img_suffix_from_mime(mime: str) -> str:
    return {"image/jpeg": ".jpg", "image/jpg": ".jpg",
            "image/gif": ".gif", "image/webp": ".webp"}.get(mime, ".png")


# ─── PPTX ────────────────────────────────────────────────────────────────────

def generate_pptx(data: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    T = get_theme(data.get("theme", DEFAULT_THEME))
    is_dark = data.get("theme") == "dark"
    primary  = RGBColor(*T["primary"])
    accent   = RGBColor(*T["accent"])
    bg_rgb   = RGBColor(*T["bg"])
    txt_rgb  = RGBColor(*T["text"])
    sub_rgb  = RGBColor(*T["subtitle"])
    tbl_alt  = RGBColor(*T["table_alt"])
    white    = RGBColor(255, 255, 255)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_text    = data.get("title", "Document")
    subtitle_text = data.get("subtitle", "")
    author        = data.get("author", "YukpoIA")
    logo_b64      = data.get("logo_base64", "")
    logo_mime     = data.get("logo_mime", "image/png")

    # Temp file for logo
    logo_path = None
    if logo_b64:
        logo_path = save_tmp_image(logo_b64, img_suffix_from_mime(logo_mime))

    def add_bg(slide, color: RGBColor):
        s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        el = s._element; slide.shapes._spTree.remove(el); slide.shapes._spTree.insert(2, el)

    def add_rect(slide, x, y, w, h, color: RGBColor, border: RGBColor | None = None):
        s = slide.shapes.add_shape(1, x, y, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        if border:
            s.line.color.rgb = border
        else:
            s.line.fill.background()
        return s

    def add_text(slide, x, y, w, h, text: str, size: float, bold=False,
                 color: RGBColor | None = None, align=PP_ALIGN.LEFT, italic=False) -> None:
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color or txt_rgb

    def add_footer(slide):
        # Bottom accent line
        add_rect(slide, 0, Inches(7.1), prs.slide_width, Inches(0.06), accent)
        add_text(slide, Inches(0.4), Inches(7.18), Inches(9.0), Inches(0.28),
                 f"{title_text}  ·  YukpoIA", 8.5, color=sub_rgb)

    def add_header(slide, slide_title: str, hdr_color: RGBColor | None = None):
        add_rect(slide, 0, 0, prs.slide_width, Inches(1.1), hdr_color or primary)
        add_text(slide, Inches(0.45), Inches(0.14), Inches(11.5), Inches(0.85),
                 slide_title, 22, bold=True, color=white)

    # ── COVER slide ────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_rgb)

    # Left accent bar + top stripe
    add_rect(slide, 0, 0, Inches(0.18), prs.slide_height, primary)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.07), accent)

    # Logo in top-right
    if logo_path:
        try:
            slide.shapes.add_picture(logo_path, Inches(10.8), Inches(0.3), width=Inches(2.0))
        except Exception:
            pass

    add_text(slide, Inches(0.8), Inches(2.0), Inches(11.3), Inches(2.0),
             title_text, 38, bold=True, color=primary)
    if subtitle_text:
        add_text(slide, Inches(0.8), Inches(4.15), Inches(11.0), Inches(1.0),
                 subtitle_text, 18, italic=True, color=sub_rgb)

    # Divider + author
    add_rect(slide, Inches(0.8), Inches(5.25), Inches(12.3), Inches(0.035), accent)
    add_text(slide, Inches(0.8), Inches(5.45), Inches(6.0), Inches(0.5),
             f"Préparé par : {author}", 11, italic=True, color=sub_rgb)

    # ── Content slides ─────────────────────────────────────────────────────────
    for sd in data.get("slides", []):
        layout = sd.get("layout", "content")
        slide  = prs.slides.add_slide(blank)
        add_bg(slide, bg_rgb)

        # Per-slide accent override
        sl_accent = accent
        if acc_hex := sd.get("accent_color"):
            try:
                sl_accent = RGBColor(*hex_to_rgb(acc_hex))
            except Exception:
                pass

        add_header(slide, sd.get("title", ""))
        add_footer(slide)

        content = sd.get("content", "")
        bullets = sd.get("bullets", [])
        sub_bullets = sd.get("sub_bullets", [])
        icon = sd.get("icon", "▸")

        # ── image layout ───────────────────────────────────────────────────────
        if layout == "image":
            img_b64 = sd.get("image_base64", "")
            if img_b64:
                img_path = save_tmp_image(img_b64)
                if img_path:
                    try:
                        slide.shapes.add_picture(img_path,
                            Inches(0.4), Inches(1.2), width=Inches(12.5),
                            height=Inches(5.8))
                        os.unlink(img_path)
                    except Exception:
                        pass
            if cap := sd.get("image_caption"):
                add_text(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.3),
                         cap, 9.5, italic=True, color=sub_rgb, align=PP_ALIGN.CENTER)

        # ── title layout ───────────────────────────────────────────────────────
        elif layout == "title":
            y = Inches(1.5)
            if content:
                add_text(slide, Inches(0.7), y, Inches(11.9), Inches(1.8),
                         content, 28, bold=True, color=primary, align=PP_ALIGN.CENTER)
                y = Inches(3.5)
            for bi, b in enumerate(bullets):
                add_text(slide, Inches(1.5), y + Inches(bi * 0.65), Inches(10.0),
                         Inches(0.6), f"  {icon}  {b}", 16, color=txt_rgb, align=PP_ALIGN.CENTER)

        # ── two_column layout ──────────────────────────────────────────────────
        elif layout == "two_column":
            half = prs.slide_width // 2
            left_text  = sd.get("left", "")
            right_text = sd.get("right", "")

            for i, col_text in enumerate([left_text, right_text]):
                cx = Inches(0.4) if i == 0 else half + Inches(0.2)
                cw = half - Inches(0.6)
                cy = Inches(1.3)
                for j, line in enumerate(col_text.split("\n") if col_text else []):
                    if line.strip():
                        sz = 15 if j == 0 else 13
                        bld = (j == 0)
                        col = primary if j == 0 else txt_rgb
                        add_text(slide, cx, cy + Inches(j * 0.48), cw, Inches(0.46),
                                 line, sz, bold=bld, color=col)

            # Vertical divider
            add_rect(slide, half - Inches(0.04), Inches(1.2), Inches(0.04), Inches(5.8), sl_accent)

        # ── kpi layout ─────────────────────────────────────────────────────────
        elif layout == "kpi":
            kpis = sd.get("kpis", [])
            n = max(len(kpis), 1)
            cw = Inches(12.5) / n
            for ki, kpi in enumerate(kpis):
                kx = Inches(0.4) + int(ki * cw)
                # Card
                card_color = RGBColor(255, 255, 255) if not is_dark else RGBColor(38, 38, 60)
                add_rect(slide, kx + Inches(0.12), Inches(1.45),
                         cw - Inches(0.24), Inches(4.8), card_color, sl_accent)
                # Top accent band on card
                add_rect(slide, kx + Inches(0.12), Inches(1.45),
                         cw - Inches(0.24), Inches(0.08), sl_accent)
                # Value
                add_text(slide, kx + Inches(0.2), Inches(2.0),
                         cw - Inches(0.4), Inches(1.6),
                         str(kpi.get("value", "")), 44, bold=True,
                         color=primary, align=PP_ALIGN.CENTER)
                # Label
                add_text(slide, kx + Inches(0.2), Inches(3.8),
                         cw - Inches(0.4), Inches(0.7),
                         str(kpi.get("label", "")), 13,
                         color=sub_rgb, align=PP_ALIGN.CENTER)
                # Trend
                if trend := kpi.get("trend"):
                    tc = RGBColor(46, 180, 100) if "+" in str(trend) else RGBColor(220, 60, 60)
                    add_text(slide, kx + Inches(0.2), Inches(4.6),
                             cw - Inches(0.4), Inches(0.5),
                             str(trend), 12, bold=True, color=tc, align=PP_ALIGN.CENTER)

        # ── table layout ───────────────────────────────────────────────────────
        elif layout == "table":
            td = sd.get("table", {})
            headers = td.get("headers", [])
            rows    = td.get("rows", [])
            cols = max(len(headers), max((len(r) for r in rows), default=1), 1)
            total_rows = len(rows) + (1 if headers else 0)
            if total_rows > 0 and cols > 0:
                row_h = min(0.55, (5.8 / total_rows))
                tbl = slide.shapes.add_table(
                    total_rows, cols,
                    Inches(0.4), Inches(1.25),
                    Inches(12.5), Inches(row_h * total_rows)
                ).table
                # Headers
                for ci, h in enumerate(headers[:cols]):
                    cell = tbl.cell(0, ci)
                    cell.text = h
                    cell.fill.solid(); cell.fill.fore_color.rgb = primary
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.bold = True; r.font.size = Pt(11.5)
                            r.font.color.rgb = white
                # Rows
                for ri, row in enumerate(rows):
                    bg = tbl_alt if ri % 2 == 0 else RGBColor(*T["bg"])
                    for ci, val in enumerate(row[:cols]):
                        c = tbl.cell(ri + (1 if headers else 0), ci)
                        c.text = str(val)
                        c.fill.solid(); c.fill.fore_color.rgb = bg
                        for p in c.text_frame.paragraphs:
                            for r in p.runs:
                                r.font.size = Pt(11); r.font.color.rgb = txt_rgb

        # ── default content layout ─────────────────────────────────────────────
        else:
            y = Inches(1.3)

            # Optional inline image (right side)
            img_b64 = sd.get("image_base64", "")
            text_w = Inches(12.2)
            if img_b64:
                img_path = save_tmp_image(img_b64)
                if img_path:
                    try:
                        slide.shapes.add_picture(img_path,
                            Inches(9.0), Inches(1.3), width=Inches(4.0))
                        os.unlink(img_path)
                        text_w = Inches(8.4)
                    except Exception:
                        pass

            if content:
                box = slide.shapes.add_textbox(Inches(0.6), y, text_w, Inches(1.3))
                tf = box.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run(); r.text = content
                r.font.size = Pt(14.5); r.font.color.rgb = txt_rgb
                y = Inches(2.7)

            if bullets:
                avail_h = Inches(7.5) - y - Inches(0.5)
                box = slide.shapes.add_textbox(Inches(0.6), y, text_w, avail_h)
                tf = box.text_frame; tf.word_wrap = True
                for bi, bullet in enumerate(bullets):
                    bp = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                    bp.space_before = Pt(3); bp.space_after = Pt(3)
                    br = bp.add_run()
                    br.text = f"  {icon}  {bullet}"
                    br.font.size = Pt(13.5); br.font.color.rgb = txt_rgb
                    # Sub-bullets
                    subs = sub_bullets[bi] if bi < len(sub_bullets) else []
                    for sub in subs:
                        sp2 = tf.add_paragraph()
                        sp2.space_before = Pt(1); sp2.space_after = Pt(1)
                        sp2.level = 1
                        sr = sp2.add_run()
                        sr.text = f"        ◦  {sub}"
                        sr.font.size = Pt(12); sr.font.color.rgb = sub_rgb

        # Notes
        if notes := sd.get("notes", ""):
            slide.notes_slide.notes_text_frame.text = notes

    if logo_path:
        try:
            os.unlink(logo_path)
        except Exception:
            pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── DOCX ────────────────────────────────────────────────────────────────────

def generate_docx(data: dict[str, Any]) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    T = get_theme(data.get("theme", DEFAULT_THEME))
    primary  = RGBColor(*T["primary"])
    accent   = RGBColor(*T["accent"])
    txt_rgb  = RGBColor(*T["text"])
    sub_rgb  = RGBColor(*T["subtitle"])

    logo_b64  = data.get("logo_base64", "")
    logo_mime = data.get("logo_mime", "image/png")

    doc = Document()
    for sec in doc.sections:
        sec.left_margin = sec.right_margin = Cm(2.5)
        sec.top_margin = sec.bottom_margin = Cm(2.5)

    def set_color(run, rgb: RGBColor):
        run.font.color.rgb = rgb

    def add_hr(doc_obj):
        p = doc_obj.add_paragraph()
        pPr = p._element.get_or_add_pPr()
        pBdr = OxmlElement("w:pBdr")
        bot = OxmlElement("w:bottom")
        bot.set(qn("w:val"), "single"); bot.set(qn("w:sz"), "6")
        bot.set(qn("w:space"), "1")
        bot.set(qn("w:color"), "%02X%02X%02X" % T["accent"])
        pBdr.append(bot); pPr.append(pBdr)
        return p

    def colored_heading(text: str, level: int):
        h = doc.add_heading(text, level=level)
        h.alignment = WD_ALIGN_PARAGRAPH.LEFT
        for r in h.runs:
            r.font.color.rgb = primary if level == 1 else accent
        return h

    # ── Cover ──────────────────────────────────────────────────────────────────
    # Logo
    if logo_b64:
        lpath = save_tmp_image(logo_b64, img_suffix_from_mime(logo_mime))
        if lpath:
            try:
                p_logo = doc.add_paragraph()
                p_logo.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                run_logo = p_logo.add_run()
                run_logo.add_picture(lpath, width=Inches(2.0))
                os.unlink(lpath)
            except Exception:
                pass

    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = tp.add_run(data.get("title", "Document"))
    tr.bold = True; tr.font.size = Pt(28); set_color(tr, primary)

    if sub := data.get("subtitle", ""):
        sp = doc.add_paragraph()
        sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sr = sp.add_run(sub)
        sr.italic = True; sr.font.size = Pt(16); set_color(sr, sub_rgb)

    ap = doc.add_paragraph()
    ap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    ar = ap.add_run(f"Par : {data.get('author', 'YukpoIA')}")
    ar.font.size = Pt(12); set_color(ar, sub_rgb)

    add_hr(doc)
    doc.add_paragraph()

    # ── Sections ───────────────────────────────────────────────────────────────
    for sec_data in data.get("sections", []):
        heading = sec_data.get("heading", "")
        level   = max(1, min(6, int(sec_data.get("level", 1))))
        if heading:
            colored_heading(heading, level)

        for para in sec_data.get("paragraphs", []):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            r = p.add_run(para)
            r.font.size = Pt(11.5); set_color(r, txt_rgb)

        for bullet in sec_data.get("bullets", []):
            bp = doc.add_paragraph(style="List Bullet")
            br = bp.add_run(bullet)
            br.font.size = Pt(11.5); set_color(br, txt_rgb)

        # Inline image
        img_b64 = sec_data.get("image_base64", "")
        if img_b64:
            ipath = save_tmp_image(img_b64)
            if ipath:
                try:
                    ip = doc.add_paragraph()
                    ip.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    ir = ip.add_run()
                    ir.add_picture(ipath, width=Inches(5.5))
                    os.unlink(ipath)
                    if cap := sec_data.get("image_caption"):
                        cp = doc.add_paragraph(cap)
                        cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        for cr in cp.runs:
                            cr.font.size = Pt(9.5); cr.italic = True; set_color(cr, sub_rgb)
                except Exception:
                    pass

        # Table
        td = sec_data.get("table")
        if td:
            headers = td.get("headers", [])
            rows    = td.get("rows", [])
            cols = max(len(headers), max((len(r) for r in rows), default=1), 1)
            total = len(rows) + (1 if headers else 0)
            if total > 0:
                tbl = doc.add_table(rows=total, cols=cols)
                tbl.style = "Table Grid"
                if headers:
                    for ci, h in enumerate(headers[:cols]):
                        cell = tbl.rows[0].cells[ci]
                        cell.text = ""
                        p2 = cell.paragraphs[0]
                        r2 = p2.add_run(h)
                        r2.bold = True; r2.font.size = Pt(11)
                        set_color(r2, RGBColor(255, 255, 255))
                        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
                        shd = OxmlElement("w:shd")
                        shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto")
                        shd.set(qn("w:fill"), "%02X%02X%02X" % T["primary"])
                        tcPr.append(shd)
                for ri, row in enumerate(rows):
                    ridx = ri + (1 if headers else 0)
                    for ci, val in enumerate(row[:cols]):
                        c = tbl.rows[ridx].cells[ci]
                        c.text = str(val)
                        for p3 in c.paragraphs:
                            for r3 in p3.runs:
                                r3.font.size = Pt(10.5); set_color(r3, txt_rgb)

        doc.add_paragraph()

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─── PDF ─────────────────────────────────────────────────────────────────────

def _pdf_safe(text: str) -> str:
    """Replace characters outside latin-1 with ASCII equivalents for fpdf2 Helvetica."""
    replacements = {
        "\u25b8": ">", "\u2192": "->", "\u2022": "*", "\u2013": "-", "\u2014": "--",
        "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
        "\u20ac": "EUR", "\u00d7": "x", "\u2714": "v", "\u2718": "x",
        "\u00b0": "deg", "\u00e9": "e", "\u00e8": "e", "\u00ea": "e", "\u00eb": "e",
        "\u00e0": "a", "\u00e2": "a", "\u00e4": "a", "\u00f4": "o", "\u00f6": "o",
        "\u00fc": "u", "\u00fb": "u", "\u00f9": "u", "\u00ee": "i", "\u00ef": "i",
        "\u00e7": "c", "\u00e6": "ae", "\u0153": "oe", "\u00c9": "E", "\u00c0": "A",
        "\u00c7": "C",
    }
    result = []
    for ch in text:
        result.append(replacements.get(ch, ch if ord(ch) < 256 else "?"))
    return "".join(result)


def _blend(c1: tuple, c2: tuple, alpha: float) -> tuple:
    """Blend two RGB tuples. alpha=1.0 → c1, alpha=0.0 → c2."""
    return tuple(int(c1[i] * alpha + c2[i] * (1 - alpha)) for i in range(3))


def generate_pdf(data: dict[str, Any]) -> bytes:
    try:
        from fpdf import FPDF, XPos, YPos
    except ImportError:
        from fpdf import FPDF
        XPos = None; YPos = None

    T = get_theme(data.get("theme", DEFAULT_THEME))
    P  = T["primary"]
    A  = T["accent"]
    BG = T["bg"]
    TX = T["text"]
    SB = T["subtitle"]
    TA = T["table_alt"]
    is_dark = data.get("theme") == "dark"

    logo_b64  = data.get("logo_base64", "")
    logo_mime = data.get("logo_mime", "image/png")
    title     = _pdf_safe(data.get("title", "Document"))
    subtitle  = _pdf_safe(data.get("subtitle", ""))
    author    = _pdf_safe(data.get("author", "YukpoIA"))

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_margins(20, 20, 20)

    W = 210  # A4 width mm
    CONTENT_W = W - 40  # usable width

    def set_fill(r, g, b):
        pdf.set_fill_color(r, g, b)

    def set_text(r, g, b):
        pdf.set_text_color(r, g, b)

    def draw_rect(x, y, w, h, fill_rgb, border_rgb=None):
        pdf.set_fill_color(*fill_rgb)
        if border_rgb:
            pdf.set_draw_color(*border_rgb)
        else:
            pdf.set_draw_color(*fill_rgb)
        pdf.rect(x, y, w, h, style="F")

    def header_band(page_title: str):
        draw_rect(0, 0, W, 24, P)
        # Left accent strip in header
        draw_rect(0, 0, 5, 24, A)
        pdf.set_xy(12, 6)
        pdf.set_font("Helvetica", "B", 17)
        set_text(255, 255, 255)
        pdf.cell(W - 24, 12, page_title, border=False, align="L")
        # Accent stripe bottom of header
        draw_rect(0, 24, W, 1.5, A)
        # Right decorative sidebar (full page height minus header)
        draw_rect(W - 4, 25.5, 4, 270, _blend(P, BG, 0.15))
        pdf.set_y(30)

    def footer_line():
        # Désactive le saut de page auto le temps du footer pour éviter les pages blanches
        pdf.set_auto_page_break(False)
        foot_y = 287  # 297 - 10mm du bas
        draw_rect(0, foot_y, W, 0.6, A)
        pdf.set_xy(10, foot_y + 1.5)
        pdf.set_font("Helvetica", "", 7.5)
        set_text(*SB)
        txt = f"{title}  |  YukpoIA  |  Page {pdf.page_no()}"
        pdf.cell(W - 20, 5, _pdf_safe(txt), align="C")
        pdf.set_auto_page_break(True, margin=20)

    # ── COVER page ─────────────────────────────────────────────────────────────
    pdf.add_page()
    # Full BG
    draw_rect(0, 0, W, 297, BG)
    # Large primary color block — top 55% of page
    draw_rect(0, 0, W, 160, P)
    # Accent diagonal strip overlay (simulated with thin bands)
    for i in range(8):
        alpha_val = 0.04 + i * 0.008
        draw_rect(W - 60 + i * 7, 0, 6, 160, _blend(A, P, alpha_val))
    # Bottom accent stripe at junction
    draw_rect(0, 158, W, 3, A)
    # Left sidebar strip (full height)
    draw_rect(0, 0, 5, 297, A)

    # YukpoIA watermark top-left corner
    pdf.set_xy(10, 8)
    pdf.set_font("Helvetica", "B", 10)
    set_text(255, 255, 255)
    pdf.cell(60, 7, "YukpoIA", border=False)

    # Logo top-right inside the primary block
    if logo_b64:
        lpath = save_tmp_image(logo_b64, img_suffix_from_mime(logo_mime))
        if lpath:
            try:
                pdf.image(lpath, x=W - 48, y=6, w=38)
                os.unlink(lpath)
            except Exception:
                pass

    # Title (inside primary block, centered vertically at ~60mm)
    pdf.set_xy(12, 60)
    pdf.set_font("Helvetica", "B", 34)
    set_text(255, 255, 255)
    pdf.multi_cell(W - 24, 13, title, align="L")

    # Subtitle (inside primary block)
    if subtitle:
        pdf.set_x(12)
        pdf.set_font("Helvetica", "I", 15)
        set_text(220, 235, 255)
        pdf.multi_cell(W - 24, 8, subtitle, align="L")

    # White content zone below the color block
    # Author + metadata card
    card_y = 172
    draw_rect(12, card_y, W - 24, 28, (255, 255, 255))
    draw_rect(12, card_y, 3, 28, A)
    pdf.set_xy(20, card_y + 5)
    pdf.set_font("Helvetica", "B", 11)
    set_text(*P)
    pdf.cell(W - 32, 7, _pdf_safe(f"Prepare par : {author}"), border=False, align="L")
    pdf.set_xy(20, card_y + 14)
    pdf.set_font("Helvetica", "", 10)
    set_text(*SB)
    pdf.cell(W - 32, 7, _pdf_safe("Genere avec YukpoIA — Intelligence Artificielle"), border=False, align="L")

    # Bottom large accent block
    draw_rect(0, 260, W, 37, _blend(P, BG, 0.08))
    draw_rect(0, 260, W, 1, A)
    pdf.set_xy(12, 270)
    pdf.set_font("Helvetica", "I", 9)
    set_text(*SB)
    pdf.cell(W - 24, 6, _pdf_safe("Document confidentiel — Usage interne"), border=False, align="C")

    # ── Content pages ──────────────────────────────────────────────────────────
    # Accept "pages" (PDF-native) or fall back to slides/sections
    pages_data = data.get("pages") or data.get("slides") or data.get("sections") or []

    for pg in pages_data:
        pdf.add_page()
        draw_rect(0, 0, W, 297, BG)

        page_title = _pdf_safe(pg.get("title", pg.get("heading", "")))
        layout = pg.get("layout", "content")

        header_band(page_title)
        cx = 20

        content    = _pdf_safe(pg.get("content", "") or (pg.get("paragraphs", [None])[0] or ""))
        bullets    = [_pdf_safe(b) for b in pg.get("bullets", [])]
        paragraphs = [_pdf_safe(p) for p in pg.get("paragraphs", [])]
        kpis       = pg.get("kpis", [])
        td       = pg.get("table", {})

        # Inline image (right column if text exists)
        img_b64 = pg.get("image_base64", "")
        img_path_tmp = None
        if img_b64:
            img_path_tmp = save_tmp_image(img_b64)

        if layout == "kpi" and kpis:
            n = len(kpis)
            cw = CONTENT_W / n
            card_y = pdf.get_y() + 8
            card_h = 52
            for ki, kpi in enumerate(kpis):
                kx = 20 + ki * cw
                card_bg = (255, 255, 255) if not is_dark else (40, 40, 65)
                # Shadow / depth effect
                draw_rect(kx + 3.5, card_y + 2, cw - 6, card_h, _blend(P, BG, 0.12))
                # Card background
                draw_rect(kx + 2, card_y, cw - 6, card_h, card_bg)
                # Top accent bar
                draw_rect(kx + 2, card_y, cw - 6, 3, A)
                # Bottom accent bar
                draw_rect(kx + 2, card_y + card_h - 2, cw - 6, 2, _blend(P, BG, 0.3))
                # Value
                pdf.set_xy(kx + 2, card_y + 8)
                pdf.set_font("Helvetica", "B", 28)
                set_text(*P)
                pdf.cell(cw - 6, 16, _pdf_safe(str(kpi.get("value", ""))), align="C")
                # Label
                pdf.set_xy(kx + 2, card_y + 26)
                pdf.set_font("Helvetica", "B", 9)
                set_text(*TX)
                pdf.cell(cw - 6, 8, _pdf_safe(str(kpi.get("label", ""))).upper(), align="C")
                # Trend
                if tr := kpi.get("trend"):
                    pdf.set_xy(kx + 2, card_y + 36)
                    is_pos = "+" in str(tr) or "hausse" in str(tr).lower()
                    tc = (30, 160, 80) if is_pos else (200, 50, 50)
                    set_text(*tc)
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.cell(cw - 6, 8, _pdf_safe(str(tr)), align="C")
                # Description
                if desc := kpi.get("description", kpi.get("desc", "")):
                    pdf.set_xy(kx + 4, card_y + 44)
                    pdf.set_font("Helvetica", "I", 7.5)
                    set_text(*SB)
                    pdf.multi_cell(cw - 10, 4, _pdf_safe(str(desc))[:50], align="C")
            pdf.set_y(card_y + card_h + 10)

        elif layout == "table" and td:
            _render_pdf_table(pdf, td, cx, CONTENT_W, P, TA, TX, SB, A)

        else:
            # Text content
            text_w = CONTENT_W
            if img_path_tmp:
                text_w = CONTENT_W * 0.57
                try:
                    img_x = 20 + text_w + 5
                    pdf.image(img_path_tmp, x=img_x, y=pdf.get_y(), w=CONTENT_W * 0.4)
                    if cap := pg.get("image_caption"):
                        pdf.set_xy(img_x, pdf.get_y() + CONTENT_W * 0.4 * 0.56 + 1)
                        pdf.set_font("Helvetica", "I", 8)
                        set_text(*SB)
                        pdf.cell(CONTENT_W * 0.4, 5, cap, align="C")
                    os.unlink(img_path_tmp)
                    img_path_tmp = None
                except Exception:
                    text_w = CONTENT_W

            if content:
                # Content card background
                content_y = pdf.get_y()
                card_bg = _blend(P, (255, 255, 255), 0.06) if not is_dark else (38, 38, 60)
                draw_rect(cx - 2, content_y - 2, text_w + 4, 4, card_bg)  # placeholder sized below
                pdf.set_xy(cx, content_y)
                pdf.set_font("Helvetica", "", 11.5)
                set_text(*TX)
                pdf.multi_cell(text_w, 6.5, content, align="J")
                pdf.ln(4)

            for para in paragraphs:
                if para and para != content:
                    pdf.set_x(cx)
                    pdf.set_font("Helvetica", "", 11)
                    set_text(*TX)
                    pdf.multi_cell(text_w, 6.5, para, align="J")
                    pdf.ln(3)

            if bullets:
                # Bullets section with left accent bar
                bull_start_y = pdf.get_y()
                draw_rect(cx - 2, bull_start_y, 2.5, len(bullets) * 8.5, A)
                for bullet in bullets:
                    by = pdf.get_y()
                    # Colored square bullet indicator
                    draw_rect(cx + 4, by + 2.5, 2.5, 2.5, P)
                    pdf.set_xy(cx + 10, by)
                    pdf.set_font("Helvetica", "", 11)
                    set_text(*TX)
                    bullet_safe = _pdf_safe(bullet)
                    pdf.multi_cell(text_w - 12, 7, bullet_safe, align="L")
                    pdf.ln(0.5)
                pdf.ln(2)

            # Table after text
            if td and layout != "table":
                pdf.ln(4)
                _render_pdf_table(pdf, td, cx, CONTENT_W if text_w == CONTENT_W else text_w, P, TA, TX, SB, A)

        if img_path_tmp:
            try:
                os.unlink(img_path_tmp)
            except Exception:
                pass

        # ── Decorative bottom zone if page is sparse (current y < 220mm) ──
        cur_y = pdf.get_y()
        if cur_y < 215:
            # Bottom accent band
            band_y = 268
            draw_rect(0, band_y, W, 8, _blend(P, BG, 0.08))
            draw_rect(0, band_y, W, 0.8, A)
            # Document branding in band
            pdf.set_xy(10, band_y + 1.5)
            pdf.set_font("Helvetica", "I", 8)
            set_text(*SB)
            pdf.cell(W - 20, 5, _pdf_safe(f"Document genere par YukpoIA  |  {title}"), align="C")
            # Corner accent triangle (bottom-left)
            draw_rect(0, 260, 20, 8, _blend(A, BG, 0.12))

        footer_line()

    buf = io.BytesIO()
    buf.write(pdf.output())
    return buf.getvalue()


def _render_pdf_table(pdf, td, cx, width, P, TA, TX, SB, A):
    headers = td.get("headers", [])
    rows    = td.get("rows", [])
    cols = max(len(headers), max((len(r) for r in rows), default=1), 1)
    cw = width / cols
    row_h = 8

    if headers:
        pdf.set_fill_color(*P)
        pdf.set_text_color(255, 255, 255)
        pdf.set_font("Helvetica", "B", 10)
        pdf.set_x(cx)  # position initiale une seule fois avant la boucle
        for h in headers[:cols]:
            text = _pdf_safe(str(h))[:25]
            pdf.cell(cw, row_h, text, border=1, fill=True, align="C")  # auto-avance x
        cx = 20
        pdf.ln()

    for ri, row in enumerate(rows):
        fill_rgb = TA if ri % 2 == 0 else (255, 255, 255)
        pdf.set_fill_color(*fill_rgb)
        pdf.set_text_color(*TX)
        pdf.set_font("Helvetica", "", 10)
        cx_cur = cx
        for ci, val in enumerate(row[:cols]):
            pdf.set_x(cx_cur)
            pdf.cell(cw, row_h, _pdf_safe(str(val))[:30], border=1, fill=True, align="C")
            cx_cur += cw
        pdf.ln()

    pdf.ln(2)


# ─── Main ────────────────────────────────────────────────────────────────────

# ─── Génération Excel (XLSX) ──────────────────────────────────────────────────

def generate_xlsx(data: dict) -> bytes:
    """Génère un fichier Excel d'analyse à partir de l'outline fourni par l'IA."""
    from openpyxl import Workbook
    from openpyxl.styles import (Font, PatternFill, Alignment, Border, Side,
                                  GradientFill)
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, Reference
    import io

    wb = Workbook()

    title = data.get("title", "Analyse")
    theme_color = data.get("theme_color", "1F4E79")  # bleu foncé par défaut
    accent_color = data.get("accent_color", "2E75B6")

    # ── Styles de base ────────────────────────────────────────────────────────
    header_font  = Font(bold=True, color="FFFFFF", size=11)
    header_fill  = PatternFill("solid", fgColor=theme_color)
    title_font   = Font(bold=True, color=theme_color, size=14)
    sub_font     = Font(bold=True, color=accent_color, size=11)
    section_fill = PatternFill("solid", fgColor="D6E4F0")
    alt_fill     = PatternFill("solid", fgColor="EBF3FB")
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left_align   = Alignment(horizontal="left",  vertical="center", wrap_text=True)
    thin_border  = Border(
        left=Side(style="thin", color="B0C4D8"),
        right=Side(style="thin", color="B0C4D8"),
        top=Side(style="thin", color="B0C4D8"),
        bottom=Side(style="thin", color="B0C4D8"),
    )

    def style_header_row(ws, row_idx, num_cols):
        for col in range(1, num_cols + 1):
            cell = ws.cell(row=row_idx, column=col)
            cell.font  = header_font
            cell.fill  = header_fill
            cell.alignment = center_align
            cell.border = thin_border

    def style_data_row(ws, row_idx, num_cols, alt=False):
        fill = alt_fill if alt else PatternFill()
        for col in range(1, num_cols + 1):
            cell = ws.cell(row=row_idx, column=col)
            cell.fill = fill
            cell.alignment = left_align
            cell.border = thin_border

    sheets = data.get("sheets", [])

    # Si l'IA fournit directement des "tables" (format simplifié), les convertir
    if not sheets and data.get("tables"):
        sheets = [{"name": t.get("title", f"Feuille {i+1}"), "table": t}
                  for i, t in enumerate(data["tables"])]

    # Feuille par défaut si l'IA n'a rien fourni
    if not sheets:
        sheets = [{"name": "Analyse", "table": data}]

    first = True
    for sheet_def in sheets:
        sheet_name = str(sheet_def.get("name", "Feuille"))[:31]
        if first:
            ws = wb.active
            ws.title = sheet_name
            first = False
        else:
            ws = wb.create_sheet(title=sheet_name)

        row = 1

        # Titre de la feuille
        ws.merge_cells(f"A{row}:H{row}")
        ws[f"A{row}"] = sheet_def.get("sheet_title", title)
        ws[f"A{row}"].font = title_font
        ws[f"A{row}"].alignment = center_align
        ws.row_dimensions[row].height = 28
        row += 1

        # Sous-titre / description
        desc = sheet_def.get("description", "")
        if desc:
            ws.merge_cells(f"A{row}:H{row}")
            ws[f"A{row}"] = desc
            ws[f"A{row}"].font = Font(italic=True, color="666666", size=10)
            ws[f"A{row}"].alignment = left_align
            row += 1

        row += 1  # ligne vide

        # KPIs (résumé rapide en haut)
        kpis = sheet_def.get("kpis", [])
        if kpis:
            ws.merge_cells(f"A{row}:H{row}")
            ws[f"A{row}"] = "📊 Indicateurs clés"
            ws[f"A{row}"].font = sub_font
            ws[f"A{row}"].fill = section_fill
            ws[f"A{row}"].alignment = left_align
            row += 1

            kpi_header = ["Indicateur", "Valeur", "Tendance", "Commentaire"]
            for c, h in enumerate(kpi_header, 1):
                ws.cell(row=row, column=c, value=h)
            style_header_row(ws, row, len(kpi_header))
            row += 1

            for ki, kpi in enumerate(kpis):
                ws.cell(row=row, column=1, value=kpi.get("label", ""))
                ws.cell(row=row, column=2, value=kpi.get("value", ""))
                trend = kpi.get("trend", "")
                cell_trend = ws.cell(row=row, column=3, value=trend)
                if isinstance(trend, str):
                    if "+" in trend:
                        cell_trend.font = Font(color="1D7A4B", bold=True)
                    elif "-" in trend:
                        cell_trend.font = Font(color="C00000", bold=True)
                ws.cell(row=row, column=4, value=kpi.get("comment", ""))
                style_data_row(ws, row, len(kpi_header), alt=ki % 2 == 1)
                row += 1
            row += 1

        # Table principale
        table_def = sheet_def.get("table", sheet_def)
        headers = table_def.get("headers", [])
        rows_data = table_def.get("rows", [])

        if headers:
            ws.merge_cells(f"A{row}:H{row}")
            ws[f"A{row}"] = sheet_def.get("table_title", "📋 Données détaillées")
            ws[f"A{row}"].font = sub_font
            ws[f"A{row}"].fill = section_fill
            ws[f"A{row}"].alignment = left_align
            row += 1

            for c, h in enumerate(headers, 1):
                ws.cell(row=row, column=c, value=str(h))
            style_header_row(ws, row, len(headers))
            row += 1

            for ri, data_row in enumerate(rows_data):
                for c, val in enumerate(data_row, 1):
                    cell = ws.cell(row=row, column=c, value=val)
                    # Détecter les nombres pour formatage
                    if isinstance(val, str):
                        try:
                            num = float(val.replace(" ", "").replace(",", ".").replace("%", ""))
                            cell.value = num
                            if "%" in val:
                                cell.number_format = "0.0%"
                        except ValueError:
                            pass
                style_data_row(ws, row, len(headers), alt=ri % 2 == 1)
                row += 1

            # Ligne totaux si définie
            totals = table_def.get("totals")
            if totals:
                for c, val in enumerate(totals, 1):
                    cell = ws.cell(row=row, column=c, value=val)
                    cell.font = Font(bold=True, color=theme_color)
                    cell.fill = PatternFill("solid", fgColor="D6E4F0")
                    cell.border = thin_border
                row += 1
            row += 1

        # Sections texte / analyse narrative
        for section in sheet_def.get("sections", []):
            ws.merge_cells(f"A{row}:H{row}")
            ws[f"A{row}"] = section.get("title", "")
            ws[f"A{row}"].font = sub_font
            ws[f"A{row}"].fill = section_fill
            ws[f"A{row}"].alignment = left_align
            row += 1

            for line in section.get("lines", []):
                ws.merge_cells(f"A{row}:H{row}")
                ws[f"A{row}"] = f"  • {line}"
                ws[f"A{row}"].alignment = left_align
                ws[f"A{row}"].font = Font(size=10)
                row += 1
            row += 1

        # ── Graphiques (charts) ─────────────────────────────────────────────
        # Génère les graphiques définis dans `charts` + auto-chart si data table présente
        _add_charts_to_sheet(ws, sheet_def, row)

        # Ajuster largeur colonnes
        for col_cells in ws.columns:
            max_len = 0
            col_letter = get_column_letter(col_cells[0].column)
            for cell in col_cells:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            ws.column_dimensions[col_letter].width = min(max(max_len + 2, 12), 45)

        ws.freeze_panes = "A4"

    # Onglet Métadonnées
    ws_meta = wb.create_sheet("ℹ️ À propos")
    ws_meta["A1"] = "Fichier généré par YukpoIA"
    ws_meta["A1"].font = Font(bold=True, size=12, color=theme_color)
    ws_meta["A2"] = f"Titre : {title}"
    ws_meta["A3"] = f"Date de génération : {data.get('generated_at', '')}"
    ws_meta["A4"] = "Source : Yukpo — Plateforme IA Africaine"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _add_charts_to_sheet(ws, sheet_def: dict, data_end_row: int):
    """
    Ajoute des graphiques à une feuille Excel.
    Supporte les specs explicites (`charts` dans sheet_def) et l'auto-detection.
    """
    try:
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        from openpyxl.utils import get_column_letter
    except ImportError:
        return

    table_def = sheet_def.get("table", sheet_def)
    headers = table_def.get("headers", [])
    rows_data = table_def.get("rows", [])

    # Position d'ancrage des graphiques (à droite des données ou en dessous)
    chart_anchor_col = len(headers) + 2
    chart_anchor_letter = get_column_letter(max(chart_anchor_col, 1))
    chart_row = 3  # commence sur la 3e ligne pour laisser le titre

    # ── 1. Graphiques explicitement définis par l'IA ──────────────────────
    explicit_charts = sheet_def.get("charts", [])
    for chart_def in explicit_charts:
        chart_type = chart_def.get("type", "bar").lower()
        chart_title = chart_def.get("title", "Graphique")
        data_col = chart_def.get("data_col", 2)   # colonne des valeurs (1-indexed)
        cat_col  = chart_def.get("categories_col", 1)
        from_row = chart_def.get("from_row", 2)
        to_row   = chart_def.get("to_row", max(2, data_end_row - 1))
        anchor   = chart_def.get("anchor", f"{chart_anchor_letter}{chart_row}")

        try:
            if chart_type == "pie":
                chart = PieChart()
            elif chart_type == "line":
                chart = LineChart()
            else:
                chart = BarChart()

            chart.title = chart_title
            chart.style = 10
            chart.width  = 18
            chart.height = 10

            data_ref = Reference(ws, min_col=data_col, min_row=from_row - 1, max_row=to_row)
            cats_ref = Reference(ws, min_col=cat_col,  min_row=from_row,     max_row=to_row)
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(cats_ref)
            ws.add_chart(chart, anchor)
            chart_row += 20
        except Exception as e:
            pass  # ne pas faire planter la génération pour un graphique raté

    # ── 2. Auto-chart : si table avec ≥ 3 lignes et ≥ 1 colonne numérique ──
    auto_chart = sheet_def.get("auto_chart", len(explicit_charts) == 0)
    if not auto_chart or not headers or len(rows_data) < 3:
        return

    # Détecter la première colonne numérique parmi les colonnes 2..N
    numeric_col_idx = None
    for ci in range(1, len(headers)):
        num_count = 0
        for row_vals in rows_data[:10]:
            if ci < len(row_vals):
                try:
                    float(str(row_vals[ci]).replace(" ", "").replace(",", ".").replace("%", ""))
                    num_count += 1
                except ValueError:
                    pass
        if num_count >= min(3, len(rows_data)):
            numeric_col_idx = ci
            break

    if numeric_col_idx is None:
        return

    # Trouver la plage de données dans la feuille
    # Les headers sont écrits en ligne 3 (après titre + desc), rows à partir de 4
    # On cherche la vraie ligne des headers dans la feuille
    header_row_in_ws = None
    for r_idx in range(1, min(20, ws.max_row + 1)):
        cell_val = ws.cell(row=r_idx, column=1).value
        if cell_val == headers[0]:
            header_row_in_ws = r_idx
            break

    if header_row_in_ws is None:
        return

    data_start = header_row_in_ws + 1
    data_end = data_start + len(rows_data) - 1

    chart = BarChart()
    chart.type = "col"
    chart.style = 10
    chart.title = sheet_def.get("chart_title", f"📊 {headers[numeric_col_idx]}")
    chart.y_axis.title = headers[numeric_col_idx]
    chart.x_axis.title = headers[0]
    chart.width  = 20
    chart.height = 12

    data_ref = Reference(ws,
                         min_col=numeric_col_idx + 1,
                         min_row=header_row_in_ws,
                         max_row=data_end)
    cats_ref = Reference(ws,
                         min_col=1,
                         min_row=data_start,
                         max_row=data_end)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)

    ws.add_chart(chart, f"{chart_anchor_letter}{chart_row}")


# ─── Parsing Excel entrant (pour analyse par l'IA) ────────────────────────────

def parse_excel_for_ai(data: dict) -> bytes:
    """
    Mode 'parse' : lit un fichier Excel base64, extrait les données structurées
    et retourne un JSON texte (UTF-8) décrivant chaque feuille.
    Ce JSON est ensuite injecté dans le contexte de l'IA pour analyse.
    """
    import base64
    from openpyxl import load_workbook
    import io

    excel_b64 = data.get("excel_base64", "")
    try:
        excel_bytes = base64.b64decode(excel_b64)
    except Exception as e:
        result = json.dumps({"error": f"Impossible de décoder le base64 : {e}"})
        return result.encode("utf-8")

    try:
        wb = load_workbook(io.BytesIO(excel_bytes), read_only=True, data_only=True)
    except Exception:
        # Essai avec xlrd pour les anciens .xls
        try:
            import xlrd
            book = xlrd.open_workbook(file_contents=excel_bytes)
            sheets_data = []
            for sh in book.sheets():
                headers = [str(sh.cell_value(0, c)) for c in range(sh.ncols)] if sh.nrows > 0 else []
                rows = []
                for r in range(1, min(sh.nrows, 201)):
                    rows.append([str(sh.cell_value(r, c)) for c in range(sh.ncols)])
                summary = _compute_summary(headers, rows)
                sheets_data.append({
                    "name": sh.name,
                    "num_rows": sh.nrows - 1,
                    "num_cols": sh.ncols,
                    "headers": headers,
                    "rows": rows[:200],
                    "summary": summary,
                    "chart_suggestions": _detect_chart_suggestions(headers, rows, summary),
                })
            result = json.dumps({"sheets": sheets_data, "num_sheets": len(sheets_data),
                                 "instructions_for_ia": "Analyse ces données Excel en profondeur."}, ensure_ascii=False)
            return result.encode("utf-8")
        except Exception as e2:
            result = json.dumps({"error": f"Impossible de lire le fichier Excel : {e2}"})
            return result.encode("utf-8")

    sheets_data = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_rows = list(ws.iter_rows(values_only=True))
        if not all_rows:
            continue

        # Première ligne = en-têtes
        headers = [str(h) if h is not None else f"Colonne_{i+1}"
                   for i, h in enumerate(all_rows[0])]

        data_rows = []
        for r in all_rows[1:201]:  # max 200 lignes
            data_rows.append([str(v) if v is not None else "" for v in r])

        summary = _compute_summary(headers, data_rows)
        chart_suggestions = _detect_chart_suggestions(headers, data_rows, summary)

        sheets_data.append({
            "name": sheet_name,
            "num_rows": max(0, len(all_rows) - 1),
            "num_cols": len(headers),
            "headers": headers,
            "rows": data_rows,
            "summary": summary,
            "chart_suggestions": chart_suggestions,
        })

    wb.close()

    # Générer un prompt d'analyse préformaté pour aider l'IA
    analysis_hints = []
    for s in sheets_data:
        for col, stats in s.get("summary", {}).items():
            td = stats.get("trend_direction", "stable")
            gp = stats.get("growth_pct")
            if gp is not None:
                g_str = f"+{gp}%" if gp > 0 else f"{gp}%"
                analysis_hints.append(f"• {s['name']}/{col} : {td} ({g_str}, moy={stats['mean']}, somme={stats['sum']})")
            if stats.get("anomalies"):
                analysis_hints.append(f"  ⚠️ Anomalies détectées : {stats['anomalies']}")

    result = json.dumps({
        "sheets": sheets_data,
        "num_sheets": len(sheets_data),
        "analysis_hints": analysis_hints,
        "instructions_for_ia": (
            "Analyse ces données Excel en profondeur. "
            "Identifie les tendances, anomalies et insights business clés. "
            "Si l'utilisateur demande un fichier Excel, génère un `document_generation` "
            "avec document_type=xlsx incluant des KPIs, des tableaux de données et des charts. "
            "Les `chart_suggestions` dans chaque feuille indiquent les graphiques pertinents à générer."
        )
    }, ensure_ascii=False)
    return result.encode("utf-8")


def _compute_summary(headers: list, rows: list) -> dict:
    """Calcule des statistiques avancées par colonne numérique avec tendances et anomalies."""
    summary = {}
    for i, header in enumerate(headers):
        vals = []
        for row in rows:
            if i < len(row):
                try:
                    v = float(str(row[i]).replace(",", ".").replace(" ", "").replace("%", ""))
                    vals.append(v)
                except ValueError:
                    pass
        if len(vals) >= 2:
            mean_val = sum(vals) / len(vals)
            # Tendance : % de croissance entre première et dernière valeur
            growth = None
            if vals[0] != 0:
                growth = round((vals[-1] - vals[0]) / abs(vals[0]) * 100, 1)
            # Anomalies : valeurs à plus de 2 écarts-types de la moyenne
            if len(vals) >= 4:
                variance = sum((v - mean_val) ** 2 for v in vals) / len(vals)
                std = variance ** 0.5
                anomalies = [round(v, 2) for v in vals if std > 0 and abs(v - mean_val) > 2 * std]
            else:
                anomalies = []
            # Direction de tendance pour l'IA
            trend_direction = "stable"
            if growth is not None:
                if growth > 5:
                    trend_direction = "hausse"
                elif growth < -5:
                    trend_direction = "baisse"
            summary[header] = {
                "min": round(min(vals), 2),
                "max": round(max(vals), 2),
                "mean": round(mean_val, 2),
                "sum": round(sum(vals), 2),
                "count": len(vals),
                "growth_pct": growth,
                "trend_direction": trend_direction,
                "anomalies": anomalies[:5],  # max 5
            }
    return summary


def _detect_chart_suggestions(headers: list, rows: list, summary: dict) -> list:
    """Suggère des types de graphiques pertinents selon la structure des données."""
    suggestions = []
    if not headers or not rows:
        return suggestions

    # Chercher une colonne catégorielle (étiquettes) et des colonnes numériques
    cat_cols = []
    num_cols = []
    for i, h in enumerate(headers):
        if h in summary:
            num_cols.append({"col_idx": i + 1, "header": h, "info": summary[h]})
        else:
            cat_cols.append({"col_idx": i + 1, "header": h})

    if not num_cols:
        return suggestions

    cat = cat_cols[0] if cat_cols else None

    # Suggestion graphique à barres (toujours pertinent avec catégories)
    if cat and len(rows) >= 3:
        best_num = max(num_cols, key=lambda c: abs(c["info"].get("sum", 0)))
        suggestions.append({
            "type": "bar",
            "title": f"{best_num['header']} par {cat['header']}",
            "categories_col": cat["col_idx"],
            "data_col": best_num["col_idx"],
            "reason": "comparaison catégorielle"
        })

    # Suggestion courbe si tendance temporelle détectée (label contient mois/année/jour)
    temporal_keywords = ["mois", "année", "annee", "jan", "fév", "mar", "avr", "mai", "juin",
                         "juil", "août", "aout", "sep", "oct", "nov", "déc", "dec",
                         "2020", "2021", "2022", "2023", "2024", "2025", "2026",
                         "q1", "q2", "q3", "q4", "trim", "sem"]
    is_temporal = cat and any(
        any(kw in str(r[cat["col_idx"] - 1]).lower() for kw in temporal_keywords)
        for r in rows[:5] if len(r) >= cat["col_idx"]
    )
    if is_temporal and num_cols:
        suggestions.append({
            "type": "line",
            "title": f"Évolution de {num_cols[0]['header']}",
            "categories_col": cat["col_idx"] if cat else 1,
            "data_col": num_cols[0]["col_idx"],
            "reason": "série temporelle détectée"
        })

    # Suggestion camembert si 2-8 catégories
    if cat and 2 <= len(rows) <= 8 and num_cols:
        suggestions.append({
            "type": "pie",
            "title": f"Répartition de {num_cols[0]['header']}",
            "categories_col": cat["col_idx"] if cat else 1,
            "data_col": num_cols[0]["col_idx"],
            "reason": "peu de catégories — répartition idéale"
        })

    return suggestions[:3]  # max 3 suggestions


# ─── Main ────────────────────────────────────────────────────────────────────

def main():
    try:
        raw = sys.stdin.buffer.read()
        data = json.loads(raw)
    except Exception as e:
        sys.stderr.write(f"JSON parse error: {e}\n")
        sys.exit(1)

    # Mode parse : retourne JSON structuré (texte, pas binaire)
    if data.get("mode") == "parse":
        try:
            result = parse_excel_for_ai(data)
        except Exception:
            sys.stderr.write(f"Parse error:\n{traceback.format_exc()}\n")
            sys.exit(2)
        sys.stdout.buffer.write(result)
        return

    doc_type = data.get("document_type", "pptx").lower()

    generators = {
        "pptx": (generate_pptx, "pptx"),
        "powerpoint": (generate_pptx, "pptx"),
        "docx": (generate_docx, "docx"),
        "word": (generate_docx, "docx"),
        "pdf": (generate_pdf, "pdf"),
        "xlsx": (generate_xlsx, "xlsx"),
        "excel": (generate_xlsx, "xlsx"),
    }

    entry = generators.get(doc_type)
    if not entry:
        sys.stderr.write(f"Unknown document_type: {doc_type}\n")
        sys.exit(4)

    gen_fn, _ = entry
    try:
        result = gen_fn(data)
    except Exception:
        sys.stderr.write(f"Generation error:\n{traceback.format_exc()}\n")
        sys.exit(2)

    sys.stdout.buffer.write(result)


if __name__ == "__main__":
    main()
