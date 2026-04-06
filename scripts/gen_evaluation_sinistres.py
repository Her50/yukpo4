#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génération des documents d'évaluation et validation des acquis
Formation : Gestion des Sinistres — Toutes Branches
Formateur : TALOM Eric

Fichiers générés :
  1. Evaluation_Guide_Formateur.pptx  — présentation corrigée pour le formateur
  2. Evaluation_Cahier_Participant.docx — cahier à remplir par le stagiaire
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from docx import Document
from docx.shared import Pt as DPt, Inches as DInches, RGBColor as DRGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import copy

# ── PALETTE PPT ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
BLUE   = RGBColor(0x2E, 0x86, 0xAB)
ORANGE = RGBColor(0xF1, 0x8F, 0x01)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xC0, 0x39, 0x2B)
TEAL   = RGBColor(0x16, 0xA0, 0x85)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
YELLOW = RGBColor(0xF3, 0xC3, 0x00)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# UTILITAIRES PPT
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp

def txb(slide, text, x, y, w, h, sz=14, color=WHITE, bold=False,
        italic=False, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = name
    return tb

def set_bg(slide, color):
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = color

def header(slide, title, label=""):
    add_rect(slide, 0, 0, W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), W, Inches(0.07), fill=ORANGE)
    if label:
        txb(slide, label, Inches(0.35), Inches(0.07), Inches(9), Inches(0.42),
            sz=11, color=RGBColor(0xBB, 0xCC, 0xDD))
    txb(slide, title, Inches(0.4), Inches(0.16), Inches(12.5), Inches(0.95),
        sz=24, bold=True)

def bullets_box(slide, items, x, y, w, h, sz=13, color=DGRAY, prefix="▸  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = prefix + item
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.name = "Calibri"

def section_div(prs, title, subtitle="", color=NAVY):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, color)
    add_rect(slide, 0, 0, Inches(0.28), H, fill=ORANGE)
    add_rect(slide, 0, H - Inches(0.15), W, Inches(0.15), fill=ORANGE)
    txb(slide, title, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.8),
        sz=44, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(4.2), Inches(4.55), Inches(5.0), Inches(0.08), fill=ORANGE)
    if subtitle:
        txb(slide, subtitle, Inches(1.5), Inches(4.75), Inches(10.3), Inches(1.5),
            sz=18, italic=True, align=PP_ALIGN.CENTER,
            color=RGBColor(0xCC, 0xDD, 0xEE))
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# DONNÉES — QCM 30 QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════

QCM = [
    # --- CONNAISSANCES GÉNÉRALES (20 questions) ---
    {
        "n": 1, "cat": "Connaissances Générales",
        "q": "Quel article du Code CIMA fixe le délai de 30 jours pour la proposition d'indemnisation ?",
        "a": ["Art. 12", "Art. 15", "Art. 16", "Art. 28"],
        "rep": "C", "expl": "L'art. 16 CIMA impose à l'assureur de faire une proposition d'indemnisation dans les 30 jours suivant l'accord des parties sur les montants."
    },
    {
        "n": 2, "cat": "Connaissances Générales",
        "q": "Le principe indemnitaire signifie que :",
        "a": ["L'assuré peut s'enrichir grâce à l'assurance",
              "L'indemnité ne peut dépasser la valeur du préjudice réel subi",
              "L'assureur peut refuser tout sinistre sans motif",
              "L'assuré est toujours indemnisé à 100%"],
        "rep": "B", "expl": "L'assurance ne doit pas être source d'enrichissement. L'indemnité est plafonnée à la valeur réelle du préjudice (art. CIMA + principe général)."
    },
    {
        "n": 3, "cat": "Connaissances Générales",
        "q": "Quelle est la durée légale de prescription des actions nées d'un contrat d'assurance selon le Code CIMA ?",
        "a": ["1 an", "2 ans", "5 ans", "10 ans"],
        "rep": "B", "expl": "L'art. 28 CIMA fixe la prescription biennale (2 ans) pour les actions nées du contrat d'assurance."
    },
    {
        "n": 4, "cat": "Connaissances Générales",
        "q": "La subrogation en assurance consiste à :",
        "a": ["Annuler le contrat après un sinistre",
              "Permettre à l'assureur de se substituer à l'assuré pour exercer un recours contre le tiers responsable",
              "Obliger l'assuré à rembourser l'assureur dans tous les cas",
              "Suspendre les garanties après déclaration de sinistre"],
        "rep": "B", "expl": "Après règlement du sinistre, l'assureur est subrogé dans les droits et actions de l'assuré contre le tiers responsable."
    },
    {
        "n": 5, "cat": "Connaissances Générales",
        "q": "Qu'est-ce qu'une franchise absolue ?",
        "a": ["Une somme déduite de l'indemnité, quelle que soit l'importance du sinistre",
              "Une garantie couvrant tous les sinistres sans plafond",
              "Un sinistre inférieur à la franchise qui reste entièrement à charge de l'assureur",
              "Un plafond au-delà duquel l'assureur ne couvre plus"],
        "rep": "A", "expl": "La franchise absolue est toujours déduite de l'indemnité (même si le sinistre est très élevé), contrairement à la franchise relative."
    },
    {
        "n": 6, "cat": "Connaissances Générales",
        "q": "En cas de double assurance (même risque assuré par deux assureurs différents), chaque assureur :",
        "a": ["Doit payer la totalité du sinistre",
              "Peut refuser d'indemniser",
              "Contribue proportionnellement à hauteur de sa quote-part",
              "Est totalement déchargé de sa responsabilité"],
        "rep": "C", "expl": "Principe de contribution : chaque assureur participe proportionnellement au rapport entre sa garantie et la somme totale assurée."
    },
    {
        "n": 7, "cat": "Connaissances Générales",
        "q": "Quelle est la première action du gestionnaire sinistres à réception d'une déclaration ?",
        "a": ["Mandater l'expert immédiatement",
              "Envoyer un refus de garantie",
              "Enregistrer le dossier et envoyer un accusé de réception sous 24-48h",
              "Contacter les pompiers"],
        "rep": "C", "expl": "L'accusé de réception dans les 24-48h est la première obligation légale et professionnelle du gestionnaire."
    },
    {
        "n": 8, "cat": "Connaissances Générales",
        "q": "La règle proportionnelle s'applique quand :",
        "a": ["L'assuré a déclaré la valeur exacte de ses biens",
              "La valeur assurée est inférieure à la valeur réelle du bien",
              "L'assureur refuse la garantie",
              "Le sinistre dépasse le plafond de garantie"],
        "rep": "B", "expl": "Si VA < VR, l'assuré est son propre assureur pour la différence. Formule : I = Sinistre × (VA / VR)."
    },
    {
        "n": 9, "cat": "Connaissances Générales",
        "q": "Qu'est-ce que la déchéance de garantie ?",
        "a": ["Un avenant au contrat réduisant les primes",
              "La perte du droit à indemnisation suite au non-respect d'une obligation contractuelle",
              "Une extension de garantie automatique",
              "Un remboursement partiel de la prime"],
        "rep": "B", "expl": "La déchéance sanctionne le non-respect d'obligations (délai de déclaration, mesures conservatoires, etc.) par la perte du droit à garantie."
    },
    {
        "n": 10, "cat": "Connaissances Générales",
        "q": "La vétusté en assurance représente :",
        "a": ["La prime supplémentaire due après un sinistre",
              "L'augmentation de valeur d'un bien avec le temps",
              "La dépréciation d'un bien due à l'usure et au vieillissement",
              "Un type de franchise applicable aux biens anciens"],
        "rep": "C", "expl": "La vétusté réduit l'indemnité pour tenir compte de la dépréciation du bien (usure normale). La clause 'valeur à neuf' la supprime."
    },
    {
        "n": 11, "cat": "Connaissances Générales",
        "q": "En branche automobile, la Convention IDA (Indemnisation Directe de l'Assuré) permet :",
        "a": ["À l'assureur de refuser d'indemniser en cas de responsabilité partagée",
              "À l'assuré d'être indemnisé directement par son propre assureur sans attendre le recours contre le tiers",
              "Aux deux assureurs de partager les frais d'expertise",
              "De supprimer la franchise pour les accidents"],
        "rep": "B", "expl": "La Convention IDA accélère l'indemnisation en permettant à l'assureur de l'assuré de régler directement, puis de se faire rembourser via la Convention IRSA."
    },
    {
        "n": 12, "cat": "Connaissances Générales",
        "q": "En branche incendie, les 'dégâts des eaux' (DDO) sont :",
        "a": ["Toujours couverts automatiquement sans option",
              "Une extension facultative à souscrire spécifiquement",
              "Exclus de toutes les polices incendie sans exception",
              "Couverts uniquement pour les entreprises"],
        "rep": "B", "expl": "Les DDO constituent une garantie complémentaire (option) qui doit être expressément souscrite dans la police incendie."
    },
    {
        "n": 13, "cat": "Connaissances Générales",
        "q": "En transports facultés maritimes, le connaissement (Bill of Lading) est :",
        "a": ["Un simple bon de commande",
              "Un titre de propriété de la marchandise faisant preuve du contrat de transport",
              "Un certificat phytosanitaire",
              "Une facture douanière"],
        "rep": "B", "expl": "Le B/L est à la fois un reçu, un titre de propriété et la preuve du contrat de transport. Son original est indispensable au dossier sinistre."
    },
    {
        "n": 14, "cat": "Connaissances Générales",
        "q": "En branche santé, le tiers payant signifie que :",
        "a": ["L'assuré paie une partie des soins",
              "L'assureur règle directement le prestataire sans avance de frais par l'assuré",
              "Un tiers garant se substitue à l'assuré défaillant",
              "Les frais sont remboursés 3 mois après les soins"],
        "rep": "B", "expl": "Le mécanisme du tiers payant supprime l'avance de frais pour l'assuré : l'assureur règle directement l'établissement de santé ou le prestataire."
    },
    {
        "n": 15, "cat": "Connaissances Générales",
        "q": "La consolidation en branche Individuel Accident désigne :",
        "a": ["La fusion de deux polices d'assurance",
              "La date à partir de laquelle l'état de santé de la victime est stabilisé et ne nécessite plus de soins",
              "Le remboursement total des frais médicaux",
              "La clôture administrative du dossier sinistre"],
        "rep": "B", "expl": "La consolidation est la date médicale clé à partir de laquelle le taux d'IPP peut être fixé définitivement."
    },
    {
        "n": 16, "cat": "Connaissances Générales",
        "q": "La Clause A en assurance transport facultés maritimes couvre :",
        "a": ["Uniquement les naufrages et échouages",
              "Les risques listés de manière limitative",
              "Tous les risques sauf les exclusions expressément prévues (all risks)",
              "Uniquement les avaries communes"],
        "rep": "C", "expl": "La Clause A est la plus large : elle couvre tous les risques non expressément exclus. C'est la couverture 'tous risques' des facultés."
    },
    {
        "n": 17, "cat": "Connaissances Générales",
        "q": "En cas de destruction totale d'un véhicule assuré en dommages, l'indemnité nette est :",
        "a": ["La valeur à neuf du véhicule",
              "Le coût de réparation estimé",
              "La valeur vénale au jour du sinistre moins la valeur de l'épave",
              "Le capital souscrit dans la police"],
        "rep": "C", "expl": "Pour une DT, l'indemnité = Valeur Vénale (VV) − Valeur épave, déduction faite de la franchise éventuelle."
    },
    {
        "n": 18, "cat": "Connaissances Générales",
        "q": "L'avarie commune en transport maritime se produit lorsque :",
        "a": ["Un seul colis est endommagé lors du transport",
              "Le capitaine prend une décision volontaire et raisonnable pour sauver l'expédition commune en péril",
              "Le transporteur refuse de livrer la marchandise",
              "L'assureur refuse la garantie"],
        "rep": "B", "expl": "L'avarie commune implique un acte volontaire (ex. jet à la mer) pour sauver le navire et la cargaison. Les pertes sont réparties entre tous les intéressés (Règles d'York-Anvers)."
    },
    {
        "n": 19, "cat": "Connaissances Générales",
        "q": "En branche individuel accident, l'ITT (Incapacité Temporaire Totale) ouvre droit à :",
        "a": ["Un capital unique versé en fin d'arrêt",
              "Une indemnité journalière versée pendant la durée de l'arrêt de travail (hors franchise)",
              "Une rente viagère mensuelle",
              "Le remboursement des frais médicaux uniquement"],
        "rep": "B", "expl": "L'ITT donne lieu à une indemnité journalière (IJ) × nombre de jours d'arrêt, après déduction de la franchise temporelle (généralement 3 jours)."
    },
    {
        "n": 20, "cat": "Connaissances Générales",
        "q": "Le délai légal de conservation des dossiers sinistres clos selon le Code CIMA est de :",
        "a": ["2 ans", "5 ans", "10 ans", "30 ans"],
        "rep": "C", "expl": "La prescription biennale (art. 28 CIMA) et les recours possibles imposent une conservation minimale de 10 ans après clôture du dossier."
    },
    # --- RÉGLEMENTATION CIMA (10 questions) ---
    {
        "n": 21, "cat": "Réglementation CIMA",
        "q": "Le Code CIMA a été institué par le Traité signé à :",
        "a": ["Dakar le 14 juillet 1990",
              "Yaoundé le 10 juillet 1992",
              "Abidjan le 1er janvier 1995",
              "Paris le 5 mars 1988"],
        "rep": "B", "expl": "Le Traité instituant la Conférence Interafricaine des Marchés d'Assurances (CIMA) a été signé à Yaoundé le 10 juillet 1992."
    },
    {
        "n": 22, "cat": "Réglementation CIMA",
        "q": "L'article 15 du Code CIMA concerne :",
        "a": ["Le calcul des primes d'assurance",
              "L'obligation d'accusé de réception et d'instruction de la déclaration de sinistre",
              "Les conditions de résiliation du contrat",
              "La nomination des experts agréés"],
        "rep": "B", "expl": "L'art. 15 CIMA impose à l'assureur d'accuser réception et d'instruire la déclaration de sinistre dans des délais stricts."
    },
    {
        "n": 23, "cat": "Réglementation CIMA",
        "q": "Selon l'article 18 du Code CIMA, en cas de retard de règlement, les pénalités dues par l'assureur sont :",
        "a": ["Les intérêts légaux uniquement",
              "Une pénalité forfaitaire de 10% du capital",
              "Les intérêts légaux majorés de 50%",
              "Aucune pénalité n'est prévue par la loi"],
        "rep": "C", "expl": "L'art. 18 CIMA prévoit que l'indemnité produit intérêts de plein droit au taux légal majoré de 50% en cas de retard fautif de l'assureur."
    },
    {
        "n": 24, "cat": "Réglementation CIMA",
        "q": "La CIMA regroupe actuellement combien d'États membres ?",
        "a": ["7", "10", "14", "20"],
        "rep": "C", "expl": "La CIMA regroupe 14 États membres d'Afrique subsaharienne francophone (Bénin, Burkina Faso, Cameroun, Centrafrique, Comores, Congo, Côte d'Ivoire, Gabon, Guinée Bissau, Guinée Équatoriale, Mali, Niger, Sénégal, Tchad, Togo)."
    },
    {
        "n": 25, "cat": "Réglementation CIMA",
        "q": "En assurance vie et individuel accident, la fausse déclaration intentionnelle entraîne :",
        "a": ["Une réduction proportionnelle de l'indemnité",
              "La nullité du contrat avec conservation des primes par l'assureur",
              "Une simple majoration de prime",
              "Un avertissement sans conséquence financière"],
        "rep": "B", "expl": "La fausse déclaration intentionnelle entraîne la nullité absolue du contrat selon l'art. 18 CIMA, l'assureur conservant les primes à titre de dommages-intérêts."
    },
    {
        "n": 26, "cat": "Réglementation CIMA",
        "q": "En assurance RC automobile obligatoire, la victime d'un accident peut :",
        "a": ["Seulement agir contre le conducteur responsable",
              "Agir directement contre l'assureur du responsable (action directe)",
              "Agir uniquement devant les juridictions civiles",
              "Seulement demander réparation à l'État"],
        "rep": "B", "expl": "Le Code CIMA consacre le droit d'action directe de la victime contre l'assureur RC du responsable, sans passer par l'assuré."
    },
    {
        "n": 27, "cat": "Réglementation CIMA",
        "q": "Un sinistre est déclaré 3 ans après sa survenance. En droit CIMA, cette action est :",
        "a": ["Recevable, car aucun délai n'est prévu",
              "Prescrite, car le délai légal de 2 ans est dépassé",
              "Recevable si l'assuré fournit une bonne raison",
              "Prescrite uniquement en assurance vie"],
        "rep": "B", "expl": "La prescription biennale (art. 28 CIMA) s'applique à toutes les branches. Passé 2 ans sans interruption, l'action est forclose."
    },
    {
        "n": 28, "cat": "Réglementation CIMA",
        "q": "Quelle est l'autorité de supervision des compagnies d'assurance dans l'espace CIMA ?",
        "a": ["La Banque Centrale (BCEAO / BEAC)",
              "La Commission Régionale de Contrôle des Assurances (CRCA)",
              "Le Ministère des Finances de chaque État",
              "L'UEMOA"],
        "rep": "B", "expl": "La CRCA est l'organe supranational de supervision et de contrôle des entreprises d'assurance dans l'espace CIMA."
    },
    {
        "n": 29, "cat": "Réglementation CIMA",
        "q": "La RC automobile est qualifiée d'obligatoire car :",
        "a": ["Elle est souscrite volontairement par tous les conducteurs",
              "La loi impose à tout propriétaire de véhicule de couvrir sa responsabilité envers les tiers",
              "Elle est imposée uniquement aux flottes d'entreprises",
              "Elle est exigée uniquement pour les véhicules de plus de 5 ans"],
        "rep": "B", "expl": "La RC Auto obligatoire (RCO) est imposée par la loi à tout propriétaire ou conducteur de véhicule motorisé pour protéger les victimes tierces."
    },
    {
        "n": 30, "cat": "Réglementation CIMA",
        "q": "En cas de désaccord entre l'assuré et l'assureur sur le montant de l'indemnité, l'assuré peut :",
        "a": ["Uniquement accepter la proposition de l'assureur",
              "Demander une contre-expertise aux frais exclusifs de l'assureur",
              "Demander une contre-expertise contradictoire, les frais étant partagés en cas de désaccord persistant",
              "Saisir immédiatement le tribunal sans autre recours préalable"],
        "rep": "C", "expl": "L'assuré a le droit de solliciter une contre-expertise. Si le désaccord persiste, les deux experts désignent un tiers arbitre et les frais sont partagés."
    },
]

# ══════════════════════════════════════════════════════════════════════════════
# DONNÉES — CAS PRATIQUES
# ══════════════════════════════════════════════════════════════════════════════

CAS_PRATIQUES = [
    {
        "num": 1, "branche": "AUTOMOBILE", "color": BLUE, "duree": "20 min",
        "titre": "Sinistre RC — Collision en agglomération",
        "contexte": (
            "M. KOUASSI Aimé, assuré chez ASSUR-PLUS depuis 4 ans, déclare un accident "
            "de la circulation survenu le 15 mars 2026 à 14h30, avenue de l'Indépendance à Abidjan. "
            "Son véhicule Toyota Corolla (immatriculé CI-1234-AB, valeur vénale estimée : 7 500 000 FCFA) "
            "a été percuté à l'arrière par le véhicule de M. BAMBA Seydou alors que M. KOUASSI était "
            "à l'arrêt à un feu rouge. M. BAMBA reconnaît sa responsabilité totale sur le constat amiable. "
            "Devis de réparation du garage agréé : 1 850 000 FCFA. Franchise contractuelle : 75 000 FCFA."
        ),
        "questions": [
            "Q1. Listez les documents que vous devez réclamer à M. KOUASSI pour instruire ce dossier.",
            "Q2. Quelle garantie de sa police M. KOUASSI peut-il actionner ? Justifiez.",
            "Q3. Quel est le montant net de l'indemnité à verser à M. KOUASSI ?",
            "Q4. L'assureur peut-il exercer un recours ? Contre qui et sur quel fondement ?",
            "Q5. Quel délai maximum avez-vous pour formuler votre proposition d'indemnisation ?",
        ],
        "corrige": [
            "Q1. Documents : constat amiable signé des deux parties, carte grise, permis de conduire, attestation d'assurance valide, devis de réparation du garage agréé, photos des dommages, PV de police si disponible.",
            "Q2. Garantie 'Dommages Collision' (sinistre avec tiers identifié et responsable) — ou 'Dommages Tous Accidents' si souscrite. La RC de M. BAMBA est également engagée.",
            "Q3. Indemnité nette = 1 850 000 − 75 000 (franchise) = 1 775 000 FCFA.",
            "Q4. Oui — recours subrogatoire contre l'assureur RC de M. BAMBA via la Convention IRSA. Base : responsabilité totale de M. BAMBA reconnue sur le constat.",
            "Q5. 30 jours maximum à compter de la réception de toutes les pièces justificatives (Art. 16 CIMA).",
        ]
    },
    {
        "num": 2, "branche": "INCENDIE", "color": RED, "duree": "20 min",
        "titre": "Incendie d'un commerce — Règle proportionnelle",
        "contexte": (
            "Mme DIALLO Fatou exploite une boutique de prêt-à-porter à Dakar. "
            "Elle a souscrit une police incendie avec une valeur assurée de 15 000 000 FCFA "
            "pour les marchandises. Le 02 avril 2026 à 02h15, un incendie d'origine électrique "
            "détruit 60% de ses stocks. L'expertise contradictoire évalue : "
            "— Valeur réelle des marchandises au jour du sinistre : 25 000 000 FCFA. "
            "— Dommages constatés : 12 000 000 FCFA. "
            "Franchise absolue : 250 000 FCFA. "
            "Mme DIALLO demande une indemnité de 12 000 000 FCFA."
        ),
        "questions": [
            "Q1. Quels documents constitutifs du dossier sinistre devez-vous rassembler ?",
            "Q2. La valeur assurée est-elle suffisante ? Quelle conséquence cela entraîne-t-il ?",
            "Q3. Calculez le montant de l'indemnité nette due à Mme DIALLO (montrez le calcul).",
            "Q4. Mme DIALLO conteste votre décision. Quel est son recours possible ?",
            "Q5. Comment ce dossier doit-il être archivé et pendant combien de temps ?",
        ],
        "corrige": [
            "Q1. Documents : déclaration écrite, PV pompiers, PV police si incendie criminel suspecté, inventaire des marchandises détruites, factures d'achat, rapport d'expertise contradictoire, photos des dommages.",
            "Q2. VA (15M) < VR (25M) → insuffisance d'assurance. La règle proportionnelle s'applique : Mme DIALLO est son propre assureur pour 10/25 de la valeur.",
            "Q3. Indemnité brute = 12 000 000 × (15 000 000 / 25 000 000) = 12 000 000 × 0,60 = 7 200 000 FCFA. Indemnité nette = 7 200 000 − 250 000 (franchise) = 6 950 000 FCFA.",
            "Q4. Mme DIALLO peut demander une contre-expertise contradictoire. En cas de désaccord persistant, un tiers arbitre est désigné par les deux experts.",
            "Q5. Archivage physique (chemise Incendie) + scan GED, indexé au numéro de sinistre. Conservation : 10 ans minimum après clôture (prescription biennale + délais recours).",
        ]
    },
    {
        "num": 3, "branche": "SANTÉ", "color": GREEN, "duree": "20 min",
        "titre": "Hospitalisation — Prise en charge et liquidation",
        "contexte": (
            "M. TOSSOU Rodrigue, salarié dans une entreprise abonnée au contrat collectif santé "
            "d'ASSUR-VIE, est hospitalisé en urgence le 10 avril 2026 pour appendicite aiguë. "
            "Durée d'hospitalisation : 4 jours. "
            "Facture de la clinique agréée : 890 000 FCFA (dont honoraires chirurgien : 250 000 FCFA, "
            "frais de séjour : 480 000 FCFA, médicaments : 160 000 FCFA). "
            "Tableau de garanties : hospitalisation prise en charge à 80%, médicaments à 70%. "
            "Ticket modérateur à charge de l'assuré : 20% sur séjour + honoraires, 30% sur médicaments. "
            "Forfait journalier : non applicable (clinique agréée réseau)."
        ),
        "questions": [
            "Q1. M. TOSSOU a été hospitalisé en urgence. La PEC préalable était-elle obligatoire ? Quel protocole appliquer ?",
            "Q2. Quels documents devez-vous exiger pour traiter la demande de remboursement ?",
            "Q3. Calculez le montant pris en charge par l'assureur et le ticket modérateur de M. TOSSOU.",
            "Q4. Dans quel délai devez-vous régler la facture une fois le dossier complet reçu ?",
            "Q5. Citez deux motifs légaux de refus de prise en charge en branche santé.",
        ],
        "corrige": [
            "Q1. En urgence, la PEC préalable n'est pas obligatoire avant l'admission. Le protocole est : PEC rétroactive dans les 72h suivant l'admission, sur présentation du certificat médical d'urgence.",
            "Q2. Documents : bulletin d'hospitalisation, certificat médical (urgence), facture détaillée de la clinique, ordonnances médicales, carte d'assuré valide, bulletin de prise en charge (BPC).",
            "Q3. Séjour + honoraires : (250 000 + 480 000) × 80% = 584 000 FCFA pris en charge. Médicaments : 160 000 × 70% = 112 000 FCFA. TOTAL assureur = 696 000 FCFA. Ticket modérateur M. TOSSOU = 890 000 − 696 000 = 194 000 FCFA.",
            "Q4. Délai de règlement : 10 jours ouvrables maximum après réception du dossier complet (conditions générales du contrat).",
            "Q5. Motifs légaux de refus : (a) soins antérieurs à la date de prise d'effet (délai de carence non respecté) ; (b) actes ou pathologies expressément exclus (esthétique, maladies antérieures non déclarées, etc.).",
        ]
    },
]

# ══════════════════════════════════════════════════════════════════════════════
# GRILLE MISE EN SITUATION
# ══════════════════════════════════════════════════════════════════════════════

GRILLE_SITUATION = [
    ("Communication & Accueil",        "Vocabulaire adapté, ton professionnel, empathie",                   "/4"),
    ("Écoute active",                  "Reformulation, questions pertinentes, prise de notes",               "/4"),
    ("Maîtrise réglementaire",         "Respect des délais CIMA cités, justification des décisions",         "/4"),
    ("Gestion du conflit",             "Calme, argumentation factuelle, proposition de solution",            "/4"),
    ("Rédaction courrier de règlement","Clarté, structure, mentions obligatoires, ton adapté",               "/4"),
    ("TOTAL",                          "",                                                                   "/20"),
]

# ══════════════════════════════════════════════════════════════════════════════
# BAREME GLOBAL
# ══════════════════════════════════════════════════════════════════════════════

BAREME = [
    ("QCM Théorique",       "30 questions × 0,5 pt",      "15 pts",  "10,5 pts (70%)"),
    ("Cas Pratiques",       "3 dossiers × 10 pts",        "30 pts",  "21 pts (70%)"),
    ("Mise en Situation",   "Grille comportementale",      "20 pts",  "14 pts (70%)"),
    ("Participation active","Attitude, assiduité",         "5 pts",   "Appréciation"),
    ("TOTAL GÉNÉRAL",       "",                            "70 pts",  "49 pts (70%)"),
]


# ══════════════════════════════════════════════════════════════════════════════
# CONSTRUCTION PPT — GUIDE FORMATEUR
# ══════════════════════════════════════════════════════════════════════════════

def ppt_cover():
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    add_rect(slide, 0, 0, Inches(5.1), H, fill=NAVY)
    add_rect(slide, Inches(5.1), 0, Inches(0.13), H, fill=ORANGE)
    add_rect(slide, 0, H - Inches(0.18), W, Inches(0.18), fill=ORANGE)
    txb(slide, "GUIDE FORMATEUR", Inches(0.35), Inches(0.8), Inches(4.4), Inches(0.7),
        sz=14, bold=True, align=PP_ALIGN.CENTER, color=ORANGE, name="Calibri Light")
    txb(slide, "EVALUATION\nDES ACQUIS", Inches(0.3), Inches(1.7), Inches(4.55), Inches(1.9),
        sz=38, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.5), Inches(3.75), Inches(4.1), Inches(0.06), fill=ORANGE)
    txb(slide, "Gestion des Sinistres\nToutes Branches", Inches(0.3), Inches(3.95),
        Inches(4.55), Inches(1.0), sz=14, italic=True, align=PP_ALIGN.CENTER,
        color=RGBColor(0xCC, 0xDD, 0xEE))
    txb(slide, "Formateur : TALOM Eric", Inches(0.3), Inches(5.1), Inches(4.55), Inches(0.5),
        sz=15, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    txb(slide, "Avril 2026  |  Confidentiel Formateur", Inches(0.3), Inches(5.65),
        Inches(4.55), Inches(0.4), sz=11, align=PP_ALIGN.CENTER,
        color=RGBColor(0x99, 0xBB, 0xDD))

    txb(slide, "ÉVALUATION &\nVALIDATION\nDES ACQUIS",
        Inches(5.5), Inches(1.4), Inches(7.5), Inches(2.6),
        sz=34, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_rect(slide, Inches(5.5), Inches(4.15), Inches(5.5), Inches(0.07), fill=ORANGE)
    txb(slide, "Formation de Remise à Niveau\nProfessionnels de l'Assurance",
        Inches(5.5), Inches(4.35), Inches(7.4), Inches(1.2),
        sz=16, color=DGRAY, align=PP_ALIGN.LEFT)
    txb(slide, "QCM (30 questions)  |  Cas Pratiques (3 dossiers)  |  Mise en Situation",
        Inches(5.5), Inches(5.6), Inches(7.4), Inches(0.5),
        sz=12, color=ORANGE, bold=True, align=PP_ALIGN.LEFT)
    txb(slide, "Durée totale : 2h15  |  Note minimale : 70%",
        Inches(5.5), Inches(6.15), Inches(7.4), Inches(0.4),
        sz=12, color=DGRAY, align=PP_ALIGN.LEFT, italic=True)


def ppt_instructions():
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    header(slide, "Instructions Générales — Évaluation", label="ÉVALUATION — Instructions")
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=ORANGE)

    blocs = [
        (NAVY,   "PARTIE 1 — QCM",           "30 min", [
            "30 questions à choix unique (1 seule bonne réponse par question)",
            "Valeur : 0,5 point par question correcte — 0 point si fausse ou non répondue",
            "Score maximum : 15 points | Seuil de validation : 10,5 pts (70%)",
            "Aucune calculatrice, ni document autorisé",
        ]),
        (BLUE,   "PARTIE 2 — CAS PRATIQUES", "60 min", [
            "3 dossiers indépendants (Automobile, Incendie, Santé)",
            "Valeur : 10 points par dossier — Score maximum : 30 points",
            "Seuil de validation par dossier : 7/10 | Global : 21/30",
            "Justifiez toujours vos calculs et référencez les articles CIMA pertinents",
        ]),
        (TEAL,   "PARTIE 3 — MISE EN SITUATION", "45 min", [
            "Simulation en binômes (assuré / gestionnaire) sur 3 scénarios",
            "Évaluation par grille comportementale : 20 points",
            "Seuil de validation : 14/20 | Critères : communication, maîtrise, gestion conflit",
            "Rotation des rôles toutes les 15 minutes",
        ]),
    ]
    for i, (col, t, dur, items) in enumerate(blocs):
        bx = Inches(0.35) + i * Inches(4.29)
        by = Inches(1.48)
        bw = Inches(4.1)
        bh = Inches(5.85)
        add_rect(slide, bx, by, bw, bh, fill=LGRAY)
        add_rect(slide, bx, by, bw, Inches(0.85), fill=col)
        txb(slide, t, bx + Inches(0.08), by + Inches(0.04), bw - Inches(0.16),
            Inches(0.48), sz=14, bold=True, align=PP_ALIGN.CENTER)
        txb(slide, f"Duree : {dur}", bx + Inches(0.08), by + Inches(0.54),
            bw - Inches(0.16), Inches(0.3), sz=11, italic=True,
            align=PP_ALIGN.CENTER, color=WHITE)
        bullets_box(slide, items, bx + Inches(0.15), by + Inches(0.92),
                    bw - Inches(0.3), bh - Inches(1.05), sz=12)


def ppt_qcm_slide(questions_group, label):
    """Affiche 5-6 questions QCM avec réponses cochées"""
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    header(slide, f"QCM — {label}", label=f"PARTIE 1 — QCM THEORIQUE")
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=BLUE)

    row_h = Inches(5.85) / len(questions_group)
    for i, q in enumerate(questions_group):
        qy = Inches(1.45) + i * row_h
        bg = LGRAY if i % 2 == 0 else WHITE
        add_rect(slide, Inches(0.22), qy, Inches(12.89), row_h - Inches(0.04), fill=bg)

        # Numéro
        add_rect(slide, Inches(0.22), qy, Inches(0.55), row_h - Inches(0.04), fill=BLUE)
        txb(slide, str(q["n"]), Inches(0.22), qy + Inches(0.05),
            Inches(0.55), row_h - Inches(0.12), sz=14, bold=True, align=PP_ALIGN.CENTER)

        # Question
        txb(slide, q["q"], Inches(0.85), qy + Inches(0.04),
            Inches(5.7), row_h - Inches(0.1), sz=11, color=NAVY, bold=True)

        # Choix + réponse correcte
        choices_text = ""
        for j, choice in enumerate(q["a"]):
            letter = chr(65 + j)
            mark = "  ✓" if letter == q["rep"] else ""
            choices_text += f"{letter}. {choice}{mark}\n"

        txb(slide, choices_text.strip(), Inches(6.65), qy + Inches(0.04),
            Inches(5.9), row_h - Inches(0.1), sz=10, color=DGRAY)

        # Réponse surlignée
        rep_colors = {"A": RGBColor(0xE8, 0xF5, 0xE9), "B": RGBColor(0xE8, 0xF5, 0xE9),
                      "C": RGBColor(0xE8, 0xF5, 0xE9), "D": RGBColor(0xE8, 0xF5, 0xE9)}


def ppt_corrige_qcm():
    """Slide corrigé QCM — tableau récapitulatif"""
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    header(slide, "Corrige QCM — Tableau des Reponses Correctes", label="PARTIE 1 — CORRIGE")
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=GREEN)

    # Tableau 3 colonnes × 10 lignes pour les 30 questions
    data_col1 = [(q["n"], q["rep"]) for q in QCM[:10]]
    data_col2 = [(q["n"], q["rep"]) for q in QCM[10:20]]
    data_col3 = [(q["n"], q["rep"]) for q in QCM[20:30]]

    for ci, col_data in enumerate([data_col1, data_col2, data_col3]):
        cx = Inches(0.4) + ci * Inches(4.3)
        cy = Inches(1.45)
        cw = Inches(4.1)
        ch = Inches(5.9)
        n_rows = len(col_data) + 1
        tbl = slide.shapes.add_table(n_rows, 3, cx, cy, cw, ch).table
        for j, hdr in enumerate(["N°", "Bonne Reponse", "Categorie"]):
            cell = tbl.cell(0, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = hdr; r.font.size = Pt(11); r.font.bold = True
            r.font.color.rgb = WHITE; r.font.name = "Calibri"
        for ri, (num, rep) in enumerate(col_data):
            q_obj = QCM[num - 1]
            for j, val in enumerate([str(num), rep, q_obj["cat"][:12]]):
                cell = tbl.cell(ri + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LGRAY if ri % 2 == 0 else WHITE
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = val; r.font.size = Pt(12)
                r.font.color.rgb = GREEN if j == 1 else DGRAY
                r.font.bold = (j == 1)
                r.font.name = "Calibri"


def ppt_cas_pratique(cas):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    header(slide,
           f"Cas Pratique {cas['num']} — Branche {cas['branche']} ({cas['duree']})",
           label=f"PARTIE 2 — CAS PRATIQUES")
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=cas["color"])

    # Titre
    txb(slide, cas["titre"], Inches(0.35), Inches(1.37), Inches(12.6), Inches(0.5),
        sz=16, bold=True, color=cas["color"])

    # Contexte (fond coloré)
    add_rect(slide, Inches(0.3), Inches(1.9), Inches(12.73), Inches(2.1), fill=LGRAY)
    add_rect(slide, Inches(0.3), Inches(1.9), Inches(0.12), Inches(2.1), fill=cas["color"])
    txb(slide, "CONTEXTE :", Inches(0.5), Inches(1.95), Inches(2.0), Inches(0.35),
        sz=11, bold=True, color=cas["color"])
    txb(slide, cas["contexte"], Inches(0.5), Inches(2.3), Inches(12.5), Inches(1.65),
        sz=11, color=DGRAY, italic=False)

    # Questions
    txb(slide, "QUESTIONS :", Inches(0.35), Inches(4.1), Inches(3.0), Inches(0.35),
        sz=12, bold=True, color=NAVY)
    bullets_box(slide, cas["questions"], Inches(0.35), Inches(4.48),
                Inches(12.6), Inches(2.8), sz=12, color=DGRAY, prefix="")


def ppt_corrige_cas(cas):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    header(slide,
           f"Corrige — Cas Pratique {cas['num']} : {cas['branche']}",
           label="PARTIE 2 — CORRIGE FORMATEUR")
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=GREEN)

    txb(slide, cas["titre"], Inches(0.35), Inches(1.37), Inches(12.6), Inches(0.45),
        sz=15, bold=True, color=cas["color"])

    tb = slide.shapes.add_textbox(Inches(0.35), Inches(1.88), Inches(12.6), Inches(5.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ligne in enumerate(cas["corrige"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = ligne
        r.font.size = Pt(13)
        r.font.color.rgb = NAVY if ligne.startswith("Q") else DGRAY
        r.font.bold = ligne.startswith("Q")
        r.font.name = "Calibri"


def ppt_grille_situation():
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    header(slide, "Grille d'Evaluation — Mise en Situation", label="PARTIE 3 — MISE EN SITUATION")
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=PURPLE)

    tbl = slide.shapes.add_table(
        len(GRILLE_SITUATION) + 1, 4,
        Inches(0.35), Inches(1.48), Inches(12.6), Inches(5.85)).table

    hdrs = ["Critere d'Evaluation", "Descripteurs", "Bareme", "Note Obtenue"]
    col_ws = [Inches(3.5), Inches(6.0), Inches(1.5), Inches(1.6)]
    for j, w in enumerate(col_ws):
        tbl.columns[j].width = w

    for j, h in enumerate(hdrs):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Calibri"

    for i, (crit, desc, bar) in enumerate(GRILLE_SITUATION):
        is_total = (i == len(GRILLE_SITUATION) - 1)
        for j, val in enumerate([crit, desc, bar, ""]):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_total else (LGRAY if i % 2 == 0 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j in [2, 3] else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val; r.font.size = Pt(12)
            r.font.color.rgb = ORANGE if is_total else DGRAY
            r.font.bold = is_total; r.font.name = "Calibri"


def ppt_bareme():
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, WHITE)
    header(slide, "Bareme de Notation — Recapitulatif General", label="SYNTHESE — BAREME")
    add_rect(slide, 0, Inches(1.32), Inches(0.18), H - Inches(1.32), fill=ORANGE)

    tbl = slide.shapes.add_table(
        len(BAREME) + 1, 4,
        Inches(1.0), Inches(1.55), Inches(11.33), Inches(4.5)).table

    hdrs = ["Partie", "Detail du Bareme", "Points Maximum", "Seuil de Validation"]
    col_ws = [Inches(3.2), Inches(4.0), Inches(2.0), Inches(2.13)]
    for j, w in enumerate(col_ws):
        tbl.columns[j].width = w

    for j, h in enumerate(hdrs):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Calibri"

    for i, row in enumerate(BAREME):
        is_total = (i == len(BAREME) - 1)
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if is_total else (LGRAY if i % 2 == 0 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val; r.font.size = Pt(13)
            r.font.color.rgb = ORANGE if is_total else DGRAY
            r.font.bold = is_total; r.font.name = "Calibri"

    # Note minimale
    txb(slide, "Note minimale de validation : 49 / 70  (70%)",
        Inches(1.0), Inches(6.2), Inches(11.33), Inches(0.55),
        sz=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.0), Inches(6.78), Inches(11.33), Inches(0.06), fill=ORANGE)
    txb(slide, "Formateur : TALOM Eric  |  Avril 2026",
        Inches(1.0), Inches(6.85), Inches(11.33), Inches(0.4),
        sz=12, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ── CONSTRUCTION PPT ─────────────────────────────────────────────────────────
ppt_cover()
ppt_instructions()

section_div(prs, "PARTIE 1\nQCM THEORIQUE", "30 questions — Connaissances & Reglementation CIMA", color=NAVY)

# QCM par tranches de 5
groups = [
    (QCM[0:5],   "Questions 1 a 5  — Connaissances Generales"),
    (QCM[5:10],  "Questions 6 a 10 — Connaissances Generales"),
    (QCM[10:15], "Questions 11 a 15 — Connaissances Generales"),
    (QCM[15:20], "Questions 16 a 20 — Connaissances Generales"),
    (QCM[20:25], "Questions 21 a 25 — Reglementation CIMA"),
    (QCM[25:30], "Questions 26 a 30 — Reglementation CIMA"),
]
for grp, lbl in groups:
    ppt_qcm_slide(grp, lbl)

ppt_corrige_qcm()

section_div(prs, "PARTIE 2\nCAS PRATIQUES", "3 dossiers — Automobile | Incendie | Sante", color=BLUE)

for cas in CAS_PRATIQUES:
    ppt_cas_pratique(cas)
    ppt_corrige_cas(cas)

section_div(prs, "PARTIE 3\nMISE EN SITUATION", "Grille comportementale — Simulation", color=PURPLE)
ppt_grille_situation()
ppt_bareme()

prs.save("Evaluation_Guide_Formateur_TALOM_Eric.pptx")
print(f"PPT OK : {len(prs.slides)} diapositives generees")


# ══════════════════════════════════════════════════════════════════════════════
# CONSTRUCTION WORD — CAHIER PARTICIPANT
# ══════════════════════════════════════════════════════════════════════════════

doc = Document()

# Marges
for section in doc.sections:
    section.top_margin    = Cm(2.0)
    section.bottom_margin = Cm(2.0)
    section.left_margin   = Cm(2.5)
    section.right_margin  = Cm(2.5)

# ── Helpers DOCX ─────────────────────────────────────────────────────────────

def set_cell_bg(cell, hex_color):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement("w:shd")
    shd.set(qn("w:val"),   "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"),  hex_color)
    tcPr.append(shd)

def set_cell_border(cell, border="bottom", color="1B3A5C", sz=6):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = tcPr.first_child_found_in("w:tcBorders")
    if tcBorders is None:
        tcBorders = OxmlElement("w:tcBorders")
        tcPr.append(tcBorders)
    bd = OxmlElement(f"w:{border}")
    bd.set(qn("w:val"),   "single")
    bd.set(qn("w:sz"),    str(sz))
    bd.set(qn("w:space"), "0")
    bd.set(qn("w:color"), color)
    tcBorders.append(bd)

def d_heading(doc, text, level=1, color="1B3A5C"):
    p   = doc.add_heading(text, level=level)
    run = p.runs[0] if p.runs else p.add_run(text)
    run.font.color.rgb = DRGBColor(
        int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))
    return p

def d_para(doc, text, sz=11, bold=False, italic=False, color="444444", space_before=6):
    p   = doc.add_paragraph()
    p.paragraph_format.space_before = DPt(space_before)
    p.paragraph_format.space_after  = DPt(2)
    run = p.add_run(text)
    run.font.size  = DPt(sz)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = DRGBColor(
        int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))
    return p

def d_rule(doc):
    p   = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pb  = OxmlElement("w:pBdr")
    bot = OxmlElement("w:bottom")
    bot.set(qn("w:val"),   "single")
    bot.set(qn("w:sz"),    "6")
    bot.set(qn("w:space"), "1")
    bot.set(qn("w:color"), "F1900l".replace("l", "1"))
    pb.append(bot)
    pPr.append(pb)
    return p

def d_answer_box(doc, lines=3, label="Reponse :"):
    d_para(doc, label, sz=10, bold=True, color="2E86AB")
    for _ in range(lines):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = DPt(2)
        p.paragraph_format.space_after  = DPt(0)
        run = p.add_run("_" * 110)
        run.font.size = DPt(9)
        run.font.color.rgb = DRGBColor(0xCC, 0xCC, 0xCC)

def d_checkbox(doc, text):
    p   = doc.add_paragraph()
    p.paragraph_format.left_indent   = Cm(0.8)
    p.paragraph_format.space_before  = DPt(2)
    p.paragraph_format.space_after   = DPt(2)
    run = p.add_run("☐  " + text)
    run.font.size = DPt(11)
    run.font.color.rgb = DRGBColor(0x44, 0x44, 0x44)


# ── PAGE DE GARDE ─────────────────────────────────────────────────────────────
d_heading(doc, "CAHIER D'EVALUATION", level=1, color="1B3A5C")
d_heading(doc, "Formation : Gestion des Sinistres — Toutes Branches", level=2, color="2E86AB")
d_para(doc, "Formateur : TALOM Eric  |  Avril 2026", sz=12, bold=True, color="F19001")
doc.add_paragraph()

# Cadre identité stagiaire
tbl_id = doc.add_table(rows=4, cols=2)
tbl_id.style = "Table Grid"
tbl_id.alignment = WD_TABLE_ALIGNMENT.CENTER
infos = [
    ("Nom & Prenom :", ""),
    ("Entreprise / Structure :", ""),
    ("Poste occupe :", ""),
    ("Date de la formation :", "Avril 2026"),
]
for i, (label, val) in enumerate(infos):
    c0, c1 = tbl_id.cell(i, 0), tbl_id.cell(i, 1)
    set_cell_bg(c0, "1B3A5C")
    r0 = c0.paragraphs[0].add_run(label)
    r0.font.bold = True; r0.font.size = DPt(11)
    r0.font.color.rgb = DRGBColor(0xFF, 0xFF, 0xFF)
    r1 = c1.paragraphs[0].add_run(val)
    r1.font.size = DPt(11)
doc.add_paragraph()

d_para(doc,
    "Instructions : Lisez attentivement chaque question avant de repondre. "
    "Justifiez vos reponses aux cas pratiques. Duree totale : 2h15. "
    "Note minimale de validation : 49/70 (70%).",
    sz=10, italic=True, color="666666")
doc.add_page_break()

# ── PARTIE 1 — QCM ────────────────────────────────────────────────────────────
d_heading(doc, "PARTIE 1 — QCM THEORIQUE  (30 questions / 15 points)", level=1, color="1B3A5C")
d_para(doc, "Une seule bonne reponse par question. Cochez la case correspondante.",
       sz=10, italic=True, color="888888")
doc.add_paragraph()

current_cat = ""
for q in QCM:
    if q["cat"] != current_cat:
        current_cat = q["cat"]
        d_heading(doc, current_cat, level=2, color="2E86AB")

    d_para(doc, f"Question {q['n']}  (0,5 pt)   {q['q']}", sz=11, bold=True, color="1B3A5C", space_before=10)
    for j, choice in enumerate(q["a"]):
        d_checkbox(doc, f"{chr(65+j)}.  {choice}")
    doc.add_paragraph()

doc.add_page_break()

# ── PARTIE 2 — CAS PRATIQUES ──────────────────────────────────────────────────
d_heading(doc, "PARTIE 2 — CAS PRATIQUES  (3 dossiers / 30 points)", level=1, color="1B3A5C")
d_para(doc, "Traitez chaque dossier de maniere independante. Justifiez vos calculs.",
       sz=10, italic=True, color="888888")

for cas in CAS_PRATIQUES:
    doc.add_paragraph()
    d_heading(doc, f"Cas Pratique {cas['num']} — Branche {cas['branche']}  ({cas['duree']} / 10 pts)",
              level=2, color="1B3A5C")
    d_para(doc, f"Intitule : {cas['titre']}", sz=11, bold=True, color="2E86AB")
    doc.add_paragraph()

    # Contexte encadre
    tbl_ctx = doc.add_table(rows=1, cols=1)
    tbl_ctx.style = "Table Grid"
    cell = tbl_ctx.cell(0, 0)
    set_cell_bg(cell, "F0F4F8")
    r = cell.paragraphs[0].add_run("CONTEXTE  —  " + cas["contexte"])
    r.font.size = DPt(10)
    r.font.color.rgb = DRGBColor(0x44, 0x44, 0x44)
    doc.add_paragraph()

    for question in cas["questions"]:
        d_para(doc, question, sz=11, bold=True, color="1B3A5C", space_before=8)
        d_answer_box(doc, lines=4)
        doc.add_paragraph()

    doc.add_page_break()

# ── PARTIE 3 — MISE EN SITUATION ──────────────────────────────────────────────
d_heading(doc, "PARTIE 3 — MISE EN SITUATION  (20 points)", level=1, color="1B3A5C")
d_para(doc, "Simulation en binomes — roles : gestionnaire sinistres / assure.",
       sz=10, italic=True, color="888888")
doc.add_paragraph()

d_para(doc, "Scenarios proposes (tires au sort) :", sz=11, bold=True, color="1B3A5C")
scenarios = [
    "Scenario A : Reception telephonique d'une declaration de sinistre automobile",
    "Scenario B : Annonce a un assure d'un refus de garantie (exclusion contractuelle)",
    "Scenario C : Remise d'un courrier de reglement partiel (regle proportionnelle incendie)",
]
for sc in scenarios:
    d_para(doc, "▸  " + sc, sz=11, color="444444")

doc.add_paragraph()
d_para(doc, "Auto-evaluation du participant :", sz=11, bold=True, color="2E86AB")

tbl_g = doc.add_table(rows=len(GRILLE_SITUATION) + 1, cols=4)
tbl_g.style = "Table Grid"
tbl_g.alignment = WD_TABLE_ALIGNMENT.CENTER
hdrs_g = ["Critere", "Descripteurs", "Bareme", "Auto-note"]
for j, h in enumerate(hdrs_g):
    cell = tbl_g.cell(0, j)
    set_cell_bg(cell, "1B3A5C")
    r = cell.paragraphs[0].add_run(h)
    r.font.bold = True; r.font.size = DPt(10)
    r.font.color.rgb = DRGBColor(0xFF, 0xFF, 0xFF)
for i, (crit, desc, bar) in enumerate(GRILLE_SITUATION):
    for j, val in enumerate([crit, desc, bar, ""]):
        cell = tbl_g.cell(i + 1, j)
        if i % 2 == 0:
            set_cell_bg(cell, "F0F4F8")
        r = cell.paragraphs[0].add_run(val)
        r.font.size = DPt(10)
        if i == len(GRILLE_SITUATION) - 1:
            r.font.bold = True
            set_cell_bg(cell, "1B3A5C")
            r.font.color.rgb = DRGBColor(0xF1, 0x90, 0x01)

doc.add_page_break()

# ── FICHE DE RESULTATS ────────────────────────────────────────────────────────
d_heading(doc, "FICHE DE RESULTATS & ATTESTATION", level=1, color="1B3A5C")
doc.add_paragraph()

tbl_r = doc.add_table(rows=len(BAREME) + 1, cols=4)
tbl_r.style = "Table Grid"
tbl_r.alignment = WD_TABLE_ALIGNMENT.CENTER
hdrs_r = ["Partie", "Points Maximum", "Seuil (70%)", "Note Obtenue"]
for j, h in enumerate(hdrs_r):
    cell = tbl_r.cell(0, j)
    set_cell_bg(cell, "1B3A5C")
    r = cell.paragraphs[0].add_run(h)
    r.font.bold = True; r.font.size = DPt(11)
    r.font.color.rgb = DRGBColor(0xFF, 0xFF, 0xFF)
for i, row in enumerate(BAREME):
    for j, val in enumerate([row[0], row[2], row[3], ""]):
        cell = tbl_r.cell(i + 1, j)
        if i % 2 == 0:
            set_cell_bg(cell, "F0F4F8")
        if i == len(BAREME) - 1:
            set_cell_bg(cell, "1B3A5C")
        r = cell.paragraphs[0].add_run(val)
        r.font.size = DPt(11)
        r.font.bold = (i == len(BAREME) - 1)
        if i == len(BAREME) - 1:
            r.font.color.rgb = DRGBColor(0xF1, 0x90, 0x01)

doc.add_paragraph()
d_para(doc, "Apreciation generale :", sz=11, bold=True, color="1B3A5C")
d_answer_box(doc, lines=3, label="Commentaires du formateur :")
doc.add_paragraph()

# Signatures
tbl_s = doc.add_table(rows=2, cols=2)
tbl_s.style = "Table Grid"
labels_s = [
    ("Signature du Participant",    "Signature du Formateur"),
    ("",                            "TALOM Eric"),
]
for i, row in enumerate(labels_s):
    for j, val in enumerate(row):
        cell = tbl_s.cell(i, j)
        if i == 0:
            set_cell_bg(cell, "1B3A5C")
        r = cell.paragraphs[0].add_run(val)
        r.font.size  = DPt(11)
        r.font.bold  = True
        if i == 0:
            r.font.color.rgb = DRGBColor(0xFF, 0xFF, 0xFF)
        r.font.color.rgb = DRGBColor(0x1B, 0x3A, 0x5C) if i == 1 else r.font.color.rgb

doc.add_paragraph()
d_para(doc, "Ce document est confidentiel — Formation TALOM Eric — Avril 2026",
       sz=9, italic=True, color="AAAAAA")

doc.save("Evaluation_Cahier_Participant.docx")
print(f"DOCX OK : Evaluation_Cahier_Participant.docx")
print(f"Tous les documents ont ete generes avec succes.")
