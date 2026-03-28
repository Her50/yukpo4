# -*- coding: utf-8 -*-
"""Présentation professionnelle : IA au travail (PowerPoint).

Intro 2 slides (IA digestible) ; section ; tableaux 6 colonnes par métier
(situation simple, apport concret, où / comment en production, outil, lien, accès).
Fond discret, ligne d'accent, tableau rétracté pour laisser une zone logo sans chevauchement.
Bloc Yukpo en fin. Numérotation « n / total » sur chaque slide. Génère IA_milieu_professionnel_court.pptx à la racine du dépôt.
"""
from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

_SCRIPTS_DIR = Path(__file__).resolve().parent
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from ia_presentation_data import (
    claims_broker,
    civil_engineer,
    data_scientist,
    doctor_gp,
    industry_plastic,
    lawyer,
    legal_insurance,
    marketing_telco,
    med_rep,
    reinsurance,
    surgeon,
)

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

_LOGO_PATH: Path | None = None

# Mise en page tableaux
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.45)
MARGIN_R = Inches(0.38)
TITLE_TOP = Inches(0.26)
TITLE_H = Inches(0.88)
TABLE_LEFT = Inches(0.48)

# Logo bas droite : zone réservée (évite de recouvrir la dernière colonne du tableau)
LOGO_H = Inches(0.42)
LOGO_MARGIN = Inches(0.2)
# Largeur utile = slide − marge gauche − (logo + marge droite)
TABLE_W = Inches(10) - TABLE_LEFT - MARGIN_R - LOGO_H - LOGO_MARGIN
TITLE_W = TABLE_W  # aligné sur le tableau — jamais sous le logo

# Sous-titre = bloc plus haut : la table commence plus bas
TABLE_TOP_NO_SUB = Inches(1.28)
TABLE_TOP_WITH_SUB = Inches(1.66)

# Au-delà, scinder (6 colonnes = texte dense)
MAX_DATA_ROWS_PER_TABLE_SLIDE = 3


def resolve_logo_path() -> Path | None:
    root = Path(__file__).resolve().parent.parent
    candidate = root / "mobile" / "assets" / "icon.png"
    return candidate if candidate.is_file() else None


def apply_slide_background(slide) -> None:
    """Fond très léger — rendu type présentation corporate (sans surcharge)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFB)


def add_slide_top_accent_bar(slide) -> None:
    """Fine barre supérieure — repère visuel type template pro (sous le contenu)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0.0,
        0.0,
        float(SLIDE_W),
        float(Inches(0.055)),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    bar.line.fill.background()


def add_logo_bottom_right(slide) -> None:
    path = _LOGO_PATH
    if path is None or not path.is_file():
        return
    left = float(SLIDE_W) - float(LOGO_H) - float(MARGIN_R)
    top = float(SLIDE_H) - float(LOGO_H) - float(MARGIN_R)
    slide.shapes.add_picture(str(path), left, top, height=LOGO_H)


def add_title_accent_line(slide, y_inches: float) -> None:
    """Ligne d'accent sous la zone titre — rendu type présentation corporate."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        float(TABLE_LEFT),
        float(Inches(y_inches)),
        float(TABLE_W),
        float(Inches(0.035)),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x2E, 0x6B, 0xC4)
    line.line.fill.background()


def add_slide_footer(slide, left_text: str = "IA & métiers — présentation professionnelle") -> None:
    """Pied de page discret — au-dessus de la zone numérotation (bas gauche)."""
    box = slide.shapes.add_textbox(
        float(MARGIN_L),
        float(SLIDE_H) - float(Inches(0.72)),
        float(Inches(6.8)),
        float(Inches(0.28)),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = left_text
    p.font.size = Pt(8.5)
    p.font.color.rgb = RGBColor(0x77, 0x77, 0x77)


def stamp_slide_number(slide, num: int, total: int) -> None:
    """Numérotation bas gauche — évite tout chevauchement avec le logo (bas droite)."""
    box = slide.shapes.add_textbox(
        float(MARGIN_L),
        float(SLIDE_H) - float(Inches(0.38)),
        float(Inches(1.35)),
        float(Inches(0.28)),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p.alignment = PP_ALIGN.LEFT


def apply_numbering_to_all_slides(prs: Presentation) -> None:
    """Numérote toutes les diapositives après construction (évite les doublons)."""
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        stamp_slide_number(slide, i + 1, total)


def _title_textbox(slide, title: str, subtitle: str | None = None) -> None:
    """Titre dans la zone sûre : même largeur que le tableau (aucune zone sous le logo)."""
    box = slide.shapes.add_textbox(TABLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H if not subtitle else Inches(1.18))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p0 = tf.paragraphs[0]
    p0.text = title
    title_sz = Pt(20) if len(title) > 54 else Pt(22)
    p0.font.size = title_sz
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    if subtitle:
        p1 = tf.add_paragraph()
        p1.text = subtitle
        p1.font.size = Pt(12.5)
        p1.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        p1.space_before = Pt(4)


def _style_table_header(table) -> None:
    row0 = table.rows[0]
    for cell in row0.cells:
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(10.5)
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)


def _style_table_alternate_rows(table) -> None:
    """Lignes paires légèrement grisées — lisibilité type présentation corporate."""
    for i in range(1, len(table.rows)):
        if i % 2 == 0:
            for cell in table.rows[i].cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF9)


def _set_cell_text(cell, text: str, size_pt: float = 10.0, bold: bool = False) -> None:
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(0)
    cell.vertical_anchor = MSO_ANCHOR.TOP


def _set_cell_hyperlink(cell, text: str, url: str, size_pt: float = 10.0) -> None:
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = RGBColor(0x0B, 0x57, 0xA6)
    run.hyperlink.address = url


@dataclass(frozen=True)
class TableRow:
    problem: str
    ai_value: str
    deployment: str
    tool_name: str
    tool_url: str
    access: str


def _title_bottom_inches(subtitle: bool) -> float:
    """Position basse (pouces) du bloc titre — pour placer la ligne d'accent sans chevauchement."""
    top = 0.26
    if subtitle:
        return top + 1.18
    return top + 0.88


def add_table_profession_slides(
    prs: Presentation,
    title: str,
    subtitle: str | None,
    rows: list[TableRow],
) -> None:
    """Découpe en plusieurs diapositives si le tableau dépasse la hauteur lisible."""
    if not rows:
        return
    chunks = [
        rows[i : i + MAX_DATA_ROWS_PER_TABLE_SLIDE]
        for i in range(0, len(rows), MAX_DATA_ROWS_PER_TABLE_SLIDE)
    ]
    n = len(chunks)
    for i, chunk in enumerate(chunks):
        suff = f" — partie {i + 1}/{n}" if n > 1 else ""
        if subtitle and n > 1:
            sub = f"{subtitle} — tableau {i + 1}/{n}"
        else:
            sub = subtitle
        add_table_profession_slide(prs, f"{title}{suff}", sub, chunk)


def _set_cell_bold_label(cell, label: str, body: str, label_pt: float = 9.5, body_pt: float = 9.5) -> None:
    """Label en gras + corps — lisibilité des colonnes « apport » / « installation »."""
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r0 = p.add_run()
    r0.text = label
    r0.font.bold = True
    r0.font.size = Pt(label_pt)
    r0.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    r1 = p.add_run()
    r1.text = body
    r1.font.size = Pt(body_pt)
    r1.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    cell.vertical_anchor = MSO_ANCHOR.TOP


def add_table_profession_slide(
    prs: Presentation,
    title: str,
    subtitle: str | None,
    rows: list[TableRow],
) -> None:
    """Tableau 6 colonnes : situation, apport, déploiement, outil, lien, accès."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    apply_slide_background(slide)
    add_slide_top_accent_bar(slide)
    _title_textbox(slide, title, subtitle)

    accent_y = _title_bottom_inches(bool(subtitle)) + 0.04
    add_title_accent_line(slide, accent_y)

    table_top = TABLE_TOP_WITH_SUB if subtitle else TABLE_TOP_NO_SUB

    nrows = len(rows) + 1
    ncols = 6
    row_h = Inches(0.88)
    footer_reserve = float(LOGO_H) + float(Inches(0.52))
    max_table_bottom = float(SLIDE_H) - footer_reserve
    table_h = min(max_table_bottom - float(table_top), float(row_h * nrows))

    shape = slide.shapes.add_table(nrows, ncols, TABLE_LEFT, table_top, TABLE_W, table_h)
    table = shape.table

    headers = (
        "Situation (concret)",
        "Ce que change l'IA",
        "Où ça s'installe / comment ça tourne",
        "Outil",
        "Lien",
        "Accès",
    )
    widths = (
        Inches(1.22),
        Inches(1.78),
        Inches(1.5),
        Inches(1.18),
        Inches(1.32),
        Inches(1.52),
    )
    for j, w in enumerate(widths):
        table.columns[j].width = w

    hdr_sz = 9.0
    for j, h in enumerate(headers):
        _set_cell_text(table.cell(0, j), h, hdr_sz, True)
    _style_table_header(table)

    for i, tr in enumerate(rows, start=1):
        table.rows[i].height = row_h
        _set_cell_text(table.cell(i, 0), tr.problem, 10.0)
        _set_cell_bold_label(
            table.cell(i, 1),
            "Effet : ",
            tr.ai_value,
            9.5,
            10.0,
        )
        _set_cell_bold_label(
            table.cell(i, 2),
            "Installation : ",
            tr.deployment,
            9.0,
            9.5,
        )
        _set_cell_hyperlink(table.cell(i, 3), tr.tool_name, tr.tool_url, 9.5)
        short_link = tr.tool_url.replace("https://", "").replace("http://", "")
        if len(short_link) > 32:
            short_link = short_link[:29] + "…"
        _set_cell_hyperlink(table.cell(i, 4), short_link, tr.tool_url, 8.5)
        _set_cell_text(table.cell(i, 5), tr.access, 9.5)

    _style_table_alternate_rows(table)
    add_slide_footer(slide)
    add_logo_bottom_right(slide)


def add_bullet_slide_professional(
    prs: Presentation,
    title: str,
    bullets: list[str],
    subtitle: str | None = None,
    font_pt: float = 15.5,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    add_slide_top_accent_bar(slide)
    _title_textbox(slide, title, subtitle)
    has_sub = bool(subtitle)
    accent_y = _title_bottom_inches(has_sub) + 0.04
    add_title_accent_line(slide, accent_y)
    body_top_in = accent_y + 0.12
    box_h_in = max(2.0, 7.5 - body_top_in - 0.58)
    box = slide.shapes.add_textbox(TABLE_LEFT, Inches(body_top_in), TABLE_W, Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    for k, line in enumerate(bullets):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_pt)
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(10)
        p.level = 0
    add_slide_footer(slide)
    add_logo_bottom_right(slide)


def add_section_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Slide de rupture visuelle avant les fiches sectorielles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0.0,
        float(Inches(2.15)),
        float(SLIDE_W),
        float(Inches(1.55)),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    band.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.75), Inches(2.35), Inches(8.5), Inches(1.35))
    tf = box.text_frame
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(26)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p0.alignment = PP_ALIGN.CENTER
    p1 = tf.add_paragraph()
    p1.text = subtitle
    p1.font.size = Pt(14)
    p1.font.color.rgb = RGBColor(0xDD, 0xE5, 0xF5)
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(10)
    add_slide_footer(slide, "Partie 2 — Fiches sectorielles (IA)")
    add_logo_bottom_right(slide)


def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    author_block: str | None = None,
) -> None:
    """Page de garde / clôture — zone auteur optionnelle sous le sous-titre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    add_slide_top_accent_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.85), Inches(8.5), Inches(3.35))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(30)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    p0.alignment = PP_ALIGN.CENTER
    for i, line in enumerate(subtitle.split("\n")):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13.5)
        p.font.color.rgb = RGBColor(0x45, 0x45, 0x45)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(16) if i == 0 else Pt(6)
    if author_block:
        for j, line in enumerate(author_block.split("\n")):
            pa = tf.add_paragraph()
            pa.text = line
            pa.font.size = Pt(11.5) if j == 0 else Pt(10.5)
            pa.font.bold = j == 0
            pa.font.color.rgb = RGBColor(0x2E, 0x4E, 0x72)
            pa.alignment = PP_ALIGN.CENTER
            pa.space_before = Pt(22) if j == 0 else Pt(5)
    add_slide_footer(slide, "IA & métiers — présentation professionnelle")
    add_logo_bottom_right(slide)


def add_yukpo_hero_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    add_slide_top_accent_bar(slide)
    path = _LOGO_PATH
    if path and path.is_file():
        slide.shapes.add_picture(str(path), Inches(3.55), Inches(0.95), height=Inches(2.75))
    box = slide.shapes.add_textbox(Inches(0.75), Inches(3.95), Inches(8.5), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "Yukpo"
    p0.font.size = Pt(40)
    p0.font.bold = True
    p0.alignment = PP_ALIGN.CENTER
    p1 = tf.add_paragraph()
    p1.text = "L'écoute qui comprend vraiment"
    p1.font.size = Pt(20)
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Plateforme intelligente — besoins et solutions"
    p2.font.size = Pt(15)
    p2.alignment = PP_ALIGN.CENTER
    add_slide_footer(slide, "Yukpo — écosystème Yukpomnang")


def add_bullet_slide_bold_leads(
    prs: Presentation,
    title: str,
    items: list[tuple[str, str]],
    subtitle: str | None = None,
    lead_pt: float = 15.5,
    body_pt: float = 14.0,
) -> None:
    """Puces avec début de phrase en gras."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    add_slide_top_accent_bar(slide)
    _title_textbox(slide, title, subtitle)
    has_sub = bool(subtitle)
    accent_y = _title_bottom_inches(has_sub) + 0.04
    add_title_accent_line(slide, accent_y)
    body_top_in = accent_y + 0.12
    box_h_in = max(2.0, 7.5 - body_top_in - 0.58)
    box = slide.shapes.add_textbox(TABLE_LEFT, Inches(body_top_in), TABLE_W, Inches(box_h_in))
    tf = box.text_frame
    tf.word_wrap = True
    for k, (lead, body) in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        r0 = p.add_run()
        r0.text = lead
        r0.font.bold = True
        r0.font.size = Pt(lead_pt)
        r0.font.color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
        r1 = p.add_run()
        r1.text = " " + body
        r1.font.size = Pt(body_pt)
        r1.font.color.rgb = RGBColor(0x28, 0x28, 0x28)
        p.space_after = Pt(12)
    add_slide_footer(slide)
    add_logo_bottom_right(slide)


def add_tools_lines_slide(
    prs: Presentation,
    title: str,
    lines: list[tuple[str, str, str, str]],
) -> None:
    """Lignes nom (lien), accès, détail — pour Yukpo fin."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_background(slide)
    add_slide_top_accent_bar(slide)
    _title_textbox(slide, title, None)
    accent_y = _title_bottom_inches(False) + 0.04
    add_title_accent_line(slide, accent_y)
    box = slide.shapes.add_textbox(TABLE_LEFT, Inches(accent_y + 0.12), TABLE_W, Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, (name, url, access, detail) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = name
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.hyperlink.address = url
        r2 = p.add_run()
        r2.text = f"  —  {access}  —  {detail}"
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(14)
    add_slide_footer(slide)
    add_logo_bottom_right(slide)


def main() -> None:
    global _LOGO_PATH
    _LOGO_PATH = resolve_logo_path()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ----- Page de garde -----
    add_title_slide(
        prs,
        "Intelligence artificielle & métiers",
        "Apports concrets par secteur — outils et mise en œuvre\n"
        "Mars 2026",
        author_block="Hernandez LELE",
    )

    # ----- Rappel IA (langage accessible) -----
    add_bullet_slide_professional(
        prs,
        "Qu'est-ce que l'intelligence artificielle ?",
        [
            "En une phrase : des programmes qui apprennent à partir d'exemples (textes, images, chiffres) pour repérer des motifs et proposer des réponses — ce n'est pas une conscience : du calcul statistique à très grande échelle.",
            "Ce n'est pas infaillible : ces systèmes peuvent se tromper ou produire un texte « crédible » mais faux — relire et valider reste indispensable.",
            "Différence utile : une automatisation classique suit des règles fixes ; l'IA s'ajuste quand on lui montre de nouveaux cas (images, texte, prévisions, anomalies…).",
            "Au travail : l'IA réduit le temps sur le répétitif et le volumineux ; la décision qui engage l'entreprise, le patient ou le client reste humaine et responsable.",
        ],
        subtitle="Pour tout public",
        font_pt=15.0,
    )

    add_bullet_slide_professional(
        prs,
        "Les grands usages de l'IA au travail",
        [
            "Voir : repérer un défaut, une anomalie ou un objet sur une image ou une vidéo.",
            "Lire et écrire : résumer, reformuler, extraire des informations, produire un brouillon — toujours relire.",
            "Anticiper : prévoir une panne, un volume ou un risque ; signaler ce qui sort de l'ordinaire.",
            "Prioriser : trier dossiers, comptes ou actions à partir de l'historique.",
            "Dans les tableaux : situation du quotidien → effet utile → où l'outil se branche techniquement → nom du service → lien officiel → type d'accès.",
            "Données personnelles : ne les saisir dans un service en ligne qu'avec cadre adapté (RGPD, contrat, hébergement santé HDS le cas échéant).",
        ],
        subtitle="Lecture des fiches",
        font_pt=13.5,
    )

    add_section_slide(
        prs,
        "Fiches par secteur",
        "Situation → effet utile → installation / production → outil → lien",
    )

    # ----- Métiers (données : scripts/ia_presentation_data.py) -----
    add_table_profession_slides(
        prs,
        "Industriel — plastique & extrusion",
        "Contrôle visuel, capteurs, prévision — ligne de production",
        [TableRow(*t) for t in industry_plastic()],
    )

    add_table_profession_slides(
        prs,
        "Médecin généraliste",
        "Dossier, imagerie, aide au raisonnement — décision médicale toujours humaine",
        [TableRow(*t) for t in doctor_gp()],
    )

    add_table_profession_slides(
        prs,
        "Chirurgien & imagerie",
        "Planification 3D, urgences, bloc, file d'attente imagerie",
        [TableRow(*t) for t in surgeon()],
    )

    add_table_profession_slides(
        prs,
        "Responsable juridique — assurance",
        "Contrats et veille",
        [TableRow(*t) for t in legal_insurance()],
    )

    add_table_profession_slides(
        prs,
        "Avocat",
        "Recherche et rédaction",
        [TableRow(*t) for t in lawyer()],
    )

    add_table_profession_slides(
        prs,
        "Délégué médical",
        "Littérature scientifique et CRM",
        [TableRow(*t) for t in med_rep()],
    )

    add_table_profession_slides(
        prs,
        "Responsable marketing — télécom",
        "Campagnes et CRM",
        [TableRow(*t) for t in marketing_telco()],
    )

    add_table_profession_slides(
        prs,
        "Ingénieur génie civil & construction",
        "Planning, chantier, BIM",
        [TableRow(*t) for t in civil_engineer()],
    )

    add_table_profession_slides(
        prs,
        "Responsable sinistres & courtage — assurance",
        "Priorisation, fraude, pièces, chiffrage",
        [TableRow(*t) for t in claims_broker()],
    )

    add_table_profession_slides(
        prs,
        "Responsable réassurance",
        "Tarification, modèles, comités",
        [TableRow(*t) for t in reinsurance()],
    )

    add_table_profession_slides(
        prs,
        "Statisticien / data scientist",
        "Modélisation et mise en production",
        [TableRow(*t) for t in data_scientist()],
    )

    # ----- Action -----
    add_bullet_slide_professional(
        prs,
        "Passer à l'action",
        [
            "Choisir une ligne pertinente dans votre secteur et contacter l'éditeur pour une démonstration ou un pilote encadré.",
            "Tester sur un petit cas interne avec données anonymisées ; mesurer le gain de temps sur deux semaines.",
            "Obtenir l'accord juridique, IT et qualité avant toute mise en production sur données réelles.",
        ],
        subtitle="Prochaines étapes",
        font_pt=15.0,
    )

    # ----- Yukpo -----
    add_bullet_slide_bold_leads(
        prs,
        "Yukpo — positionnement",
        [
            (
                "Une plateforme O2O :",
                "elle relie ce qui se passe en ligne (recherche, offre, commande) à la réalité terrain (livraison, service, contact local).",
            ),
            (
                "L'intelligence au centre :",
                "recherche multimodale (texte, voix, image), assistants pour créer ou compléter une offre, parcours guidés pour aller plus vite du besoin à la solution.",
            ),
            (
                "Pour les professionnels et les communautés :",
                "visibilité numérique des commerces et services locaux, mise en relation directe, expérience unifiée mobile et web sous la marque Yukpomnang.",
            ),
        ],
        subtitle="Yukpo Company — écosystème Yukpomnang",
        lead_pt=16.0,
        body_pt=14.5,
    )

    add_bullet_slide_bold_leads(
        prs,
        "Ce que les utilisateurs et les pros y font",
        [
            (
                "Trouver :",
                "recherche intelligente combinant mots-clés, oral, photo et géolocalisation pour découvrir produits et prestataires à proximité.",
            ),
            (
                "Publier :",
                "création accélérée de fiches produits ou services (texte, médias) avec aide à la mise en forme et à la complétude.",
            ),
            (
                "Échanger :",
                "messagerie temps réel, suivi de commandes et de livraisons, parcours d'achat ou de demande de service de bout en bout.",
            ),
            (
                "Fidéliser :",
                "visibilité renforcée pour les vendeurs et créateurs, logique de place de marché avec livraison intégrée lorsque le contexte le permet.",
            ),
        ],
        subtitle="Fonctionnalités majeures",
        lead_pt=16.0,
        body_pt=14.5,
    )

    add_bullet_slide_bold_leads(
        prs,
        "Confiance et déploiement",
        [
            (
                "Sécurité et conformité :",
                "traitement des données dans le respect des cadres applicables ; transparence sur les usages des fonctionnalités intelligentes.",
            ),
            (
                "Accès :",
                "application mobile Yukpo (Android / iOS selon les stores), site Yukpomnang pour l'information et le contact.",
            ),
            (
                "Ambition :",
                "réduire la fracture entre besoins exprimés (y compris en langues locales) et solutions disponibles sur le marché réel.",
            ),
        ],
        subtitle="Cadre et canaux",
        lead_pt=16.0,
        body_pt=14.5,
    )

    add_yukpo_hero_slide(prs)

    add_tools_lines_slide(
        prs,
        "Yukpo — liens officiels",
        [
            (
                "Yukpomnang.com",
                "https://www.yukpomnang.com",
                "Information & contact",
                "Site institutionnel de la plateforme.",
            ),
            (
                "Google Play — Yukpo",
                "https://play.google.com/store/apps/details?id=com.yukpomnang.mobile",
                "Téléchargement Android",
                "Vérifier l'éditeur Yukpomnang sur la fiche store.",
            ),
            (
                "App Store — recherche « Yukpo »",
                "https://apps.apple.com/search?term=Yukpo",
                "iOS",
                "Confirmer l'éditeur avant installation.",
            ),
        ],
    )

    add_title_slide(
        prs,
        "Merci",
        "Questions et échanges\n"
        "yukpomnang.com",
        author_block="Hernandez LELE\n"
        "CEO — Yukpo Company\n"
        "Consultant en statistique · suivi-évaluation de projets de développement · conception et développement full-stack",
    )

    apply_numbering_to_all_slides(prs)

    out_dir = Path(__file__).resolve().parent.parent
    out_path = out_dir / "IA_milieu_professionnel_court.pptx"
    prs.save(out_path)
    print(f"Fichier créé : {out_path}")


if __name__ == "__main__":
    main()
