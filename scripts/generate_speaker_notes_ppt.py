from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def parse_notes(md_path: Path) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_title = ""
    bullets: list[str] = []
    for raw in md_path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("# "):
            continue
        if line.startswith("## "):
            if current_title:
                sections.append((current_title, bullets))
            current_title = line[3:].strip()
            bullets = []
            continue
        if line.startswith("- "):
            bullets.append(line[2:].strip())
    if current_title:
        sections.append((current_title, bullets))
    return sections


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.6), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "Yukpo - Speaker Notes"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x12, 0x2A, 0x4A)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.6), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "Version pitch oral (6-8 minutes)"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def add_section_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x12, 0x2A, 0x4A)

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(8.4), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        pb = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pb.text = f"- {b}"
        pb.font.size = Pt(19 if len(bullets) <= 4 else 17)
        pb.space_after = Pt(8)
        pb.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    src = root / "Yukpo_Speaker_Notes.md"
    out = root / "Yukpo_Speaker_Notes.pptx"

    sections = parse_notes(src)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs)
    for title, bullets in sections:
        add_section_slide(prs, title, bullets)
    prs.save(out)
    print(f"Fichier cree : {out}")


if __name__ == "__main__":
    main()

