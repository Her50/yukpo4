#!/usr/bin/env python3
"""
Script pour créer une présentation PowerPoint professionnelle du Pitch Deck
basée sur le document de financement finalisé
Optimisé pour présentations investisseurs - Design moderne et impactant
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installation de python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

# Couleurs professionnelles Yukpomnang
COLORS = {
    'primary': RGBColor(99, 102, 241),      # #6366F1 (Indigo)
    'secondary': RGBColor(102, 126, 234),   # #667EEA (Bleu-violet)
    'accent': RGBColor(118, 75, 162),       # #764BA2 (Violet)
    'success': RGBColor(16, 185, 129),      # #10B981 (Vert)
    'warning': RGBColor(255, 193, 7),       # #FFC107 (Jaune)
    'text': RGBColor(44, 62, 80),           # #2C3E50 (Gris foncé)
    'light': RGBColor(128, 128, 128),       # Gris clair
    'background': RGBColor(255, 255, 255),  # Blanc
    'highlight': RGBColor(227, 242, 253),   # Bleu clair
}

def add_gradient_background(slide, color1, color2):
    """Ajoute un fond dégradé professionnel"""
    try:
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135.0
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
    except:
        pass  # Si gradient non supporté, on continue

def create_title_slide(prs, title, subtitle="", tagline=""):
    """Crée une slide de titre professionnelle avec design moderne"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond dégradé subtil
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['background']
    except:
        pass
    
    # Titre principal
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(64)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.name = 'Calibri'
    
    # Sous-titre
    if subtitle:
        subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = COLORS['text']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.name = 'Calibri'
    
    # Tagline
    if tagline:
        tagline_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.6))
        tagline_frame = tagline_shape.text_frame
        tagline_frame.text = tagline
        tagline_frame.paragraphs[0].font.size = Pt(18)
        tagline_frame.paragraphs[0].font.color.rgb = COLORS['light']
        tagline_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tagline_frame.paragraphs[0].font.name = 'Calibri'
        tagline_frame.paragraphs[0].font.italic = True
    
    return slide

def create_content_slide(prs, title, content_list, highlight_first=False):
    """Crée une slide de contenu avec titre et liste à puces"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre avec ligne de séparation
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].font.name = 'Calibri'
    
    # Ligne de séparation colorée
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = COLORS['primary']
    line.line.width = Pt(4)
    
    # Contenu avec puces
    content_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.2))
    content_frame = content_shape.text_frame
    content_frame.word_wrap = True
    content_frame.margin_left = Inches(0.2)
    content_frame.margin_right = Inches(0.2)
    
    for i, point in enumerate(content_list):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.level = 0
        p.space_after = Pt(14)
        p.left_indent = Inches(0.3)
        
        # Traitement du texte avec formatage
        if '**' in point:
            parts = point.split('**')
            for j, part in enumerate(parts):
                if part.strip():
                    run = p.add_run()
                    run.text = part.strip() + (' ' if j < len(parts) - 1 else '')
                    run.font.size = Pt(20)
                    run.font.name = 'Calibri'
                    if j % 2 == 1:  # Texte entre **
                        run.font.bold = True
                        run.font.color.rgb = COLORS['primary']
                    else:
                        run.font.color.rgb = COLORS['text']
        else:
            run = p.add_run()
            run.text = point
            run.font.size = Pt(20)
            run.font.color.rgb = COLORS['text'] if not (highlight_first and i == 0) else COLORS['primary']
            run.font.name = 'Calibri'
            if highlight_first and i == 0:
                run.font.bold = True
    
    return slide

def create_stat_slide(prs, title, stats):
    """Crée une slide avec statistiques en format visuel"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].font.name = 'Calibri'
    
    # Stats en grille
    y_start = 1.5
    x_positions = [1.5, 5.5]
    width = 3.5
    height = 1.2
    
    for i, (label, value) in enumerate(stats):
        x = x_positions[i % 2]
        y = y_start + (i // 2) * 1.5
        
        # Box avec fond coloré
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     Inches(x), Inches(y), Inches(width), Inches(height))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['highlight']
        box.line.color.rgb = COLORS['primary']
        box.line.width = Pt(2)
        
        # Valeur
        value_shape = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), 
                                               Inches(width - 0.2), Inches(0.5))
        value_frame = value_shape.text_frame
        value_frame.text = value
        value_frame.paragraphs[0].font.size = Pt(28)
        value_frame.paragraphs[0].font.bold = True
        value_frame.paragraphs[0].font.color.rgb = COLORS['primary']
        value_frame.paragraphs[0].font.name = 'Calibri'
        
        # Label
        label_shape = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.6), 
                                               Inches(width - 0.2), Inches(0.5))
        label_frame = label_shape.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].font.size = Pt(16)
        label_frame.paragraphs[0].font.color.rgb = COLORS['text']
        label_frame.paragraphs[0].font.name = 'Calibri'
    
    return slide

def create_table_slide(prs, title, headers, rows):
    """Crée une slide avec tableau"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].font.name = 'Calibri'
    
    # Tableau
    rows_count = len(rows) + 1
    cols_count = len(headers)
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS['background']
            paragraph.font.name = 'Calibri'
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Rows
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = COLORS['text']
                paragraph.font.name = 'Calibri'
    
    return slide

def create_pitch_deck():
    """Crée la présentation Pitch Deck complète et professionnelle"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1 : Titre
    create_title_slide(prs, 
                       "YUKPOMNANG",
                       "Plateforme Intelligente de Services Multi-Secteurs en Afrique",
                       "Application déjà développée | Financement : 470M-550M FCFA | ROI : 10-15x sur 5 ans")
    
    # Slide 2 : Problématique Sociétale
    create_content_slide(prs, "UNE RÉPONSE SOCIÉTALE À LA FRACTURE DIGITALE",
                        [
                            "**Millions de commerçants** exclus : 1-2M FCFA pour site web/vidéos",
                            "**Dizaines de millions de familles** : 50K-150K FCFA/mois juste pour chercher services",
                            "**5-10 applications différentes** pour besoins quotidiens",
                            "**Exclusion digitale** perpétue pauvreté et limite accès services essentiels"
                        ],
                        highlight_first=True)
    
    # Slide 3 : Solution Yukpomnang
    create_content_slide(prs, "YUKPOMNANG : LA SOLUTION",
                        [
                            "**Création automatique présence digitale par IA** : Gratuite en 5 minutes",
                            "**Génération vidéo automatisée** : Réduction 90% coûts marketing",
                            "**Recherche intelligente langage naturel** : Économie 50K-90K FCFA/mois par famille",
                            "**Optimisation livraison par IA** : Réduction 40-60% coûts",
                            "**Écosystème intégré multi-secteurs** : Une seule app pour tous services"
                        ])
    
    # Slide 4 : Marché
    create_stat_slide(prs, "MARCHÉ ADRESSABLE : 2-3 BILLIONS FCFA",
                     [
                         ("PIB Secteur Informel", "60-100 billions FCFA"),
                         ("Marché Digitalisable", "2-3 billions FCFA"),
                         ("Population Cible", "~180 millions"),
                         ("Pays Cibles", "11 pays francophones")
                     ])
    
    # Slide 5 : Fonctionnalités Intelligentes
    create_content_slide(prs, "FONCTIONNALITÉS INTELLIGENTES DÉJÀ OPÉRATIONNELLES",
                        [
                            "**Création automatique présence digitale (IA)** : Fiche complète en 5 min - GRATUIT",
                            "**Génération vidéo publicitaire** : TikTok/Reels-like en 2 min - 15K-30K vs 200K-1M",
                            "**Recherche intelligente + géolocalisation** : Langage naturel, 25+ catégories, 111+ filtres",
                            "**Livraison optimisée par IA** : Groupage multi-livraisons, trajet optimal",
                            "**Écosystème multi-secteurs** : Santé, éducation, transport, immobilier, e-commerce"
                        ])
    
    # Slide 6 : Innovation Technologique
    create_content_slide(prs, "INNOVATION TECHNOLOGIQUE YUKPO",
                        [
                            "**Backend haute performance** : Rust/Axum (10x Node.js)",
                            "**IA multi-modèles** : Orchestration intelligente (GPT-4, Mistral, Claude) avec fallback",
                            "**Base de données vectorielle** : PostgreSQL/pgvector pour recherche sémantique",
                            "**Génération vidéo automatisée** : Pipeline Remotion (Docker GPU)",
                            "**Géolocalisation intelligente** : Algorithmes optimisation routes par IA",
                            "**Multi-plateformes** : Web (React), Mobile (React Native), Cloud-native"
                        ])
    
    # Slide 7 : Projections Revenus
    create_table_slide(prs, "PROJECTIONS REVENUS 5 ANS",
                      ["Année", "Produits Actifs", "Placement", "Total Revenus"],
                      [
                          ["2026 (6m)", "40K", "330M", "378M"],
                          ["2027", "120K", "2,880M", "2,905M"],
                          ["2028", "400K", "9,600M", "9,650M"],
                          ["2029", "3,5M", "84,000M", "84,090M"],
                          ["2030", "7,5M", "180,000M", "180,150M"]
                      ])
    
    # Slide 8 : Rentabilité
    create_table_slide(prs, "RENTABILITÉ SUR 5 ANS",
                      ["Année", "Revenus", "Charges", "Résultat", "Marge"],
                      [
                          ["2026 (6m)", "378M", "276M", "102M", "27%"],
                          ["2027", "2,905M", "1,200M", "1,705M", "59%"],
                          ["2028", "9,650M", "3,500M", "6,150M", "64%"],
                          ["2029", "84,090M", "25,000M", "59,090M", "70%"],
                          ["2030", "180,150M", "45,000M", "135,150M", "75%"]
                      ])
    
    # Slide 9 : Besoins Financement
    create_stat_slide(prs, "BESOINS DE FINANCEMENT : 470M-550M FCFA",
                     [
                         ("Marketing", "264M (56%)"),
                         ("Infrastructure", "60M (13%)"),
                         ("Direction & Tech", "48M (10%)"),
                         ("Charges Opérationnelles", "39,6M (8%)"),
                         ("Équipe Commerciale", "21,6M (5%)"),
                         ("Support & Réserve", "37,2M (8%)")
                     ])
    
    # Slide 10 : Timeline Utilisation Fonds
    create_content_slide(prs, "TIMELINE D'UTILISATION DES FONDS (12 MOIS)",
                        [
                            "**Mois 1-3** : 120M (25%) - Recrutement, lancement campagnes, setup infrastructure",
                            "**Mois 4-6** : 120M (25%) - Acquisition 500-800 commerçants, 50K utilisateurs",
                            "**Mois 7-9** : 120M (25%) - Scaling, 150K utilisateurs, expansion géographique",
                            "**Mois 10-12** : 110M (23%) - 200K utilisateurs, 3,000+ commerçants, rentabilité"
                        ])
    
    # Slide 11 : Avantages Concurrentiels
    create_content_slide(prs, "AVANTAGES CONCURRENTIELS DÉCISIFS",
                        [
                            "**Intégration unique** : Une seule app vs 5-10 apps concurrentes",
                            "**IA intégrée** : Création automatique gratuite vs 1-2M FCFA concurrents",
                            "**Géolocalisation intelligente** : Optimisation routes IA (réduction 40-60% coûts)",
                            "**Modèle adapté marché** : Pay-per-use (2K/mois) vs abonnements inaccessibles",
                            "**Focus secteur informel** : 85% commerces exclus par concurrents",
                            "**First-mover avantage** : Fenêtre 2026-2027 avant géants tech"
                        ])
    
    # Slide 12 : Risques et Mitigation
    create_content_slide(prs, "RISQUES ET MITIGATION",
                        [
                            "**Adoption lente** : Marketing agressif (264M/an), freemium, support 24/7",
                            "**Concurrence** : First-mover, barrière technologique (Rust/IA), effets réseau",
                            "**Technique/scaling** : Architecture cloud-native, monitoring, équipe expérimentée",
                            "**Réglementaire** : Conformité RGPD, veille réglementaire continue",
                            "**Financier** : Modèle validé, trésorerie 3-4 mois, diversification revenus"
                        ])
    
    # Slide 13 : Plan de Sortie
    create_content_slide(prs, "PLAN DE SORTIE (EXIT STRATEGY)",
                        [
                            "**Horizon** : 5-7 ans (2029-2031)",
                            "**Valorisation estimée** : 50-100 billions FCFA",
                            "**Option 1** : Acquisition stratégique (Jumia, MTN, Orange, géants tech)",
                            "**Option 2** : IPO après 5 ans croissance",
                            "**Option 3** : Fusion avec autre plateforme tech africaine",
                            "**Garanties investisseurs** : Clauses anti-dilution, retour minimum 5x"
                        ])
    
    # Slide 14 : Fondateur
    create_content_slide(prs, "FONDATEUR : LELE SIAKA HERNANDEZ",
                        [
                            "**13+ ans d'expérience** : Data analyst, statisticien, modélisation avancée",
                            "**Expertise terrain** : Suivi-évaluation projets santé publique à fort impact",
                            "**Collaborations internationales** : Partenaire technique avec **UNICEF** et **BAD**",
                            "**Vision** : Transformer besoins quotidiens en solutions intelligentes, scalables",
                            "**Impact social** : Solutions à fort impact dans contextes complexes"
                        ])
    
    # Slide 15 : Pourquoi Investir Maintenant
    create_stat_slide(prs, "POURQUOI INVESTIR MAINTENANT ?",
                     [
                         ("Application développée", "Risque technique réduit"),
                         ("Marché en explosion", "1,4M → 2,5M habitants 2050"),
                         ("First-mover", "Fenêtre 2026-2027"),
                         ("ROI projeté", "10-15x sur 5 ans")
                     ])
    
    # Slide 16 : Vision 2030
    create_content_slide(prs, "VISION 2030",
                        [
                            "**40+ millions d'utilisateurs actifs**",
                            "**7,5+ millions de produits/services actifs**",
                            "**180+ milliards FCFA de revenus annuels**",
                            "**Leader panafricain** des services essentiels",
                            "**Impact social** : Amélioration qualité de vie millions d'Africains"
                        ])
    
    # Slide 17 : Contact
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_shape = end_slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "MERCI POUR VOTRE ATTENTION"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.name = 'Calibri'
    
    # Contact
    contact_shape = end_slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1.5))
    contact_frame = contact_shape.text_frame
    contact_frame.text = "Hernandez LELE\nFondateur & CEO\nlelehernandez2007@yahoo.fr\n+237 674 546895"
    contact_frame.paragraphs[0].font.size = Pt(20)
    contact_frame.paragraphs[0].font.color.rgb = COLORS['text']
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    contact_frame.paragraphs[0].font.name = 'Calibri'
    
    # Questions
    questions_shape = end_slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(0.8))
    questions_frame = questions_shape.text_frame
    questions_frame.text = "Questions ?"
    questions_frame.paragraphs[0].font.size = Pt(24)
    questions_frame.paragraphs[0].font.bold = True
    questions_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    questions_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    questions_frame.paragraphs[0].font.name = 'Calibri'
    
    # Sauvegarder
    pptx_file = Path(__file__).parent / 'PITCH_DECK_YUKPOMNANG_PROFESSIONAL.pptx'
    try:
        prs.save(pptx_file)
        print(f"[OK] Presentation creee : {pptx_file.name}")
        print(f"     {len(prs.slides)} slides professionnelles")
        print(f"     Prete pour presentation investisseurs")
        return True
    except Exception as e:
        print(f"[ERREUR] Erreur lors de la creation : {e}")
        return False

if __name__ == '__main__':
    print("=" * 70)
    print("Creation de la presentation PowerPoint professionnelle")
    print("YUKPOMNANG - Pitch Deck Investisseurs")
    print("=" * 70)
    print()
    
    success = create_pitch_deck()
    
    print()
    if success:
        print("=" * 70)
        print("[OK] Presentation creee avec succes !")
        print("=" * 70)
        print("\nFichier genere : PITCH_DECK_YUKPOMNANG_PROFESSIONAL.pptx")
        print("Vous pouvez maintenant l'ouvrir dans PowerPoint")
        print("et personnaliser les couleurs/designs si necessaire.")
    else:
        print("=" * 70)
        print("[ERREUR] Erreur lors de la creation")
        print("=" * 70)

