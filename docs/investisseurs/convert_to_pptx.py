#!/usr/bin/env python3
"""
Script pour convertir les documents Markdown en fichiers PowerPoint (.pptx)
avec formatage professionnel adapté aux présentations investisseurs
"""

import os
import sys
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installation de python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

# Mapping des emojis vers texte
EMOJI_MAP = {
    '📊': '', '🎯': '', '💡': '', '🚀': '', '📈': '', '💰': '', '👥': '',
    '⚙️': '', '📢': '', '🗺️': '', '📱': '', '🎨': '', '🌍': '', '🚨': '',
    '✅': '', '❌': '', '⚠️': '', '📋': '', '📍': '', '🏢': '', '💼': '',
    '📅': '', '🔄': '', '📊': '', '🎯': '', '💸': '', '📚': '', '🔍': ''
}

# Couleurs professionnelles
COLORS = {
    'primary': RGBColor(0, 51, 102),      # Bleu foncé
    'secondary': RGBColor(0, 102, 204),   # Bleu moyen
    'accent': RGBColor(255, 140, 0),      # Orange
    'success': RGBColor(0, 153, 51),      # Vert
    'text': RGBColor(51, 51, 51),         # Gris foncé
    'light': RGBColor(128, 128, 128),     # Gris clair
    'background': RGBColor(255, 255, 255), # Blanc
}

def clean_text(text):
    """Nettoie le texte des emojis et caractères problématiques"""
    if not text:
        return ""
    
    # Supprimer les emojis
    for emoji, replacement in EMOJI_MAP.items():
        text = text.replace(emoji, replacement)
    
    # Supprimer les références [^1]
    text = re.sub(r'\[\^\d+\]', '', text)
    text = re.sub(r'^\[\^\d+\]:\s*', '', text, flags=re.MULTILINE)
    
    # Nettoyer les étoiles isolées
    text = re.sub(r'\s+\*\s+', ' ', text)
    text = re.sub(r'^\*\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\s+\*$', '', text, flags=re.MULTILINE)
    
    return text.strip()

def process_bold_text(text):
    """Traite le texte en gras (**text**)"""
    parts = []
    pattern = r'\*\*([^*]+)\*\*'
    
    last_end = 0
    for match in re.finditer(pattern, text):
        if match.start() > last_end:
            parts.append(('normal', text[last_end:match.start()]))
        parts.append(('bold', match.group(1)))
        last_end = match.end()
    
    if last_end < len(text):
        parts.append(('normal', text[last_end:]))
    
    if not parts:
        parts.append(('normal', text))
    
    return parts

def add_text_to_shape(shape, text, is_title=False):
    """Ajoute du texte formaté à une forme"""
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Configurer les marges
    text_frame.margin_left = Inches(0.5)
    text_frame.margin_right = Inches(0.5)
    text_frame.margin_top = Inches(0.3)
    text_frame.margin_bottom = Inches(0.3)
    
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    
    # Traiter le formatage
    parts = process_bold_text(clean_text(text))
    
    for part_type, part_text in parts:
        if part_text:
            run = p.add_run()
            run.text = part_text
            if is_title:
                run.font.size = Pt(44)
                run.font.bold = True
                run.font.color.rgb = COLORS['primary']
            else:
                if part_type == 'bold':
                    run.font.size = Pt(18)
                    run.font.bold = True
                    run.font.color.rgb = COLORS['primary']
                else:
                    run.font.size = Pt(16)
                    run.font.color.rgb = COLORS['text']
    
    return text_frame

def add_bullet_points(slide, content, left=Inches(1), top=Inches(2), width=Inches(8.5), height=Inches(4.5)):
    """Ajoute des puces à une slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.2)
    text_frame.margin_bottom = Inches(0.2)
    
    lines = content.split('\n')
    first = True
    
    for line in lines:
        line = clean_text(line.strip())
        if not line or line.startswith('---'):
            continue
        
        # Enlever les puces markdown
        if line.startswith('- ') or line.startswith('* '):
            line = line[2:].strip()
        
        if line:
            if first:
                p = text_frame.paragraphs[0]
                first = False
            else:
                p = text_frame.add_paragraph()
            
            p.level = 0
            p.space_after = Pt(6)
            
            # Traiter le formatage
            parts = process_bold_text(line)
            for part_type, part_text in parts:
                if part_text:
                    run = p.add_run()
                    run.text = part_text
                    run.font.size = Pt(16)
                    if part_type == 'bold':
                        run.font.bold = True
                        run.font.color.rgb = COLORS['primary']
                    else:
                        run.font.color.rgb = COLORS['text']
    
    return textbox

def create_title_slide(prs, title, subtitle=""):
    """Crée une slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide
    
    # Fond avec couleur
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['background']
    
    # Titre principal
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.text = clean_text(title)
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    if subtitle:
        subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = clean_text(subtitle)
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = COLORS['light']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, content):
    """Crée une slide de contenu avec titre et contenu"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide
    
    # Titre
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = clean_text(title)
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    # Ligne de séparation sous le titre
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(9.5), Inches(1.2))
    line.line.color.rgb = COLORS['light']
    line.line.width = Pt(2)
    
    # Contenu
    add_bullet_points(slide, content, top=Inches(1.5))
    
    return slide

def markdown_to_pptx(md_file, pptx_file):
    """Convertit un fichier Markdown en fichier PowerPoint"""
    
    # Lire le fichier Markdown
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()
    
    # Créer une nouvelle présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Parser le Markdown
    lines = md_content.split('\n')
    i = 0
    current_slide_title = None
    current_slide_content = []
    
    # Slide de titre
    title = Path(md_file).stem.replace('_', ' ').title()
    subtitle = "Document Investisseurs - Janvier 2026"
    create_title_slide(prs, title, subtitle)
    
    while i < len(lines):
        line = lines[i].strip()
        
        if not line:
            i += 1
            continue
        
        # Titre niveau 1 (#) - Nouvelle slide
        if line.startswith('# '):
            # Sauvegarder la slide précédente si elle existe
            if current_slide_title and current_slide_content:
                create_content_slide(prs, current_slide_title, '\n'.join(current_slide_content))
            
            current_slide_title = clean_text(line[2:].strip())
            current_slide_content = []
            i += 1
        
        # Titre niveau 2 (##) - Nouvelle slide
        elif line.startswith('## '):
            # Sauvegarder la slide précédente
            if current_slide_title and current_slide_content:
                create_content_slide(prs, current_slide_title, '\n'.join(current_slide_content))
            
            current_slide_title = clean_text(line[3:].strip())
            current_slide_content = []
            i += 1
        
        # Titre niveau 3 (###) - Sous-titre dans la slide
        elif line.startswith('### '):
            if current_slide_content:
                current_slide_content.append('')
            current_slide_content.append('**' + clean_text(line[4:].strip()) + '**')
            i += 1
        
        # Ligne horizontale (---) - Ignorer
        elif line.startswith('---'):
            i += 1
        
        # Tableau - Simplifier en liste
        elif '|' in line and line.count('|') >= 2:
            # Ignorer les lignes de séparation
            if not line.startswith('|---'):
                cells = [c.strip() for c in line.split('|') if c.strip()]
                if cells and len(cells) >= 2:
                    # Prendre les 2 premières colonnes comme format liste
                    if len(cells) >= 2:
                        current_slide_content.append(f"- **{cells[0]}** : {cells[1]}")
            i += 1
        
        # Liste à puces
        elif line.startswith('- ') or line.startswith('* '):
            text = clean_text(line[2:].strip())
            if text:
                current_slide_content.append(text)
            i += 1
        
        # Paragraphe normal
        else:
            text = clean_text(line)
            if text and not text.startswith('[') and not text.startswith('*'):
                current_slide_content.append(text)
            i += 1
    
    # Ajouter la dernière slide
    if current_slide_title and current_slide_content:
        create_content_slide(prs, current_slide_title, '\n'.join(current_slide_content))
    
    # Slide de fin
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    end_shape = end_slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    end_frame = end_shape.text_frame
    end_frame.text = "Merci pour votre attention"
    end_frame.paragraphs[0].font.size = Pt(36)
    end_frame.paragraphs[0].font.bold = True
    end_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    end_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Sauvegarder
    try:
        prs.save(pptx_file)
        print(f"[OK] Converti : {Path(md_file).name} -> {Path(pptx_file).name}")
    except PermissionError:
        print(f"[ERREUR] Permission refusee : {Path(pptx_file).name} est peut-etre ouvert dans PowerPoint")
        print(f"         Veuillez fermer le fichier et reessayer")
    except Exception as e:
        print(f"[ERREUR] Erreur lors de la conversion de {md_file}: {e}")

def create_executive_summary_pptx():
    """Crée une présentation PowerPoint spéciale pour l'Executive Summary"""
    md_file = Path(__file__).parent / '01_EXECUTIVE_SUMMARY.md'
    pptx_file = Path(__file__).parent / '01_EXECUTIVE_SUMMARY.pptx'
    
    if not md_file.exists():
        print(f"[ATTENTION] Fichier non trouve : {md_file}")
        return
    
    # Lire le contenu
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1 : Titre
    create_title_slide(prs, "Yukpomnang", "Plateforme Intelligente de Services Multi-Secteurs en Afrique")
    
    # Slide 2 : Vision
    create_content_slide(prs, "Vision", 
        "Yukpomnang transforme l'accès aux services essentiels en Afrique\n"
        "• Intelligence artificielle avancée\n"
        "• Géolocalisation intelligente\n"
        "• Création vidéo automatisée\n"
        "• Multi-secteurs intégrés")
    
    # Slide 3 : Problématique 1
    create_content_slide(prs, "Problématique #1 : Invisibilité Digitale",
        "85% des commerces locaux africains n'ont aucune présence digitale\n"
        "• 70% de l'économie informelle invisible\n"
        "• Perte de 30-40% du chiffre d'affaires potentiel\n"
        "• Coûts marketing inaccessibles (500K-2M FCFA)\n"
        "• Fracture digitale massive")
    
    # Slide 4 : Solution 1
    create_content_slide(prs, "Notre Solution #1",
        "Démocratisation de l'accès au digital\n"
        "• Création automatique en 5 minutes (GRATUIT)\n"
        "• Vidéos publicitaires en 2 minutes (15K-30K FCFA)\n"
        "• Visibilité géolocalisée automatique\n"
        "• Outils marketing accessibles")
    
    # Slide 5 : Problématique 2
    create_content_slide(prs, "Problématique #2 : Fracture d'Accès",
        "Services fragmentés, coûteux et inaccessibles\n"
        "• 5-10 applications différentes pour besoins quotidiens\n"
        "• Coûts cachés : 50K-150K FCFA/mois par ménage\n"
        "• 60% de la population exclue (barrières linguistiques)\n"
        "• Manque de transparence")
    
    # Slide 6 : Solution 2
    create_content_slide(prs, "Notre Solution #2",
        "Une seule app pour tous les besoins\n"
        "• Livraisons, santé, éducation, transport, immobilier\n"
        "• Géolocalisation intelligente\n"
        "• Accessibilité multilingue (15+ langues africaines)\n"
        "• Transparence totale en temps réel")
    
    # Slide 7 : Marché
    create_content_slide(prs, "Opportunité de Marché",
        "TAM : 90+ billions FCFA\n"
        "• Croissance 15-20% par an\n"
        "• 1,4 milliard d'habitants → 2,5 milliards en 2050\n"
        "• 70% pénétration smartphone\n"
        "• 85% couverture 4G")
    
    # Slide 8 : Revenus 2026
    create_content_slide(prs, "Projections 2026",
        "Revenus : 1,2 milliard FCFA\n"
        "• 30,000 produits/services actifs\n"
        "• 200,000 utilisateurs actifs\n"
        "• 800,000 livraisons\n"
        "• GMV : 12 milliards FCFA")
    
    # Slide 9 : Financement
    create_content_slide(prs, "Besoins de Financement",
        "Série A : 2,1 - 3 milliards FCFA\n"
        "• Timeline : Levée avant fin février 2026\n"
        "• Recrutement : Avant 10 mars 2026\n"
        "• Lancement : 01 avril 2026")
    
    # Slide 10 : Merci
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    end_shape = end_slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    end_frame = end_shape.text_frame
    end_frame.text = "Merci pour votre attention"
    end_frame.paragraphs[0].font.size = Pt(36)
    end_frame.paragraphs[0].font.bold = True
    end_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    end_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    try:
        prs.save(pptx_file)
        print(f"[OK] Cree : 01_EXECUTIVE_SUMMARY.pptx (presentation speciale)")
    except Exception as e:
        print(f"[ERREUR] Erreur : {e}")

def main():
    """Fonction principale"""
    current_dir = Path(__file__).parent
    
    # Liste des fichiers à convertir
    md_files = [
        '02_BUSINESS_PLAN.md',
        '03_PROJECTIONS_FINANCIERES.md',
        '04_PLAN_MARKETING.md',
        '05_ANALYSE_MARCHE.md',
        '06_PLAN_OPERATIONNEL.md',
        '07_PITCH_DECK.md',
        '08_PLAN_FINANCIER_3_ANS.md',
    ]
    
    print("=" * 60)
    print("Conversion des documents Markdown en PowerPoint")
    print("Formatage professionnel pour presentations")
    print("=" * 60)
    print()
    
    # Créer la présentation spéciale pour Executive Summary
    create_executive_summary_pptx()
    
    success_count = 0
    error_count = 0
    
    for md_file in md_files:
        md_path = current_dir / md_file
        if md_path.exists():
            pptx_file = md_path.stem + '.pptx'
            pptx_path = current_dir / pptx_file
            try:
                markdown_to_pptx(md_path, pptx_path)
                success_count += 1
            except Exception as e:
                print(f"[ERREUR] Erreur lors de la conversion de {md_file}: {e}")
                error_count += 1
        else:
            print(f"[ATTENTION] Fichier non trouve : {md_file}")
            error_count += 1
    
    print()
    print("=" * 60)
    print(f"Conversion terminee : {success_count + 1} succes, {error_count} erreurs")
    print("=" * 60)

if __name__ == '__main__':
    main()

