#!/usr/bin/env python3
"""
Script pour créer une présentation PowerPoint professionnelle du Pitch Deck
optimisée pour les présentations investisseurs
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installation de python-pptx...")
    os.system(f"{sys.executable} -m pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

# Couleurs professionnelles
COLORS = {
    'primary': RGBColor(0, 51, 102),      # Bleu foncé
    'secondary': RGBColor(0, 102, 204),   # Bleu moyen
    'accent': RGBColor(255, 140, 0),      # Orange
    'success': RGBColor(0, 153, 51),      # Vert
    'text': RGBColor(51, 51, 51),         # Gris foncé
    'light': RGBColor(128, 128, 128),     # Gris clair
    'background': RGBColor(255, 255, 255), # Blanc
}

def create_title_slide(prs, title, subtitle=""):
    """Crée une slide de titre professionnelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre principal
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    if subtitle:
        subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = COLORS['light']
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, bullet_points):
    """Crée une slide de contenu avec titre et puces"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    # Ligne de séparation
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(9.5), Inches(1.2))
    line.line.color.rgb = COLORS['light']
    line.line.width = Pt(3)
    
    # Contenu avec puces
    content_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5))
    content_frame = content_shape.text_frame
    content_frame.word_wrap = True
    content_frame.margin_left = Inches(0.2)
    content_frame.margin_right = Inches(0.2)
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.level = 0
        p.space_after = Pt(12)
        p.left_indent = Inches(0.3)
        
        # Séparer le texte en gras si nécessaire
        if '**' in point:
            parts = point.split('**')
            for j, part in enumerate(parts):
                if part.strip():
                    run = p.add_run()
                    run.text = part.strip() + (' ' if j < len(parts) - 1 else '')
                    run.font.size = Pt(20)
                    if j % 2 == 1:  # Texte entre **
                        run.font.bold = True
                        run.font.color.rgb = COLORS['primary']
                    else:
                        run.font.color.rgb = COLORS['text']
        else:
            run = p.add_run()
            run.text = point
            run.font.size = Pt(20)
            run.font.color.rgb = COLORS['text']
    
    return slide

def create_two_column_slide(prs, title, left_col, right_col):
    """Crée une slide avec deux colonnes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Titre
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    # Colonne gauche
    left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_shape.text_frame
    left_frame.word_wrap = True
    left_frame.text = left_col
    for paragraph in left_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = COLORS['text']
        paragraph.space_after = Pt(8)
    
    # Colonne droite
    right_shape = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_shape.text_frame
    right_frame.word_wrap = True
    right_frame.text = right_col
    for paragraph in right_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = COLORS['text']
        paragraph.space_after = Pt(8)
    
    return slide

def create_pitch_deck():
    """Crée la présentation Pitch Deck complète"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1 : Titre
    create_title_slide(prs, "Yukpomnang", 
                       "Plateforme Intelligente de Services Multi-Secteurs en Afrique\n"
                       "Levée de fonds Série A : 2,1 - 3 milliards FCFA")
    
    # Slide 2 : Problème #1
    create_content_slide(prs, "Problème #1 : Invisibilité Digitale",
                        [
                            "85% des commerces locaux africains = AUCUNE présence digitale",
                            "70% de l'économie informelle = Invisible",
                            "Perte de 30-40% du chiffre d'affaires potentiel",
                            "Coûts marketing inaccessibles : 500K-2M FCFA"
                        ])
    
    # Slide 3 : Problème #2
    create_content_slide(prs, "Problème #2 : Fracture d'Accès",
                        [
                            "5-10 applications différentes pour besoins quotidiens",
                            "Coûts cachés : 50K-150K FCFA/mois par ménage",
                            "60% de la population exclue (barrières linguistiques)",
                            "Manque de transparence : prix, disponibilité, qualité"
                        ])
    
    # Slide 4 : Solution
    create_content_slide(prs, "Notre Solution",
                        [
                            "Une seule app pour tous les services essentiels",
                            "Création automatique présence digitale (5 min, GRATUIT)",
                            "Génération vidéo publicitaire automatisée (2 min, 15K-30K FCFA)",
                            "Géolocalisation intelligente + Accessibilité multilingue (15+ langues)"
                        ])
    
    # Slide 5 : Marché
    create_content_slide(prs, "Opportunité de Marché Massive",
                        [
                            "TAM : 90+ billions FCFA (croissance 15-20%/an)",
                            "1,4 milliard d'habitants → 2,5 milliards en 2050",
                            "70% pénétration smartphone, 85% couverture 4G",
                            "Classe moyenne émergente : +15M personnes/an"
                        ])
    
    # Slide 6 : Modèle Économique
    create_content_slide(prs, "Modèle Économique",
                        [
                            "Achat de tokens (pay-per-use) : 60% des revenus",
                            "Commissions livraisons : 10%",
                            "Commissions ventes : 20%",
                            "Publicité & vidéos : 8%",
                            "Services spécialisés : 2%"
                        ])
    
    # Slide 7 : Traction & Métriques
    create_content_slide(prs, "Projections 2026 (Cameroun)",
                        [
                            "30,000 produits/services actifs",
                            "200,000 utilisateurs actifs",
                            "800,000 livraisons",
                            "Revenus : 1,2 milliard FCFA",
                            "GMV : 12 milliards FCFA"
                        ])
    
    # Slide 8 : Avantages Concurrentiels
    create_content_slide(prs, "Avantages Concurrentiels",
                        [
                            "Première plateforme intégrée multi-secteurs en Afrique",
                            "IA propriétaire : création automatique produits/vidéos",
                            "Géolocalisation avancée avec optimisation routes",
                            "Accessibilité multilingue (15+ langues africaines)",
                            "Premier mover advantage"
                        ])
    
    # Slide 9 : Stratégie de Croissance
    create_content_slide(prs, "Stratégie de Croissance",
                        [
                            "2026 : Cameroun (point d'entrée livraisons)",
                            "2027 : Côte d'Ivoire + Sénégal",
                            "2028 : Expansion 10+ pays francophones",
                            "2029 : Marchés anglophones/lusophones",
                            "2030 : Leader panafricain"
                        ])
    
    # Slide 10 : Équipe
    create_content_slide(prs, "Équipe & Gouvernance",
                        [
                            "Phase initiale : 5 personnes maximum (réaliste)",
                            "Priorité : Agents commerciaux (2-3 personnes)",
                            "Extension progressive selon croissance",
                            "2027 : 50-80 personnes (expansion 3 pays)",
                            "2028 : 150-200 personnes (expansion francophone)"
                        ])
    
    # Slide 11 : Financement
    create_content_slide(prs, "Besoins de Financement",
                        [
                            "Série A : 2,1 - 3 milliards FCFA",
                            "Timeline critique :",
                            "  • Levée de fonds : Avant fin février 2026",
                            "  • Recrutement : Avant 10 mars 2026",
                            "  • Lancement : 01 avril 2026"
                        ])
    
    # Slide 12 : Projections Financières
    create_two_column_slide(prs, "Projections Financières 3 Ans",
                           "2026\n• Revenus : 1,2 milliard FCFA\n• Charges : 1,4 milliard\n• EBITDA : -200M FCFA\n\n2027\n• Revenus : 6,5 milliards\n• Charges : 5,85 milliards\n• EBITDA : +650M FCFA",
                           "2028\n• Revenus : 18 milliards\n• Charges : 14,4 milliards\n• EBITDA : +3,6 milliards\n\nMarge EBITDA\n• 2026 : -17%\n• 2027 : +10%\n• 2028 : +20%")
    
    # Slide 13 : Valorisation & Sortie
    create_content_slide(prs, "Valorisation & Sortie",
                        [
                            "Valorisation pré-money Série A : 10-15 milliards FCFA",
                            "Projection valorisation 2028 : 150-200 milliards FCFA",
                            "Options de sortie :",
                            "  • Acquisition par géant tech (Jumia, Glovo, etc.)",
                            "  • IPO sur bourse africaine",
                            "  • Consolidation avec autres plateformes"
                        ])
    
    # Slide 14 : Risques & Mitigation
    create_content_slide(prs, "Risques & Mitigation",
                        [
                            "Risque marché : Diversification géographique",
                            "Risque concurrence : Avantage premier mover + IA",
                            "Risque opérationnel : Équipe expérimentée + processus",
                            "Risque réglementaire : Conformité proactive",
                            "Risque technologique : Architecture scalable"
                        ])
    
    # Slide 15 : Vision
    create_content_slide(prs, "Vision 2030",
                        [
                            "Leader panafricain des services essentiels",
                            "40+ millions d'utilisateurs actifs",
                            "7,5+ millions de produits/services actifs",
                            "132+ milliards FCFA de revenus",
                            "Impact social : Amélioration qualité de vie millions d'Africains"
                        ])
    
    # Slide 16 : Appel à l'Action
    create_content_slide(prs, "Rejoignez-Nous",
                        [
                            "Opportunité unique de transformer l'accès aux services en Afrique",
                            "Marché de 90+ billions FCFA en croissance",
                            "Équipe dédiée, technologie éprouvée, modèle validé",
                            "Timeline agressive mais réaliste",
                            "Contactez-nous pour en savoir plus"
                        ])
    
    # Slide 17 : Merci
    end_slide = prs.slides.add_slide(prs.slide_layouts[6])
    end_shape = end_slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    end_frame = end_shape.text_frame
    end_frame.text = "Merci pour votre attention\n\nQuestions ?"
    end_frame.paragraphs[0].font.size = Pt(40)
    end_frame.paragraphs[0].font.bold = True
    end_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    end_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Sauvegarder
    pptx_file = Path(__file__).parent / '07_PITCH_DECK_PRESENTATION.pptx'
    try:
        prs.save(pptx_file)
        print(f"[OK] Cree : 07_PITCH_DECK_PRESENTATION.pptx (17 slides)")
    except Exception as e:
        print(f"[ERREUR] Erreur : {e}")

if __name__ == '__main__':
    print("=" * 60)
    print("Creation de la presentation Pitch Deck professionnelle")
    print("=" * 60)
    print()
    create_pitch_deck()
    print()
    print("=" * 60)
    print("Presentation creee avec succes !")
    print("=" * 60)

